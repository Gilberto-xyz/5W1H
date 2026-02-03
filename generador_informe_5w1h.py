# -*- coding: utf-8 -*-
"""Genera informes 5W1H en PowerPoint a partir de datos Excel.

Incluye tablas y graficos para multiples segmentos y comparativos.
"""
# Bibliotecas necesarias
#---------------------------------------------------------------------------------------------------------------------
import pandas as pd
import numpy as np
import warnings
from datetime import datetime as dt
import os
import sys
import io
import math
import textwrap
import re
import unicodedata
import matplotlib
matplotlib.use("Agg")
from matplotlib import pyplot as plt
from matplotlib.ticker import FuncFormatter
from matplotlib import colors as mcolors
import matplotlib.patches as mpatches
from matplotlib.patches import Rectangle
from matplotlib.font_manager import FontProperties
from dataclasses import dataclass
from decimal import Decimal, ROUND_CEILING, ROUND_DOWN, ROUND_FLOOR, ROUND_HALF_DOWN, ROUND_HALF_EVEN, ROUND_HALF_UP, ROUND_UP, ROUND_05UP
from pptx import Presentation 
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from typing import Dict, Iterable, List, NamedTuple, Optional, Tuple
from pathlib import Path
from collections import OrderedDict

# Fuerza E/S en UTF-8 para soportar acentos y ñ en la terminal de Windows
def _configure_utf8_io():
    """Configura E/S para forzar UTF-8 en la consola."""
    try:
        os.environ.setdefault('PYTHONUTF8', '1')
        os.environ.setdefault('PYTHONIOENCODING', 'utf-8')
        for stream_name in ('stdout', 'stderr', 'stdin'):
            stream = getattr(sys, stream_name, None)
            if stream is not None and hasattr(stream, 'reconfigure'):
                stream.reconfigure(encoding='utf-8')
    except Exception:
        # Si la consola no soporta reconfigure, seguimos sin romper ejecucion
        pass

_configure_utf8_io()
COLOR_BLUE = '\033[94m'
COLOR_YELLOW = '\033[93m'
COLOR_GREEN = '\033[92m'
COLOR_RED = '\033[91m'
COLOR_RESET = '\033[0m'
COLOR_QUESTION = '\033[38;5;37m'
HEADER_COLOR_PRIMARY = '#286B72'
HEADER_COLOR_SECONDARY = '#3EBBC7'
HEADER_FONT_COLOR = '#FFFFFF'
HEADER_TOTAL_FILL = '#E8F0FF'
HEADER_FIRST_COL_FILL = '#F5F5F5'
TABLE_TEXT_PRIMARY = '#1F1F1F'
TABLE_POSITIVE_COLOR = '#27AE60'
TABLE_NEGATIVE_COLOR = '#C0392B'
TABLE_GRID_COLOR = '#1A1A1A'
APORTE_NEGATIVE_COLOR = "#F3948B"
APORTE_NEUTRAL_COLOR = "#FFCE7E"
VOLUME_BAR_START = '#1452B8'
VOLUME_BAR_END = '#8CB8FF'
VOLUME_BAR_ALPHA = 0.85
VOLUME_BAR_BASE_ALPHA = 0.18
LINE_COLOR_CLIENT = '#2C3E50'
LINE_COLOR_NUMERATOR = "#D4AC0D"
BAR_COLOR_VARIATION_CLIENT = '#7F8C8D'
BAR_COLOR_VARIATION_NUMERATOR = '#F1C40F'
BAR_EDGE_COLOR = 'black'
BAR_LABEL_COLOR_POS = '#1E8449'
BAR_LABEL_COLOR_NEG = '#8B0000'
BAR_LABEL_COLOR_POS_ALT = '#27AE60'
BAR_LABEL_COLOR_NEG_ALT = '#C0392B'
REFERENCE_SERIES_COLOR = '#6E6E6E'
ANNOTATION_BOX_FACE = '#F2F2F2'
ANNOTATION_BOX_EDGE = 'black'
TREND_COLOR_PALETTE = {
    'trend_01': '#1F77B4',
    'trend_02': '#FF7F0E',
    'trend_03': '#2CA02C',
    'trend_04': '#D62728',
    'trend_05': '#9467BD',
    'trend_06': '#8C564B',
    'trend_07': '#E377C2',
    'trend_08': '#7F7F7F',
    'trend_09': '#BCBD22',
    'trend_10': '#17BECF',
    'trend_11': '#AEC7E8',
    'trend_12': '#FFBB78',
    'trend_13': '#98DF8A',
    'trend_14': '#FF9896',
    'trend_15': '#C5B0D5',
    'trend_16': '#C49C94',
    'trend_17': '#F7B6D2',
    'trend_18': '#C7C7C7',
    'trend_19': '#DBDB8D',
    'trend_20': '#9EDAE5',
    'trend_21': '#393B79',
    'trend_22': '#5254A3',
    'trend_23': '#6B6ECF',
    'trend_24': '#9C9EDE',
    'trend_25': '#637939',
    'trend_26': '#8CA252',
    'trend_27': '#B5CF6B',
    'trend_28': '#CEDB9C',
    'trend_29': '#8C6D31',
    'trend_30': '#BD9E39',
    'trend_31': '#E7BA52',
    'trend_32': '#E7CB94',
    'trend_33': '#843C39',
    'trend_34': '#AD494A',
    'trend_35': '#D6616B',
    'trend_36': '#E7969C',
    'trend_37': '#7B4173',
    'trend_38': '#A55194',
    'trend_39': '#CE6DBD',
    'trend_40': '#DE9ED6',
    'trend_41': '#3182BD',
    'trend_42': '#6BAED6',
    'trend_43': '#9ECAE1',
    'trend_44': '#C6DBEF',
    'trend_45': '#E6550D',
    'trend_46': '#FD8D3C',
    'trend_47': '#FDAE6B',
    'trend_48': '#FDD0A2',
    'trend_49': '#31A354',
    'trend_50': '#74C476',
    'trend_51': '#A1D99B',
    'trend_52': '#C7E9C0',
    'trend_53': '#756BB1',
    'trend_54': '#9E9AC8',
    'trend_55': '#BCBDDC',
    'trend_56': '#DADAEB',
    'trend_57': '#636363',
    'trend_58': '#969696',
    'trend_59': '#BDBDBD',
    'trend_60': '#D9D9D9',
    'trend_61': '#393E46',
    'trend_62': '#00ADB5',
    'trend_63': '#FF5722',
    'trend_64': '#795548',
    'trend_65': '#607D8B',
    'trend_66': '#8BC34A',
    'trend_67': '#CDDC39',
    'trend_68': '#FFC107',
    'trend_69': '#FF4081',
    'trend_70': '#3F51B5'
}
TREND_COLOR_SEQUENCE = list(TREND_COLOR_PALETTE.values())
# Paleta reservada para títulos de Competencia: usa los últimos colores para no interferir con los de marcas.
COMPETITION_TITLE_PALETTE = list(reversed(TREND_COLOR_SEQUENCE[-8:]))
# Paleta separada para títulos de Categoría (segmento 8). Usa un tramo distinto
# al inicio de la paleta para evitar colisión con colores de marcas.
CATEGORY_TITLE_PALETTE = TREND_COLOR_SEQUENCE[8:16] if len(TREND_COLOR_SEQUENCE) >= 16 else TREND_COLOR_SEQUENCE
TABLE_WRAP_WIDTH = 14
DISPLAY_TREND_REFERENCE_TEXT = False
TREND_SCALE_RULES = [
    {
        "threshold": 1_000_000_000_000,
        "divisor": 1_000_000_000_000,
        "suffix": "T",
        "reference": {
            'E': 'Valores expresados en billones',
            'P': 'Valores expressos em trilhões',
            'default': 'Values expressed in trillions'
        }
    },
    {
        "threshold": 1_000_000_000,
        "divisor": 1_000_000_000,
        "suffix": "B",
        "reference": {
            'E': 'Valores expresados en miles de millones',
            'P': 'Valores expressos em bilhões',
            'default': 'Values expressed in billions'
        }
    },
    {
        "threshold": 1_000_000,
        "divisor": 1_000_000,
        "suffix": "M",
        "reference": {
            'E': 'Valores expresados en millones',
            'P': 'Valores expressos em milhões',
            'default': 'Values expressed in millions'
        }
    },
    {
        "threshold": 1_000,
        "divisor": 1_000,
        "suffix": "K",
        "reference": {
            'E': 'Valores expresados en miles',
            'P': 'Valores expressos em milhares',
            'default': 'Values expressed in thousands'
        }
    },
    {
        "threshold": 0,
        "divisor": 1,
        "suffix": "",
        "reference": {
            'E': '',
            'P': '',
            'default': ''
        }
    },
]
MONTH_NAMES = {
    'P': ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho', 'Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro'],
    'E': ['Enero', 'Febrero', 'Marzo', 'Abril', 'Mayo', 'Junio', 'Julio', 'Agosto', 'Septiembre', 'Octubre', 'Noviembre', 'Diciembre']
}
def colorize(text: str, color: str = COLOR_BLUE) -> str:
    """Devuelve el texto envuelto con el codigo ANSI de color."""
    return f"{color}{text}{COLOR_RESET}"
def print_colored(text: str, color: str = COLOR_BLUE) -> None:
    """Imprime texto coloreado en la terminal."""
    print(colorize(text, color))
def select_trend_scale_rule(max_value: float) -> dict:
    """Selecciona la regla de escala segun el valor maximo."""
    for rule in TREND_SCALE_RULES:
        if max_value >= rule["threshold"]:
            return rule
    return TREND_SCALE_RULES[-1]
def format_value_with_suffix(value: float, divisor: float, suffix: str) -> str:
    """Formatea un valor escalado con sufijo (K/M/B/T)."""
    if not np.isfinite(value):
        return ''
    scaled = value / divisor if divisor else value
    abs_scaled = abs(scaled)
    if abs_scaled >= 100:
        formatted = f"{scaled:.0f}"
    elif abs_scaled >= 10:
        formatted = f"{scaled:.1f}"
    else:
        formatted = f"{scaled:.2f}"
    if '.' in formatted:
        formatted = formatted.rstrip('0').rstrip('.')
    return f"{formatted}{suffix}"
def available_width(presentation, left=Inches(0), right=Inches(0)):
    """Calcula el ancho disponible en la diapositiva."""
    return presentation.slide_width - int(left) - int(right)
def constrain_picture_width(picture, max_width):
    """Limita el ancho de una imagen preservando la relacion de aspecto."""
    if max_width is None:
        return
    max_width = int(max_width)
    if max_width <= 0:
        return
    current_width = int(picture.width)
    if current_width <= max_width:
        return
    current_height = int(picture.height)
    scale = current_width / max_width
    picture.width = max_width
    picture.height = max(int(round(current_height / scale)), 1)
EMU_PER_INCH = 914400
DEFAULT_LINE_CHART_RATIO = 3
TABLE_TARGET_HEIGHT_CM = 4.0
TABLE_SIDE_MARGIN_CM = 1.2
TABLE_PAIR_GAP_CM = 0.5
TABLE_HEADER_FONT_SIZE = 10
TABLE_WRAP_WIDTH = 14
LINE_CHART_LEFT_MARGIN = 0.045
LINE_CHART_RIGHT_MARGIN = 0.99
LINE_CHART_BOTTOM_MARGIN = 0.1
LINE_CHART_TOP_MARGIN = 0.9
LINE_CHART_SINGLE_X_MARGIN = 0.16
LINE_CHART_MULTI_LEFT_MARGIN = 0.08
LINE_CHART_MULTI_RIGHT_MARGIN = 0.98
LINE_CHART_MULTI_BOTTOM_MARGIN = 0.12
LINE_CHART_MULTI_TOP_MARGIN = 0.85
LINE_CHART_MULTI_X_MARGIN = 0.16
DEFAULT_EXPORT_DPI = 110
TABLE_EXPORT_DPI = 220
EXPORT_PAD_INCHES = 0.08
CHART_TOP_INCH = 0.72
def emu_to_inches(value: int) -> float:
    """Convierte unidades EMU a pulgadas."""
    return float(value) / EMU_PER_INCH
def wrap_table_text(value, max_width: int = TABLE_WRAP_WIDTH):
    """
    Devuelve el valor con saltos suaves para ajustar el texto en celdas.
    Mantiene los valores no string sin cambios.
    """
    if not isinstance(value, str):
        return value
    normalized = " ".join(value.split())
    if not normalized:
        return ""
    if len(normalized) <= max_width or " " not in normalized:
        return normalized
    return textwrap.fill(normalized, width=max_width, break_long_words=False)
def figure_to_stream(
    fig: plt.Figure,
    dpi: int = DEFAULT_EXPORT_DPI,
    bbox_inches="tight",
    pad_inches: float = EXPORT_PAD_INCHES,
    transparent: bool = True,
    close: bool = True,
) -> io.BytesIO:
    """
    Renderiza una figura en memoria y la cierra para liberar recursos.
    Replica el flujo rapido usado en coverage_studio (guardar en BytesIO y cerrar).
    """
    buf = io.BytesIO()
    fig.savefig(
        buf,
        format="png",
        dpi=dpi,
        bbox_inches=bbox_inches,
        pad_inches=pad_inches,
        transparent=transparent,
    )
    buf.seek(0)
    if close:
        try:
            plt.close(fig)
        except Exception:
            pass
    return buf

# ---------------- Segmento 8: tabla PPT (intervalos) ----------------

# Valores z: limites (2.0) y error muestral (1.96).
PPT8_Z_SCORE_INTERVAL = 2.0
PPT8_Z_SCORE_ERROR = 1.96


@dataclass
class Ppt8MonthlyBlock:
    brand: str
    rows: pd.DataFrame  # columnas: 0=date, 1=Weighted PENET, 2=BUYERS


@dataclass
class Ppt8AggBlock:
    brand: str
    block: pd.DataFrame  # columnas: 4=label, 5=MAT base, 6=MAT actual


@dataclass
class Ppt8Row:
    marca: str
    mat_sep25: float
    lim_inf: float
    lim_sup: float
    compradores_prom: float
    pen12m: float
    nivel_pen: str
    error_muestral: float


def ppt8_normalize_label(label: str) -> str:
    """Normaliza etiquetas para emparejar bloques por marca."""
    normalized = unicodedata.normalize("NFKD", str(label))
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return re.sub(r"[^a-z0-9]+", "", normalized.lower())


def ppt8_is_header_label(label: str) -> bool:
    """Detecta etiquetas de encabezado en lugar de marcas."""
    raw = str(label).strip()
    if not raw:
        return True
    lower = raw.lower()
    if re.search(r"\btable\b", lower):
        return True
    if re.search(r"\bcateg", lower):
        return True
    if re.search(r"\bbrands\b", lower) and ("=" in lower or "total" in lower):
        return True
    compact = ppt8_normalize_label(label)
    return not compact


def ppt8_block_has_numeric(block: pd.DataFrame) -> bool:
    """Verifica si el bloque tiene valores numericos en las columnas 5/6."""
    if block.empty:
        return False
    for col in (5, 6):
        if block.shape[1] > col:
            numeric = pd.to_numeric(block.iloc[:, col], errors="coerce")
            if numeric.notna().any():
                return True
    return False


def ppt8_filter_agg_blocks(
    monthly_blocks: Iterable[Ppt8MonthlyBlock],
    agg_blocks: Iterable[Ppt8AggBlock],
) -> Tuple[List[Ppt8AggBlock], List[str]]:
    """Omite bloques agregados sin datos o que son encabezados."""
    monthly_norms = {ppt8_normalize_label(block.brand) for block in monthly_blocks if block.brand}
    filtered: List[Ppt8AggBlock] = []
    dropped: List[str] = []
    for block in agg_blocks:
        brand_label = block.brand
        norm = ppt8_normalize_label(brand_label)
        header_like = ppt8_is_header_label(brand_label)
        has_numeric = ppt8_block_has_numeric(block.block)
        if header_like and not has_numeric:
            dropped.append(brand_label)
            continue
        if norm and norm not in monthly_norms and not has_numeric:
            dropped.append(brand_label)
            continue
        filtered.append(block)
    return filtered, dropped


def ppt8_parse_monthly_blocks(df: pd.DataFrame) -> List[Ppt8MonthlyBlock]:
    """Encuentra bloques de datos mensuales delimitados por marca en columnas 0-2."""
    first_col = df.iloc[:, 0]
    brand_rows = df[(first_col.notna()) & df.iloc[:, 1].isna() & df.iloc[:, 2].isna()].index.tolist()
    blocks: List[Ppt8MonthlyBlock] = []
    for i, start in enumerate(brand_rows):
        end = brand_rows[i + 1] if i + 1 < len(brand_rows) else len(df)
        data = df.loc[start + 1 : end - 1]
        brand_label = str(df.iloc[start, 0]).strip()
        if ppt8_is_header_label(brand_label):
            continue
        data = data[pd.to_numeric(data.iloc[:, 1], errors="coerce").notna()]
        if data.empty:
            continue
        blocks.append(Ppt8MonthlyBlock(brand=brand_label, rows=data))
    return blocks


def ppt8_parse_agg_blocks(df: pd.DataFrame) -> List[Ppt8AggBlock]:
    """Encuentra bloques agregados por marca en columnas 4-6."""
    label_col = df.iloc[:, 4]
    agg_brand_rows = df[
        label_col.notna()
        & (~label_col.astype(str).str.contains("Weighted", case=False, na=False))
        & (~label_col.astype(str).str.contains("table", case=False, na=False))
        & (~label_col.astype(str).str.contains("categ", case=False, na=False))
        & df.iloc[:, 5].isna()
        & df.iloc[:, 6].isna()
    ].index.tolist()
    blocks: List[Ppt8AggBlock] = []
    for i, start in enumerate(agg_brand_rows):
        end = agg_brand_rows[i + 1] if i + 1 < len(agg_brand_rows) else len(df)
        block = df.loc[start + 1 : end - 1]
        block = block[block.iloc[:, 4].notna()]
        if block.empty:
            continue
        brand_label = str(df.iloc[start, 4]).strip()
        if ppt8_is_header_label(brand_label):
            continue
        blocks.append(Ppt8AggBlock(brand=brand_label, block=block))
    return blocks


def ppt8_extract_metrics(block: pd.DataFrame) -> Dict[str, Tuple[float, float]]:
    """Mapea etiqueta -> (MAT base, MAT actual) tolerando alias."""

    def canonical_metric(label: str) -> Optional[str]:
        raw_label = str(label)
        if not raw_label or not raw_label.strip():
            return None
        normalized = raw_label.replace("Weighted", "").strip().upper()
        compact = re.sub(r"[^A-Z0-9]+", "", normalized)
        if not compact:
            return None
        is_vertical = "VERT" in compact

        if is_vertical:
            return None
        if compact in {"RVOL1", "RVOL2", "VOL1", "VOL2", "VOL1P", "VOL2P", "VOLSU", "UNITS", "UNIDAD", "UNIDADES"}:
            return "R_VOL1"
        if compact.startswith("PEN"):
            return "PENET"
        if compact.startswith("FREQ"):
            return "FREQ"
        if compact.startswith("HH"):
            return "HHOLDS"
        return normalized

    metrics: Dict[str, Tuple[float, float]] = {}
    for _, row in block.iterrows():
        canonical = canonical_metric(row.iloc[4])
        if not canonical:
            continue
        try:
            base_val = float(row.iloc[5])
            act_val = float(row.iloc[6])
        except (TypeError, ValueError):
            continue
        # Si solo uno de los dos valores es numerico, replica para no perder la medida
        if not np.isfinite(act_val) and np.isfinite(base_val):
            act_val = base_val
        if not np.isfinite(base_val) and np.isfinite(act_val):
            base_val = act_val
        metrics[canonical] = (base_val, act_val)
    return metrics


def ppt8_find_mat_labels(df: pd.DataFrame) -> Tuple[str, str]:
    """Detecta las etiquetas de MAT (base y actual) en las columnas 5 y 6."""
    mat_base = "MAT Base"
    mat_current = "MAT Actual"
    for _, row in df.iterrows():
        val5 = row.iloc[5] if df.shape[1] > 5 else None
        val6 = row.iloc[6] if df.shape[1] > 6 else None
        if isinstance(val5, str) and "MAT" in val5.upper():
            mat_base = val5.strip()
        if isinstance(val6, str) and "MAT" in val6.upper():
            mat_current = val6.strip()
        if "MAT" in str(val5).upper() or "MAT" in str(val6).upper():
            if mat_base != "MAT Base" and mat_current != "MAT Actual":
                break
    return mat_base, mat_current


def ppt8_classify_penetration(p: float) -> str:
    """Clasifica la penetracion (0-1)."""
    if p <= 0:
        return "Fuera de rango"
    if 0.01 <= p <= 0.10:
        return "Marca Pequeña (1% a 10%)"
    if 0.10 < p <= 0.30:
        return "Marca Mediana (11% a 30%)"
    if 0.30 < p <= 0.50:
        return "Marca Grande (31% a 50%)"
    if 0.50 < p <= 0.99:
        return "Super Marca (51% a 99%)"
    return "Fuera de rango"


def ppt8_compute_rows(
    monthly_blocks: Iterable[Ppt8MonthlyBlock], agg_blocks: Iterable[Ppt8AggBlock]
) -> Tuple[List[Ppt8Row], List[str]]:
    """Combina bloques mensuales y agregados para construir filas del cuadro PPT."""
    rows: List[Ppt8Row] = []
    errors: List[str] = []
    required_metrics = {"R_VOL1", "PENET", "FREQ", "HHOLDS"}
    agg_lookup = {ppt8_normalize_label(block.brand): block for block in agg_blocks}

    for idx, mblock in enumerate(monthly_blocks):
        try:
            ablock = agg_lookup.get(ppt8_normalize_label(mblock.brand))
            if ablock is None:
                ablock = agg_blocks[idx] if idx < len(agg_blocks) else None
            if ablock is None:
                raise ValueError(f"{mblock.brand}: no se encontro bloque agregado correspondiente")
            monthly = mblock.rows
            if monthly.empty:
                raise ValueError(f"{mblock.brand}: sin datos mensuales validos")
            last12_pen = pd.to_numeric(monthly.iloc[:, 1], errors="coerce").tail(12)
            last12_buy = pd.to_numeric(monthly.iloc[:, 2], errors="coerce").tail(12)
            if last12_pen.dropna().empty or last12_buy.dropna().empty:
                raise ValueError(f"{mblock.brand}: faltan valores para promedio 12m")
            pen_avg = float(last12_pen.mean())
            buyers_avg = float(last12_buy.mean())

            metrics = ppt8_extract_metrics(ablock.block)
            missing = [m for m in required_metrics if m not in metrics]
            if missing:
                raise ValueError(f"{mblock.brand}: faltan metricas {', '.join(missing)}")
            vol1, vol2 = metrics["R_VOL1"]
            pen1, pen2 = metrics["PENET"]
            freq1, freq2 = metrics["FREQ"]
            hh_tuple = metrics.get("HHOLDS", (np.nan, np.nan))
            hh = hh_tuple[1]
            if not np.isfinite(hh):
                hh = hh_tuple[0]

            for name, value in {
                "vol1": vol1,
                "vol2": vol2,
                "pen1": pen1,
                "pen2": pen2,
                "freq1": freq1,
                "freq2": freq2,
                "hh": hh,
            }.items():
                if value is None or not np.isfinite(value):
                    raise ValueError(f"{mblock.brand}: valor invalido en {name}")
            if hh <= 0:
                raise ValueError(f"{mblock.brand}: HHOLDS no es positivo")
            if vol1 == 0:
                raise ValueError(f"{mblock.brand}: MAT base es cero")

            dif = vol2 - vol1
            b = (pen1 + pen2) / 200
            w = (freq1 + freq2) / 2
            m_val = b * w
            if m_val <= 0:
                raise ValueError(f"{mblock.brand}: valor m invalido")
            se_rel = math.sqrt(2 / (m_val * hh))
            avg_vol = (vol1 + vol2) / 2
            int_dif = PPT8_Z_SCORE_INTERVAL * se_rel * avg_vol
            lim_dif_minus = dif - int_dif
            lim_dif_plus = dif + int_dif
            lim_cant2_minus = vol1 + lim_dif_minus
            lim_cant2_plus = vol1 + lim_dif_plus

            evol = (vol2 / vol1 - 1) * 100
            lim_evol_minus = (lim_cant2_minus / vol1 - 1) * 100
            lim_evol_plus = (lim_cant2_plus / vol1 - 1) * 100

            p_ratio = pen_avg / 100
            if p_ratio <= 0:
                err = np.nan
            else:
                err = PPT8_Z_SCORE_ERROR * math.sqrt((p_ratio * (1 - p_ratio)) / hh) / p_ratio
            nivel = ppt8_classify_penetration(p_ratio)

            rows.append(
                Ppt8Row(
                    marca=mblock.brand,
                    mat_sep25=evol,
                    lim_inf=lim_evol_minus,
                    lim_sup=lim_evol_plus,
                    compradores_prom=buyers_avg,
                    pen12m=pen_avg,
                    nivel_pen=nivel,
                    error_muestral=err,
                )
            )
        except ValueError as exc:
            errors.append(str(exc))
    return rows, errors


def render_ppt8_table(
    rows: List[Ppt8Row],
    figure_id: int,
    brand_color_lookup: Dict[str, str],
    mat_label_current: str,
) -> Tuple[io.BytesIO, Tuple[float, float]]:
    """Renderiza la tabla PPT del segmento 8 y devuelve un stream PNG."""
    COLOR_HEADER = "#1F4E78"
    COLOR_B = "#DCE6F1"
    COLOR_C = "#F2DCDB"
    COLOR_D = "#E2EFDA"
    COLOR_G = "#DCE6F1"
    COLOR_H = "#E2EFDA"
    COLOR_ERR_GREEN = "#A8D08D"
    COLOR_ERR_YELLOW = "#FFD966"
    COLOR_ERR_RED = "#C00000"

    def _interp_color(start_hex: str, end_hex: str, t: float) -> str:
        t = max(0.0, min(1.0, t))
        start_rgb = mcolors.to_rgb(start_hex)
        end_rgb = mcolors.to_rgb(end_hex)
        blended = tuple(s + (e - s) * t for s, e in zip(start_rgb, end_rgb))
        return mcolors.to_hex(blended)

    def _error_cell_color(err_ratio: float) -> str:
        """Color degradado según error muestral (%)."""
        if err_ratio is None or not np.isfinite(err_ratio):
            return COLOR_H
        val = err_ratio * 100.0
        if val <= 1.0:
            return COLOR_ERR_GREEN
        if val <= 5.0:
            t = (val - 1.0) / 4.0
            return _interp_color(COLOR_ERR_GREEN, COLOR_ERR_YELLOW, t)
        if val <= 51.0:
            t = (val - 5.0) / 46.0
            return _interp_color(COLOR_ERR_YELLOW, COLOR_ERR_RED, t)
        return COLOR_ERR_RED

    def _text_color_for_bg(hex_color: str) -> str:
        r, g, b = mcolors.to_rgb(hex_color)
        if 0.299 * r + 0.587 * g + 0.114 * b < 0.45:
            return "#FFFFFF"
        return "#1F1F1F"

    header_bottom = [
        "Marca",
        f"% VAR {mat_label_current}" if mat_label_current else "% VAR MAT",
        "% VAR LIM INFERIOR",
        "% VAR LIM SUPERIOR",
        "Nivel de penetracion",
        "Error muestral",
    ]

    data = [
        [
            r.marca,
            f"{r.mat_sep25:.1f}" if np.isfinite(r.mat_sep25) else "-",
            f"{r.lim_inf:.1f}" if np.isfinite(r.lim_inf) else "-",
            f"{r.lim_sup:.1f}" if np.isfinite(r.lim_sup) else "-",
            r.nivel_pen,
            f"{r.error_muestral:.1%}" if np.isfinite(r.error_muestral) else "-",
        ]
        for r in rows
    ]

    fig_height = 1.4 + 0.55 * len(rows)
    fig, ax = plt.subplots(num=figure_id, figsize=(12, fig_height), dpi=TABLE_EXPORT_DPI)
    ax.axis("off")

    table = ax.table(
        cellText=[header_bottom] + data,
        colLabels=None,
        cellLoc="center",
        loc="center",
    )

    table.auto_set_font_size(False)
    table.set_fontsize(13)
    table.scale(1.0, 1.25)

    widths = [0.24, 0.14, 0.14, 0.14, 0.24, 0.14]
    for i, w in enumerate(widths):
        table.auto_set_column_width(i)
        table._cells[(0, i)].set_width(w)
    for (r_idx, c_idx), cell in table._cells.items():
        cell.set_edgecolor("black")
        cell.set_linewidth(1.5 if r_idx == 0 else 1.2)
        cell.PAD = 0.08 if r_idx == 0 else 0.06

    for col in range(len(header_bottom)):
        cell = table[0, col]
        cell.set_facecolor(COLOR_HEADER)
        cell.set_text_props(color="white", weight="bold")

    start_body = 1
    for r_idx in range(start_body, len(data) + start_body):
        table[r_idx, 1].set_facecolor(COLOR_B)
        table[r_idx, 2].set_facecolor(COLOR_C)
        table[r_idx, 3].set_facecolor(COLOR_D)
        table[r_idx, 4].set_facecolor(COLOR_G)
        err_val = rows[r_idx - start_body].error_muestral
        err_color = _error_cell_color(err_val)
        table[r_idx, 5].set_facecolor(err_color)
        table[r_idx, 5].set_text_props(color=_text_color_for_bg(err_color))
        brand_label = rows[r_idx - start_body].marca
        brand_color = assign_brand_palette_color(brand_label, brand_color_lookup, TREND_COLOR_SEQUENCE)
        if brand_color:
            table[r_idx, 0].set_text_props(color=brand_color, weight="normal")

    plt.tight_layout()
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    table_bbox = table.get_tightbbox(renderer).transformed(fig.dpi_scale_trans.inverted())
    fig_size = fig.get_size_inches()
    bbox_width = float(table_bbox.width) if table_bbox.width else float(fig_size[0])
    bbox_height = float(table_bbox.height) if table_bbox.height else float(fig_size[1])
    return (
        figure_to_stream(
            fig,
            dpi=TABLE_EXPORT_DPI,
            bbox_inches=table_bbox,
            pad_inches=0.02,
            transparent=False,
        ),
        (bbox_width, bbox_height),
    )

pd.set_option('future.no_silent_downcasting', True)
pd.set_option('mode.chained_assignment', None)
warnings.filterwarnings('ignore')
# Segmento 1: utilidades para el grafico MAT (cuando).
# Funcion que prepara los datos para la creacion del grafico MAT
def df_mat(df,p):
    """Prepara el dataframe MAT con acumulados y variaciones."""
    v1 = pd.DataFrame([(df.iloc[i-12-p:i-p,1].sum()/df.iloc[i-24-p:i-12-p,1].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-in'])
    v2 = pd.DataFrame([(df.iloc[i-12:i,2].sum()/df.iloc[i-24:i-12,2].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-out'])
    ac1, ac2 = [pd.DataFrame(df[col].rolling(window=12).sum().dropna()) for col in df.columns[1:]]
    v1.index,v2.index=v1.index+11,v2.index+11
    v1 ,v2 = v1.dropna(), v2.dropna()
    mat =pd.DataFrame(df.iloc[:,0].copy())
    mat= mat.join(ac1.join(ac2)[-len(v1):].join([v1]).join([v2]))
    mat= mat[mat.notna().all(axis=1).idxmax():]
    return mat
#Grafico de MAT 
def graf_mat(mat, c_fig, p):
    """Genera el grafico MAT y devuelve una imagen en memoria."""
    altura = 7  # Aumentamos aun mas la altura
    fig, ax = plt.subplots(num=c_fig, figsize=(15, altura), dpi=DEFAULT_EXPORT_DPI)
    fig.subplots_adjust(bottom=0.2, top=0.9)
    ran = [x.strftime('%m-%y') for x in mat.iloc[:, 0].copy()]
    ac1 = mat.iloc[:, 1].copy()
    ac2 = mat.iloc[:, 2].copy()
    v1 = mat.iloc[:, 3].copy()
    v2 = mat.iloc[:, 4].copy()
    l1 = ax.plot(ran, ac1, color=LINE_COLOR_CLIENT, linewidth=4.2, label='Acumulado Cliente')
    l2 = ax.plot(ran, ac2, color=LINE_COLOR_NUMERATOR, linewidth=4.2, label='Acumulado Numerator')
    ax2 = ax.twinx()
    b1 = ax2.bar(
        np.arange(len(ran)) - 0.3,
        v1.values,
        0.3,
        color=BAR_COLOR_VARIATION_CLIENT,
        edgecolor=BAR_EDGE_COLOR,
        label='Var % Cliente Pipeline: ' + str(p) + ' ' + labels[(lang, 'Var MAT')],
    )
    b2 = ax2.bar(
        np.arange(len(ran)) + 0.3,
        v2.values,
        0.3,
        color=BAR_COLOR_VARIATION_NUMERATOR,
        edgecolor=BAR_EDGE_COLOR,
        label='Var % Numerator ' + labels[(lang, 'Var MAT')],
    )
    ax.set_xticks(np.arange(len(ran)))
    ax.set_xticklabels(ran, rotation=30)
    ax2.tick_params(
        left=False,
        labelleft=False,
        top=False,
        labeltop=False,
        right=False,
        labelright=False,
        bottom=False,
        labelbottom=False,
    )
    bar_positions = np.arange(len(ran))
    label_specs = [
        (v1, -0.3, BAR_LABEL_COLOR_POS, BAR_LABEL_COLOR_NEG),
        (v2, 0.3, BAR_LABEL_COLOR_POS_ALT, BAR_LABEL_COLOR_NEG_ALT),
    ]
    for serie, x_offset, color_pos, color_neg in label_specs:
        for idx, y in enumerate(serie):
            if not np.isfinite(y):
                continue
            y_pos = y + 0.01 if y >= 0 else y - 0.01  # leve desplazamiento para dejar espacio
            va_align = 'bottom' if y >= 0 else 'top'
            ax2.text(
                bar_positions[idx] + x_offset,
                y_pos,
                f"{y*100:.1f}%",
                ha='center',
                va=va_align,
                fontsize=8,
                fontweight='bold',
                color=color_pos if y >= 0 else color_neg,
            )
    y2_min, y2_max = ax2.get_ylim()
    padding = max(abs(y2_min), abs(y2_max)) * 0.15 if np.isfinite(y2_min) and np.isfinite(y2_max) else 0
    ax2.set_ylim(y2_min - padding, y2_max + padding * 2)
    lns = l1 + l2 + [b1, b2]
    labs = [l.get_label() for l in l1 + l2] + [b1.get_label(), b2.get_label()]
    ax.legend(
        lns,
        labs,
        loc='upper center',
        bbox_to_anchor=(0.5, -0.15),
        borderaxespad=0.05,
        frameon=True,
        prop={'size': 12},
        ncol=2,
    )
    ax.set_title(labels[(lang, 'MAT')] + ' | ' + w[2:], size=18, pad=20)
    fig.tight_layout(rect=[0, 0.05, 1, 0.95])
    return figure_to_stream(fig, bbox_inches=None)
#Grafico de Lineas

def format_title_with_bold_suffix(base_title: str, suffix: str) -> str:
    """Devuelve un titulo con sufijo en negrita usando mathtext."""
    base_title = (base_title or '').strip()
    suffix = (suffix or '').strip()
    if not suffix:
        return base_title
    safe_suffix = suffix.replace('\\', '\\\\').replace('$', '\\$')
    spacer = ' ' if base_title else ''
    return f"{base_title}{spacer}$\\bf{{{safe_suffix}}}$"

def line_graf(
    df,
    p,
    title,
    c_fig,
    ven,
    width_emu=None,
    height_emu=None,
    multi_chart=None,
    share_lookup=None,
    y_axis_percent=False,
    top_value_annotations=0,
    color_collector=None,
    color_overrides=None,
    linestyle_overrides=None,
    show_title=True,
    include_origin_suffix=True,
    force_right_labels=False,
):
    """
    Renderiza la gráfica de tendencias y devuelve un buffer PNG listo para insertar en la slide.
    Cuando se generan múltiples gráficos en la misma diapositiva (multi_chart=True), se aplican
    márgenes más amplios para evitar recortes visuales.
    """
    if width_emu is not None:
        try:
            planned_width = emu_to_inches(width_emu)
        except Exception:
            planned_width = None
    else:
        planned_width = None
    if multi_chart is None:
        detected_multi = planned_width is not None and planned_width < 6.0
    else:
        detected_multi = bool(multi_chart)
    if detected_multi:
        chart_top_margin = LINE_CHART_MULTI_TOP_MARGIN
        chart_bottom_margin = LINE_CHART_MULTI_BOTTOM_MARGIN
        chart_left_margin = LINE_CHART_MULTI_LEFT_MARGIN
        chart_right_margin = LINE_CHART_MULTI_RIGHT_MARGIN
        chart_x_margin = LINE_CHART_MULTI_X_MARGIN
    else:
        chart_top_margin = LINE_CHART_TOP_MARGIN
        chart_bottom_margin = LINE_CHART_BOTTOM_MARGIN
        chart_left_margin = LINE_CHART_LEFT_MARGIN
        chart_right_margin = LINE_CHART_RIGHT_MARGIN
        chart_x_margin = LINE_CHART_SINGLE_X_MARGIN
    # Ajusta o tamanho da figura para combinar com o espaco reservado no slide
    width_inches = emu_to_inches(width_emu) if width_emu is not None else None
    height_inches = emu_to_inches(height_emu) if height_emu is not None else None
    if width_inches is None and height_inches is None:
        # Ancho dinámico según cantidad de datos (en pulgadas)
        n_points = len(df)
        ancho = max(15, n_points * 0.5)  # 0.5 puede ajustarse para más/menos espacio (pulgadas)
        # Altura fija de 10 cm -> convertir a pulgadas
        altura = emu_to_inches(Cm(10))
        figsize = (ancho, altura)
    else:
        if width_inches is None:
            width_inches = height_inches * DEFAULT_LINE_CHART_RATIO
        elif height_inches is None:
            height_inches = width_inches / DEFAULT_LINE_CHART_RATIO
        figsize = (width_inches, height_inches)
    fig, ax = plt.subplots(num=c_fig, figsize=figsize, dpi=DEFAULT_EXPORT_DPI)
    # Escalar elementos (fuentes, linewidth) según la altura para que el contenido se adapte
    # Se toma como referencia la altura original de 5 pulgadas usada previamente
    ref_height = 5.0
    actual_height = figsize[1]
    scale = actual_height / ref_height if ref_height else 1.0
    base_linewidth = 2.0 * max(0.6, scale)
    title_base_size = max(10, int(18 * scale))
    legend_base_size = max(8, int(10 * scale))
    xtick_size = max(6, int(8 * scale))
    xtick_font = FontProperties(family='DejaVu Sans', size=xtick_size, weight='regular')
    xtick_pad = max(6, int(8 * scale))
    marker_size = max(3.5, 4.5 * scale)
    share_lookup = share_lookup or {}
    aux = df.copy()
    if aux.shape[1] > 1:
        kept_columns = [aux.columns[0]]
        for col in aux.columns[1:]:
            series = pd.to_numeric(aux[col], errors='coerce')
            if series.replace(0, np.nan).dropna().empty:
                continue
            kept_columns.append(col)
        if len(kept_columns) == 1:
            aux = aux.iloc[:, :1]
        else:
            aux = aux.loc[:, kept_columns]
    if aux.empty or aux.shape[1] <= 1:
        return figure_to_stream(fig, bbox_inches=None)
    ran = [x.strftime('%m-%y') for x in aux.iloc[:, 0]]
    data_len = len(ran)
    if data_len == 0:
        return figure_to_stream(fig, bbox_inches=None)
    start_idx = max(0, min(p, data_len - 1))
    x_positions = np.arange(data_len)
    colunas = list(aux.columns[1:])
    palette_values = TREND_COLOR_SEQUENCE or [mcolors.to_hex(c) for c in plt.get_cmap('tab20').colors]
    color_mapping: dict[str, str] = {}
    color_key_lookup: dict[str, str] = {}
    def _register_color_keys(label: str, color_value: str, key_list: Optional[list[str]] = None) -> None:
        keys = key_list if key_list is not None else generate_color_lookup_keys(label)
        for key in keys:
            if key and key not in color_key_lookup:
                color_key_lookup[key] = color_value
    if color_overrides:
        for name, hex_color in color_overrides.items():
            color_mapping[name] = hex_color
            _register_color_keys(name, hex_color)
    color_index = 0
    def _resolve_or_assign_color(label: str) -> str:
        nonlocal color_index
        existing_color = color_mapping.get(label)
        if existing_color:
            return existing_color
        keys = generate_color_lookup_keys(label)
        for key in keys:
            reused_color = color_key_lookup.get(key)
            if reused_color:
                color_mapping[label] = reused_color
                _register_color_keys(label, reused_color, keys)
                return reused_color
        assigned_color = palette_values[color_index % len(palette_values)]
        color_index += 1
        color_mapping[label] = assigned_color
        _register_color_keys(label, assigned_color, keys)
        return assigned_color
    lns = []
    legend_labels = []
    numeric_series_list = []
    plotted_columns = []
    series_points = {}
    for col in colunas:
        if ven > 1:
            estilo = '-' if '.v' in col.lower() else '--'
        else:
            estilo = '-'
        cor = _resolve_or_assign_color(col)
        numeric_series = pd.to_numeric(aux[col], errors='coerce')
        numeric_series_list.append(numeric_series)
        y = numeric_series.values
        x_slice = x_positions[start_idx:]
        y_slice = y[start_idx:]
        valid_points = [(x_idx, y_val) for x_idx, y_val in zip(x_slice, y_slice) if pd.notna(y_val)]
        if not valid_points:
            continue
        series_points[col] = valid_points
        suffix_map = {".c": "Compras", ".v": "Ventas", "_c": "Compras", "_v": "Ventas", "-c": "Compras", "-v": "Ventas"}
        lower_col = col.lower()
        legend_label = col
        for suffix, suffix_text in suffix_map.items():
            if lower_col.endswith(suffix):
                base_label = col[:-len(suffix)].strip()
                if base_label:
                    legend_label = f"{base_label} ({suffix_text})" if include_origin_suffix else base_label
                break
        if linestyle_overrides and col in linestyle_overrides:
            estilo = linestyle_overrides[col]
        line, = ax.plot(
            x_slice,
            y_slice,
            color=cor,
            linewidth=base_linewidth,
            linestyle=estilo,
            label=col,
            marker='o',
            markersize=marker_size,
            markerfacecolor='white',
            markeredgewidth=max(0.7, base_linewidth / 2),
            markeredgecolor=cor,
        )
        lns.append(line)
        legend_labels.append(legend_label)
        plotted_columns.append(col)
    if data_len and start_idx < data_len:
        ax.set_xticks(x_positions[start_idx:])
        ax.set_xticklabels(
            ran[start_idx:],
            rotation=90,
            fontproperties=xtick_font,
            ha='center',
            va='center',
        )
        ax.tick_params(axis='x', pad=xtick_pad)
        ax.set_xlim(x_positions[start_idx], x_positions[-1])
    else:
        ax.set_xticks([])
        ax.set_xticklabels([])
    y_tick_size = max(8, int(10 * scale))
    ax.tick_params(axis='y', labelsize=y_tick_size)
    axis_percent = bool(y_axis_percent)
    max_abs_value = 0.0
    for series in numeric_series_list:
        numeric_values = series.to_numpy()
        if numeric_values.size == 0:
            continue
        finite_values = numeric_values[np.isfinite(numeric_values)]
        if finite_values.size == 0:
            continue
        series_max = float(np.max(np.abs(finite_values)))
        if series_max > max_abs_value:
            max_abs_value = series_max
    if axis_percent:
        def _trend_tick_formatter(value, _):
            return f"{value:.1f}%"
        ax.yaxis.set_major_formatter(FuncFormatter(_trend_tick_formatter))
    else:
        scale_rule = select_trend_scale_rule(max_abs_value)
        divisor = scale_rule["divisor"] or 1
        suffix = scale_rule["suffix"]
        def _trend_tick_formatter(value, _):
            return format_value_with_suffix(value, divisor, suffix)
        ax.yaxis.set_major_formatter(FuncFormatter(_trend_tick_formatter))
        if DISPLAY_TREND_REFERENCE_TEXT and suffix:
            reference_template = scale_rule.get("reference", {})
            reference_text = reference_template.get(lang, reference_template.get('default', '')).strip()
            if reference_text:
                reference_font_size = max(8, int(9 * scale))
                ax.text(
                    0.02,
                    0.95,
                    reference_text,
                    transform=ax.transAxes,
                    ha='left',
                    va='top',
                    fontsize=reference_font_size,
                    color='#4D4D4D',
                    alpha=0.85,
                )
    if axis_percent:
        def _format_axis_value(value):
            return f"{value:.1f}%"
    else:
        def _format_axis_value(value):
            return format_value_with_suffix(value, divisor, suffix)
    ax.set_ylim(bottom=0)
    # Margen de los datos dentro del eje
    ax.margins(x=chart_x_margin, y=0.08)

    # --- Parámetros de separación en coordenadas de figura ---
    # para gráficos "sencillos" vs "multi" (muchas series)
    base_gap = 0.075 if not detected_multi else 0.070         # espacio visual típico
    min_axes_legend_gap = 0.070 if not detected_multi else 0.080  # mínimo absoluto entre eje y leyenda
    margin_buffer = 0.020                                      # margen al borde inferior de la figura

    legend_bottom_margin = chart_bottom_margin
    legend_columns = 1
    legend_rows = 1
    legend_height_fraction = 0.0

    if lns:
        # Cálculo de columnas/filas en la leyenda
        max_columns = 3 if detected_multi else 4
        legend_columns = max(1, min(len(legend_labels), max_columns))
        legend_rows = max(1, math.ceil(len(legend_labels) / legend_columns))

        # Altura de la leyenda en pulgadas -> fracción de la figura
        legend_font_points = legend_base_size
        legend_line_height_points = legend_font_points * 1.35
        legend_height_inches = (legend_line_height_points / 72.0) * legend_rows
        figure_height_inches = fig.get_size_inches()[1]
        legend_height_fraction = (
            legend_height_inches / figure_height_inches if figure_height_inches else 0.0
        )

        # Espacio mínimo entre eje y leyenda
        legend_clearance = max(base_gap, min_axes_legend_gap)

        # Aseguramos espacio suficiente abajo para leyenda + gap + buffer
        legend_bottom_margin = max(
            chart_bottom_margin,
            legend_height_fraction + legend_clearance + margin_buffer,
        )
    else:
        legend_clearance = 0.0  # no hay leyenda, no hace falta espacio extra

    # --- Márgenes superior/inferior del eje ---
    effective_top_margin = chart_top_margin
    if not (show_title and title):
        # Si no hay título, no "regalemos" tanto espacio arriba
        effective_top_margin = min(0.97, chart_top_margin + 0.05)
    elif title:
        plt.title(title, size=title_base_size, pad=10)

    bottom_margin = min(0.9, max(0.05, legend_bottom_margin))
    if bottom_margin >= effective_top_margin:
        bottom_margin = max(0.05, effective_top_margin - 0.05)

    fig.subplots_adjust(
        top=effective_top_margin,
        bottom=bottom_margin,
        left=chart_left_margin,
        right=chart_right_margin,
    )

    # --- Posición final de la leyenda ---
    if lns:
        axes_box = ax.get_position()  # caja del eje en coordenadas de figura

        # Queremos que la parte superior de la leyenda quede legend_clearance
        # por debajo del borde inferior del eje.
        legend_top = axes_box.y0 - legend_clearance

        legend = ax.legend(
            lns,
            legend_labels,
            loc="upper left",
            bbox_to_anchor=(axes_box.x0, legend_top, axes_box.width, 0.0),
            bbox_transform=fig.transFigure,
            borderaxespad=0.0,
            frameon=True,
            prop={"size": legend_base_size},
            ncol=legend_columns,
            mode="expand",
        )

        frame = legend.get_frame()
        frame.set_facecolor("white")
        frame.set_edgecolor("#D3D3D3")
        frame.set_alpha(0.85)

    if plotted_columns and (share_lookup or force_right_labels):
        legend_label_map = {col: label for col, label in zip(plotted_columns, legend_labels)}
        annotation_candidates = []
        share_lookup_lower = {str(key).lower(): value for key, value in share_lookup.items()} if share_lookup else {}
        for col, line in zip(plotted_columns, lns):
            points = series_points.get(col)
            if not points:
                continue
            last_x, last_y = points[-1]
            share_value = share_lookup.get(col) if share_lookup else None
            if share_value is None and share_lookup_lower:
                share_value = share_lookup_lower.get(str(col).lower())
            annotation_candidates.append({
                "column": col,
                "line": line,
                "legend": legend_label_map.get(col, col),
                "share": share_value,
                "last_x": last_x,
                "last_y": last_y,
                "color": line.get_color(),
            })
        total_series = len(annotation_candidates)
        if total_series:
            if total_series < 10:
                selected_series = annotation_candidates
            else:
                with_share = [info for info in annotation_candidates if info["share"] is not None]
                without_share = [info for info in annotation_candidates if info["share"] is None]
                with_share.sort(key=lambda info: info["share"], reverse=True)
                if force_right_labels:
                    without_share.sort(key=lambda info: info["last_y"], reverse=True)
                    selected_series = []
                    if with_share:
                        selected_series.extend(with_share[:3])
                    if without_share:
                        selected_series.extend(without_share[:2])
                    if len(selected_series) < 5:
                        remaining = [info for info in annotation_candidates if info not in selected_series]
                        remaining.sort(key=lambda info: info["last_y"], reverse=True)
                        selected_series.extend(remaining[: max(0, 5 - len(selected_series))])
                else:
                    selected_series = with_share[:5]
                    if len(selected_series) < 5:
                        remaining = [info for info in annotation_candidates if info not in selected_series]
                        selected_series.extend(remaining[: max(0, 5 - len(selected_series))])
            if selected_series:
                x_limits = ax.get_xlim()
                y_limits = ax.get_ylim()
                y_range = y_limits[1] - y_limits[0] if y_limits[1] > y_limits[0] else 1.0
                min_gap = y_range * 0.04
                min_margin = y_range * 0.02
                x_range = x_limits[1] - x_limits[0] if x_limits[1] > x_limits[0] else 1.0
                x_offset = max(0.8, x_range * 0.06)
                arranged_series = sorted(selected_series, key=lambda info: info["last_y"], reverse=True)
                placed_positions = []
                y_max = y_limits[1] - min_margin
                y_min = y_limits[0] + min_margin
                for info in arranged_series:
                    y_pos = max(min(info["last_y"], y_max), y_min)
                    for prev_y in placed_positions:
                        if y_pos > prev_y - min_gap:
                            y_pos = prev_y - min_gap
                    y_pos = max(y_pos, y_min)
                    placed_positions.append(y_pos)
                    info["label_y"] = y_pos
                    info["label_x"] = info["last_x"] + x_offset
                max_label_x = max(info["label_x"] for info in arranged_series)
                if max_label_x > x_limits[1]:
                    ax.set_xlim(x_limits[0], max_label_x + x_offset * 1.15)
                size_factor = 0.7 if detected_multi else 0.8
                min_size = 6 if detected_multi else 7
                label_font_size = max(min_size, int(legend_base_size * size_factor))
                bbox_props = dict(facecolor='white', edgecolor='none', alpha=0.85, boxstyle='round,pad=0.2')
                for info in arranged_series:
                    text_value = info["legend"]
                    ax.text(
                        info["label_x"],
                        info["label_y"],
                        text_value,
                        ha='left',
                        va='center',
                        color=info["color"],
                        fontsize=label_font_size,
                        fontweight='bold',
                        bbox=bbox_props,
                        clip_on=False,
                    )
    if top_value_annotations and plotted_columns and lns:
        point_candidates = []
        line_lookup = {col: line for col, line in zip(plotted_columns, lns)}
        for col in plotted_columns:
            if str(col).strip().lower() == 'total':
                continue
            line = line_lookup.get(col)
            if line is None:
                continue
            for x_idx, y_val in series_points.get(col, []):
                if not np.isfinite(y_val):
                    continue
                point_candidates.append({
                    "column": col,
                    "line": line,
                    "x": x_idx,
                    "y": float(y_val),
                })
        if point_candidates:
            point_candidates.sort(key=lambda item: item["y"], reverse=True)
            y_limits = ax.get_ylim()
            y_range = y_limits[1] - y_limits[0] if y_limits[1] > y_limits[0] else 1.0
            label_offset = max(y_range * 0.03, 0.5 if axis_percent else y_range * 0.03)
            label_font_size = max(8, int(10 * scale))
            annotated = 0
            seen_points = set()
            for point in point_candidates:
                point_key = (point["x"], round(point["y"], 6))
                if point_key in seen_points:
                    continue
                seen_points.add(point_key)
                label_y = point["y"] + label_offset
                if label_y > y_limits[1]:
                    label_y = min(y_limits[1], point["y"] + label_offset * 0.6)
                text_value = _format_axis_value(point["y"])
                ax.annotate(
                    text_value,
                    xy=(point["x"], point["y"]),
                    xytext=(point["x"], label_y),
                    textcoords='data',
                    ha='center',
                    va='bottom',
                    fontsize=label_font_size,
                    color=point["line"].get_color(),
                    fontweight='bold',
                    bbox=dict(
                        boxstyle='round,pad=0.2',
                        facecolor='white',
                        edgecolor=point["line"].get_color(),
                        linewidth=0.6,
                        alpha=0.85
                    ),
                    arrowprops=dict(
                        arrowstyle='-',
                        color=point["line"].get_color(),
                        linewidth=max(0.6, base_linewidth * 0.5)
                    ),
                    clip_on=False,
                )
                annotated += 1
                if annotated >= top_value_annotations:
                    break
    if isinstance(color_collector, dict):
        color_collector.clear()
        for col, line in zip(plotted_columns, lns):
            color_collector[col] = line.get_color()
    return figure_to_stream(fig)
#Normaliza etiquetas de periodo eliminando prefijos genericos
def normalize_period_label(label) -> str:
    """Normaliza etiquetas de periodo eliminando prefijos genericos."""
    if label is None:
        return ""
    text = str(label).strip()
    if not text:
        return text
    parts = text.split(None, 1)
    if parts and parts[0].lower().rstrip('.') in {'vol', 'volume', 'volumen'} and len(parts) > 1:
        return parts[1].strip()
    return text
#Grafico de barras apiladas 100% para share por periodo
def stacked_share_chart(period_label, share_values, color_mapping, c_fig, title=None):
    """
    Genera un gráfico 100% apilado a partir de la participación de un periodo puntual.
    """
    if not isinstance(share_values, (dict, OrderedDict)):
        share_values = {}
    palette_values = TREND_COLOR_SEQUENCE or [mcolors.to_hex(c) for c in plt.get_cmap('tab20').colors]
    normalized_items = []
    total = 0.0
    for key, value in share_values.items():
        try:
            numeric_value = float(value)
        except (TypeError, ValueError):
            try:
                numeric_value = float(str(value).replace('%', '').replace(',', '.')) / 100.0
            except Exception:
                numeric_value = 0.0
        if not np.isfinite(numeric_value) or numeric_value < 0:
            numeric_value = 0.0
        total += numeric_value
        normalized_items.append([key, numeric_value])
    if total <= 0:
        total = 1.0
    for item in normalized_items:
        item[1] = item[1] / total
    normalized_items = [item for item in normalized_items if item[1] > 0]
    segments = len(normalized_items)
    base_height = 6.2
    incremental_height = 0.6
    figure_height = max(base_height, 4.5 + segments * incremental_height)
    figure_height = min(7.0, figure_height)
    figure_width = 4.6
    fig, ax = plt.subplots(num=c_fig, figsize=(figure_width, figure_height), dpi=DEFAULT_EXPORT_DPI)
    bottom = 0.0
    bar_width = 0.65
    label_infos = []
    for idx, (key, value) in enumerate(normalized_items):
        color = color_mapping.get(key)
        if color is None:
            color = palette_values[idx % len(palette_values)]
            color_mapping[key] = color
        ax.bar(0, value, width=bar_width, bottom=bottom, color=color, edgecolor='white')
        center_y = bottom + value / 2 if value > 0 else bottom
        share_label = f"{value * 100:.1f}%"
        label_infos.append(
            {
                "key": key,
                "text": f"{key} ({share_label})",
                "color": color,
                "center_y": center_y,
            }
        )
        bottom += value
    if label_infos:
        label_infos.sort(key=lambda info: info["center_y"])
        min_gap = max(0.035, min(0.12, 0.95 / max(segments, 1)))
        min_y = 0.04
        max_y = 0.96
        adjusted_positions = []
        for info in label_infos:
            proposed = np.clip(info["center_y"], min_y, max_y)
            if adjusted_positions and proposed - adjusted_positions[-1] < min_gap:
                proposed = min(max_y, adjusted_positions[-1] + min_gap)
            adjusted_positions.append(proposed)
        for idx in range(len(adjusted_positions) - 2, -1, -1):
            if adjusted_positions[idx + 1] - adjusted_positions[idx] < min_gap:
                adjusted_positions[idx] = max(min_y, adjusted_positions[idx + 1] - min_gap)
        for info, position in zip(label_infos, adjusted_positions):
            info["label_y"] = position
        max_text_len = max(len(info["text"]) for info in label_infos)
        label_offset = 1.1 + max(0, (max_text_len - 22) * 0.014)
        label_offset = min(label_offset, 1.8)
        label_x = bar_width / 2 + label_offset
    else:
        label_x = bar_width / 2 + 1.15
    connector_x = bar_width / 2 + 0.1
    label_font_size = 9 if segments <= 10 else (8 if segments <= 16 else 7)
    if label_infos:
        for info in label_infos:
            ax.plot(
                [connector_x, label_x - 0.1],
                [info["center_y"], info["label_y"]],
                color=info["color"],
                linewidth=1.2,
                alpha=0.65,
            )
            ax.text(
                label_x,
                info["label_y"],
                info["text"],
                ha='left',
                va='center',
                color=info["color"],
                fontsize=label_font_size,
                fontweight='bold'
            )
    ax.set_ylim(0, 1)
    axis_right = label_x + 1.0 if label_infos else (bar_width / 2 + 1.4)
    ax.set_xlim(-0.75, axis_right)
    ax.set_xticks([])
    ax.set_ylabel('')
    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: f"{x * 100:.0f}%"))
    ax.set_title(title or period_label, fontsize=15, pad=24, fontweight='bold')
    ax.set_yticks(np.linspace(0, 1, 5))
    ax.yaxis.tick_left()
    ax.grid(axis='y', linestyle='--', alpha=0.25)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.subplots_adjust(left=0.05, right=0.92, top=0.995, bottom=0.02)
    fig_size = fig.get_size_inches()
    return figure_to_stream(fig), fig_size
# Segmento 2: helpers para arbol de medidas (por que).
def _limpiar_tabla_excel(df):
    """Detecta la fila de encabezados (MAT ...) tolerando filas adicionales al inicio."""
    header_idx = None
    for idx, row in df.iterrows():
        non_null = [str(v) for v in row.tolist() if pd.notna(v)]
        periodos = [val for val in non_null[1:] if str(val).startswith("MAT")]
        tiene_table = any("table" in str(val).lower() for val in non_null)
        if periodos or tiene_table:
            header_idx = idx
            break
    if header_idx is None:
        return None, None
    header_row = df.iloc[header_idx].tolist()
    periodos = [str(c).strip() for c in header_row[1:] if pd.notna(c)]
    if len(periodos) < 2:
        return None, None
    tabla = df.iloc[header_idx + 1:, : len(periodos) + 1].copy()
    tabla.columns = ["Metric"] + periodos
    tabla = tabla.dropna(how="all")
    tabla = tabla.dropna(subset=["Metric"])
    tabla["Metric"] = tabla["Metric"].astype(str).str.strip()
    for col in periodos:
        tabla[col] = pd.to_numeric(tabla[col], errors="coerce")
    return tabla.reset_index(drop=True), periodos
def _nombres_unidad(unidad):
    """Devuelve las etiquetas del arbol segun la unidad indicada."""
    unidad_lower = unidad.lower()
    nombre_vol = {
        "litros": ("LTs", "ML"),
        "units": ("Units", "Unit"),
        "kilos": ("KGs", "KG"),
        "toneladas": ("Tons", "Ton"),
        "rollos": ("Rollos", "Rollo"),
        "metros": ("MTs", "MT"),
        "hojas": ("Hojas", "Hoja"),
    }.get(unidad_lower, (unidad, unidad))
    vol_label, sub_label = nombre_vol
    precio_unit = vol_label[:-1] if vol_label.endswith("s") else vol_label
    return {
        "valor": "Valor $ 000s",
        "volumen": f"Volumen 000 {vol_label}",
        "precio": f"Precio Promedio $/{precio_unit}",
        "gasto": "Gasto Promedio $/Buyer",
        "compradores": "Compradores 000s",
        "volumen_prom": f"Volumen Promedio {sub_label}/Comprador",
        "penetracion": "% Penetracion",
        "hholds": "Total HHolds 000s",
        "frecuencia": "Frecuencia",
        "volumen_viaje": f"Volumen por Viaje {sub_label}/Viaje",
        "ticket": "Gasto por Ticket $/Viaje",
    }
def calcular_cambios(df, periodo_inicial, periodo_final, unidad='Units'):
    """Calcula niveles y variaciones porcentuales entre dos periodos."""
    etiquetas = _nombres_unidad(unidad)
    def calculate_percentage_change(old, new):
        if old is None or new is None:
            return 0
        if not np.isfinite(old) or not np.isfinite(new) or old == 0:
            return 0
        return ((new - old) / old) * 100
    factores = {
        'Units': 1,
        'Litros': 1,
        # La base viene en kilos: en kilos no se escala, en toneladas se divide por 1000.
        'Kilos': 1,
        # Para toneladas la base viene en kilos, se divide por 1000 en volumen
        # y se multiplica por 1000 en precios; factor < 1 corrige el inflado observado.
        'Toneladas': 0.001,
        'Rollos': 10,
        'Metros': 100,
        'Hojas': 1
    }
    factor = factores.get(unidad, 1)
    # Los promedios por comprador/viaje ya vienen en la unidad correcta
    # (sin escalar) para todas las unidades.
    factor_volumen_prom = 1
    factor_volumen_viaje = 1
    def _normalize_metric_label(label):
        normalized = unicodedata.normalize("NFKD", str(label))
        normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
        normalized = re.sub(r"\bweighted\b", "", normalized, flags=re.IGNORECASE)
        normalized = normalized.strip().upper()
        return re.sub(r"[^A-Z0-9]+", "", normalized)
    metric_index = df["Metric"].map(_normalize_metric_label).fillna("")
    metric_vertical_mask = metric_index.str.contains("VERT", na=False)
    def _metric_candidates(metrico):
        if isinstance(metrico, (list, tuple, set)):
            raw = list(metrico)
        else:
            raw = [metrico]
        return [_normalize_metric_label(m) for m in raw if m is not None]
    def leer_metric(metrico, columna):
        candidates = _metric_candidates(metrico)
        for candidate in candidates:
            mask = (metric_index == candidate) & (~metric_vertical_mask)
            if mask.any():
                return df.loc[mask, columna].values[0]
        for candidate in candidates:
            if not candidate:
                continue
            mask = metric_index.str.contains(candidate, na=False) & (~metric_vertical_mask)
            if mask.any():
                return df.loc[mask, columna].values[0]
        if isinstance(metrico, (list, tuple, set)):
            display_name = ", ".join(str(m) for m in metrico)
        else:
            display_name = str(metrico)
        raise KeyError(f"No se encontro la metrica '{display_name}' en los datos.")
    def build_entry(label, metrica, transform=lambda x, _: x):
        inicial = transform(leer_metric(metrica, periodo_inicial), periodo_inicial)
        final = transform(leer_metric(metrica, periodo_final), periodo_final)
        return {
            "valor_origen": inicial,
            "value": final,
            "change": calculate_percentage_change(inicial, final)
        }
    metrics_calculated = {
        etiquetas["valor"]: build_entry(etiquetas["valor"], 'Weighted VAL_LC'),
        etiquetas["volumen"]: build_entry(
            etiquetas["volumen"],
            ["Weighted R_VOL1", "Weighted R_VOL2", "Weighted VOL1_P", "Weighted VOL2_P", "Weighted VOLSU"],
            lambda v, _: v * factor
        ),
        etiquetas["precio"]: build_entry(
            etiquetas["precio"],
            ["Weighted PM1_LC", "Weighted PMSU_LC", "Weighted PM2_LC"],
            lambda v, _: v / factor
        ),
        etiquetas["gasto"]: build_entry(etiquetas["gasto"], 'Weighted VALC_BUY'),
        etiquetas["compradores"]: build_entry(etiquetas["compradores"], 'Weighted BUYERS'),
        etiquetas["volumen_prom"]: build_entry(
            etiquetas["volumen_prom"],
            ["Weighted VO1_BUY", "Weighted VOSU_BUY", "Weighted VO2_BUY"],
            lambda v, _: v * factor_volumen_prom
        ),
        etiquetas["penetracion"]: build_entry(etiquetas["penetracion"], 'Weighted PENET'),
        etiquetas["hholds"]: build_entry(etiquetas["hholds"], 'Weighted HHOLDS'),
        etiquetas["frecuencia"]: build_entry(etiquetas["frecuencia"], 'Weighted FREQ'),
        etiquetas["volumen_viaje"]: build_entry(
            etiquetas["volumen_viaje"],
            ["Weighted VO1_DAY", "Weighted VOSU_DAY", "Weighted VO2_DAY"],
            lambda v, _: v * factor_volumen_viaje
        ),
        etiquetas["ticket"]: build_entry(etiquetas["ticket"], 'Weighted VALC_DAY'),
    }
    return metrics_calculated
def _format_val(label_key, value):
    """Formato sensible a cada metrica para asemejarse al ejemplo de referencia."""
    reglas = {
        "Valor $ 000s": {"scale": 1 / 1000, "decimals": 0, "thousands": True},
        "Volumen 000": {"scale": 1, "decimals": 0, "thousands": True},
        "Precio Promedio": {"scale": 1, "decimals": 2, "thousands": False},
        "Gasto Promedio": {"scale": 1, "decimals": 2, "thousands": False},
        "Compradores 000s": {"scale": 1 / 1000, "decimals": 0, "thousands": True},
        # Mostrar con 1 decimal para que se vea 0.2 en lugar de 0 ó 170
        "Volumen Promedio": {"scale": 1, "decimals": 2, "thousands": False},
        "% Penetracion": {"scale": 1, "decimals": 1, "thousands": False},
        "Total HHolds 000s": {"scale": 1 / 1000, "decimals": 0, "thousands": True},
        "Frecuencia": {"scale": 1, "decimals": 1, "thousands": False},
        # Mostrar con 1 decimal (p.ej. 0.1) en lugar de redondear a entero
        "Volumen por Viaje": {"scale": 1, "decimals": 2, "thousands": False},
        "Gasto por Ticket": {"scale": 1, "decimals": 2, "thousands": False},
    }
    regla = next((v for k, v in reglas.items() if label_key.startswith(k)), {"scale": 1, "decimals": 2, "thousands": False})
    scaled = value * regla["scale"]
    if regla["thousands"]:
        fmt = f"{{:,.{regla['decimals']}f}}" if regla["decimals"] else "{:,.0f}"
    else:
        fmt = f"{{:.{regla['decimals']}f}}"
    return fmt.format(round(scaled, regla["decimals"]))
def graficar_arbol(metrics_calculated, volumen_unidad='Units', output_dir=None, hoja='', mostrar=False):
    """Dibuja el arbol de variables y retorna un stream PNG listo para insertar en la PPT."""
    etiquetas = _nombres_unidad(volumen_unidad)
    def wrap_label(text, max_chars=14):
        return textwrap.fill(text, width=max_chars, break_long_words=False, break_on_hyphens=False)
    level4_ids = {"penetracion", "hholds", "frecuencia1", "frecuencia2", "volumen_viaje", "ticket"}
    def card_size(node_id):
        if node_id in level4_ids:
            return 0.12, 0.11
        return 0.15, 0.12
    nodes = [
        {"id": "valor", "label": etiquetas["valor"], "metric": etiquetas["valor"], "pos": (0.50, 0.86), "color": "#8cc63e"},
        {"id": "volumen", "label": etiquetas["volumen"], "metric": etiquetas["volumen"], "pos": (0.33, 0.63), "color": "#8fd1c4"},
        {"id": "precio", "label": etiquetas["precio"], "metric": etiquetas["precio"], "pos": (0.70, 0.63), "color": "#bf5a9b"},
        {"id": "compradores", "label": etiquetas["compradores"], "metric": etiquetas["compradores"], "pos": (0.18, 0.46), "color": "#83a860"},
        {"id": "volumen_prom", "label": etiquetas["volumen_prom"], "metric": etiquetas["volumen_prom"], "pos": (0.50, 0.46), "color": "#6495c4"},
        {"id": "gasto", "label": etiquetas["gasto"], "metric": etiquetas["gasto"], "pos": (0.82, 0.46), "color": "#c671b7"},
        {"id": "penetracion", "label": etiquetas["penetracion"], "metric": etiquetas["penetracion"], "pos": (0.08, 0.29), "color": "#27a347"},
        {"id": "hholds", "label": etiquetas["hholds"], "metric": etiquetas["hholds"], "pos": (0.26, 0.29), "color": "#bdaa7e"},
        {"id": "frecuencia1", "label": etiquetas["frecuencia"], "metric": etiquetas["frecuencia"], "pos": (0.41, 0.29), "color": "#00aeef"},
        {"id": "volumen_viaje", "label": etiquetas["volumen_viaje"], "metric": etiquetas["volumen_viaje"], "pos": (0.58, 0.29), "color": "#0074c1"},
        {"id": "frecuencia2", "label": etiquetas["frecuencia"], "metric": etiquetas["frecuencia"], "pos": (0.74, 0.29), "color": "#00aeef"},
        {"id": "ticket", "label": etiquetas["ticket"], "metric": etiquetas["ticket"], "pos": (0.90, 0.29), "color": "#8a63ad"},
    ]
    edges = [
        ("valor", "volumen", "solid"),
        ("valor", "precio", "dashed"),
        ("valor", "gasto", "dashed"),
        ("volumen", "compradores", "solid"),
        ("volumen", "volumen_prom", "solid"),
        ("compradores", "penetracion", "solid"),
        ("compradores", "hholds", "solid"),
        ("volumen_prom", "frecuencia1", "solid"),
        ("volumen_prom", "volumen_viaje", "solid"),
        ("gasto", "frecuencia2", "solid"),
        ("gasto", "ticket", "solid"),
    ]
    pos_change_color = "#0bbdaf"
    neg_change_color = "#ef5350"
    neu_change_color = "#9e9e9e"
    def change_color(ch):
        return pos_change_color if ch > 3 else (neg_change_color if ch < -3 else neu_change_color)
    def draw_card(ax, center, node):
        metric_info = metrics_calculated[node["metric"]]
        value_txt = _format_val(node["label"], metric_info["value"])
        change_txt = f"{metric_info['change']:+.1f}%"
        label_wrapped = wrap_label(node["label"])
        label_lines = label_wrapped.split("\n")
        x, y = center
        width, height = card_size(node["id"])
        header_h = height * (0.36 + 0.08 * (len(label_lines) - 1))
        rounding = 0.02
        shadow = mpatches.FancyBboxPatch(
            (x - width / 2 + 0.008, y - height / 2 - 0.015),
            width, height,
            boxstyle=f"round,pad=0.01,rounding_size={rounding}",
            linewidth=0, facecolor="#9e9e9e", alpha=0.35, zorder=1)
        ax.add_patch(shadow)
        body = mpatches.FancyBboxPatch(
            (x - width / 2, y - height / 2),
            width, height,
            boxstyle=f"round,pad=0.01,rounding_size={rounding}",
            linewidth=1, edgecolor="#c1c1c1", facecolor="white", zorder=2)
        ax.add_patch(body)
        header = mpatches.FancyBboxPatch(
            (x - width / 2, y + height / 2 - header_h),
            width, header_h,
            boxstyle=f"round,pad=0.01,rounding_size={rounding}",
            linewidth=1, edgecolor=node["color"], facecolor=node["color"], zorder=3)
        ax.add_patch(header)
        ax.text(x, y + height / 2 - header_h / 2, label_wrapped, va="center", ha="center",
                fontsize=11, color="white", weight="bold", zorder=4)
        value_x = x - width * 0.22
        change_x = x + width * 0.22
        values_y = y - 0.035
        ax.text(value_x, values_y, value_txt, va="center", ha="center",
                fontsize=12, color="#222222", weight="bold", zorder=4)
        ax.text(change_x, values_y, change_txt, va="center", ha="center",
                fontsize=11, color=change_color(metric_info["change"]), weight="bold", zorder=4)
    def draw_edge(ax, from_id, to_id, style):
        x0, y0 = next(n for n in nodes if n["id"] == from_id)["pos"]
        x1, y1 = next(n for n in nodes if n["id"] == to_id)["pos"]
        dash_style = '--' if style == "dashed" else "solid"
        if from_id == "valor" and to_id == "gasto":
            w_valor, h_valor = card_size("valor")
            start_x = x0 + w_valor / 2
            start_y = y0 - h_valor * 0.15
            ax.annotate("",
                        xy=(x1, start_y), xytext=(start_x, start_y),
                        arrowprops=dict(arrowstyle='-', lw=1.2,
                                        linestyle=dash_style,
                                        color="#8c8c8c", shrinkA=0, shrinkB=0))
            ax.annotate("",
                        xy=(x1, y1 + 0.055), xytext=(x1, start_y),
                        arrowprops=dict(arrowstyle='-|>', lw=1.2,
                                        linestyle=dash_style,
                                        color="#8c8c8c", shrinkA=0, shrinkB=4))
            return
        _, h_from = card_size(from_id)
        split_offset = 0.68
        start_y = y0 - h_from * split_offset
        origin_offset = 0.60
        origin_y = y0 - h_from * origin_offset
        ax.annotate("",
                    xy=(x0, start_y), xytext=(x0, origin_y),
                    arrowprops=dict(arrowstyle='-', lw=1.2,
                                    linestyle=dash_style,
                                    color="#8c8c8c", shrinkA=0, shrinkB=0))
        ax.annotate("",
                    xy=(x1, y1 + 0.055), xytext=(x0, start_y),
                    arrowprops=dict(arrowstyle='-|>', lw=1.2,
                                    linestyle=dash_style,
                                    color="#8c8c8c", shrinkA=0, shrinkB=4))
    def draw_key(ax):
        ax.text(0.05, 0.88, "KEY", ha="left", va="center", fontsize=10, color="#565656", weight="bold")
        ax.add_patch(mpatches.Rectangle((0.09, 0.87), 0.02, 0.02, facecolor=pos_change_color, edgecolor="none"))
        ax.text(0.115, 0.88, "= > 3% Cambio", ha="left", va="center", fontsize=9, color=pos_change_color)
        ax.add_patch(mpatches.Rectangle((0.21, 0.87), 0.02, 0.02, facecolor=neg_change_color, edgecolor="none"))
        ax.text(0.235, 0.88, "= < -3% Cambio", ha="left", va="center", fontsize=9, color=neg_change_color)
    def draw_attribution_bar(ax, metrics):
        valor_info = metrics[etiquetas["valor"]]
        v0 = valor_info["valor_origen"]
        v1 = valor_info["value"]
        delta_v = v1 - v0
        if v0 <= 0 or v1 <= 0 or delta_v == 0:
            return
        drivers = [
            ("Frecuencia", etiquetas["frecuencia"], "#00b0f0"),
            ("% Penetracion", etiquetas["penetracion"], "#27a347"),
            ("Total HHolds 000s", etiquetas["hholds"], "#bdaa7e"),
            ("Volumen por Viaje", etiquetas["volumen_viaje"], "#1f78b4"),
            ("Precio Promedio", etiquetas["precio"], "#b4559c"),
        ]
        driver_logs = []
        for display, key, color in drivers:
            info = metrics[key]
            a, b = info["valor_origen"], info["value"]
            if a <= 0 or b <= 0:
                driver_logs.append((display, key, color, 0.0))
            else:
                driver_logs.append((display, key, color, math.log(b / a)))
        total_log = sum(item[3] for item in driver_logs)
        if total_log == 0:
            total_log = math.log(v1 / v0)
            if total_log == 0:
                return
        contribs = []
        for display, key, color, log_ratio in driver_logs:
            if log_ratio == 0:
                contribs.append((display, 0.0, 0.0, color))
                continue
            contrib_val = delta_v * (log_ratio / total_log)
            chart_pct = (contrib_val / abs(delta_v)) * 100
            contribs.append((display, chart_pct, contrib_val, color))
        total_abs = sum(abs(c[1]) for c in contribs)
        if total_abs == 0:
            return
        neg_abs = sum(abs(c[1]) for c in contribs if c[1] < 0)
        zero_x = neg_abs / total_abs
        bar_ax = fig.add_axes([0.065, 0.80, 0.36, 0.10])
        bar_ax.set_axis_off()
        bar_ax.set_xlim(0, 1)
        bar_ax.set_ylim(0, 1)
        bar_ax.plot([zero_x, zero_x], [0.30, 0.70], color="#5f5f5f", linewidth=1)
        x_left = zero_x
        x_right = zero_x
        for display, pct, _, color in contribs:
            width = abs(pct) / total_abs
            if pct < 0:
                x_left -= width
                bar_ax.add_patch(mpatches.Rectangle((x_left, 0.35), width, 0.30, facecolor=color, edgecolor="none"))
            else:
                bar_ax.add_patch(mpatches.Rectangle((x_right, 0.35), width, 0.30, facecolor=color, edgecolor="none"))
                x_right += width
        ax.text(0.06, 0.92, "Atribucion del cambio en el gasto", ha="left", va="center",
                fontsize=14, weight="bold", color="#333333")
        draw_key(ax)
    fig = plt.figure(figsize=(21, 9), dpi=DEFAULT_EXPORT_DPI)
    ax = plt.gca()
    ax.set_axis_off()
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    for p, c, style in edges:
        draw_edge(ax, p, c, style)
    for node in nodes:
        draw_card(ax, node["pos"], node)
    draw_attribution_bar(ax, metrics_calculated)
    plt.title(f"Arbol de Variables ({volumen_unidad})", fontsize=10, loc="right", color="#666")
    fig_size = tuple(fig.get_size_inches()) if fig is not None else None
    buf = figure_to_stream(fig, dpi=DEFAULT_EXPORT_DPI, bbox_inches='tight', pad_inches=0.02, close=False)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        suffix = f"_{hoja}" if hoja else ""
        image_path = os.path.join(output_dir, f"arbol_{volumen_unidad}{suffix}.png")
        with open(image_path, "wb") as handler:
            handler.write(buf.getvalue())
        buf.seek(0)
    if mostrar:
        plt.show()
    plt.close(fig)
    return buf, fig_size
def _unidad_desde_nombre_hoja(nombre_hoja):
    """
    Si el nombre cumple con el patron '2_NombreMarca_U', retorna la unidad solicitada.
    Ejemplo: '2_Coca Cola_L' -> 'Litros'. Si no coincide o la letra no es valida, devuelve None.
    """
    match = re.match(r"^\d+_(.+)_([A-Za-z])$", nombre_hoja.strip())
    if not match:
        return None
    letra = match.group(2).upper()
    mapa = {
        "U": "Units",
        "L": "Litros",
        "K": "Kilos",
        "T": "Toneladas",
        "R": "Rollos",
        "M": "Metros",
        "H": "Hojas",
    }
    return mapa.get(letra)
# Segmentos 3-6: tabla de aportes y calculos de shares/variaciones.
def aporte(df,p,lang,tipo):
        """Construye la tabla de aporte y calcula shares/variaciones."""
        aux = df.copy()
        removed_non_numeric: list[str] = []
        if aux.empty:
            apo = pd.DataFrame(columns=[tipo])
            apo.attrs["removed_columns"] = []
            apo.attrs["share_period_values"] = []
            apo.attrs["share_mat_values"] = {}
            apo.attrs["skip_table"] = True
            apo.attrs["skip_table_reason"] = "empty_data"
            return apo
        first_col_name = aux.columns[0]
        first_col = aux.iloc[:, 0].copy()
        numeric_data = OrderedDict()
        for idx in range(1, aux.shape[1]):
            col_name = aux.columns[idx]
            series = aux.iloc[:, idx]
            if np.issubdtype(series.dtype, np.datetime64):
                removed_non_numeric.append(str(col_name))
                continue
            numeric_series = pd.to_numeric(series, errors='coerce')
            if numeric_series.dropna().empty:
                removed_non_numeric.append(str(col_name))
                continue
            numeric_data[col_name] = numeric_series
        if numeric_data:
            sanitized_aux = pd.concat(
                [first_col.to_frame(name=first_col_name), pd.DataFrame(numeric_data)],
                axis=1
            )
        else:
            sanitized_aux = first_col.to_frame(name=first_col_name)
        sanitized_aux.reset_index(drop=True, inplace=True)
        aux = sanitized_aux
        if aux.shape[1] <= 1:
            apo = pd.DataFrame(columns=[tipo])
            apo.attrs["removed_columns"] = list(dict.fromkeys(removed_non_numeric))
            apo.attrs["share_period_values"] = []
            apo.attrs["share_mat_values"] = {}
            apo.attrs["skip_table"] = True
            apo.attrs["skip_table_reason"] = "no_numeric_columns"
            return apo
        try:
            pipeline_offset = int(p) if p is not None else 0
        except (TypeError, ValueError):
            pipeline_offset = 0
        if pipeline_offset < 0:
            pipeline_offset = 0
        p = pipeline_offset
        required_rows = 24 + pipeline_offset
        if len(aux) < required_rows:
            apo = pd.DataFrame(columns=[tipo])
            apo.attrs["removed_columns"] = list(dict.fromkeys(removed_non_numeric))
            apo.attrs["share_period_values"] = []
            apo.attrs["share_mat_values"] = {}
            apo.attrs["skip_table"] = True
            apo.attrs["skip_table_reason"] = "insufficient_periods"
            return apo
        date_col_label = labels.get((lang, 'Data'))
        date_col_name = date_col_label if date_col_label in aux.columns else first_col_name
        aux_numeric = aux.iloc[:, 1:].apply(pd.to_numeric, errors='coerce')
        aux['Total'] = aux_numeric.sum(axis=1, skipna=True)
        apo=pd.DataFrame(columns=[tipo] + aux.columns[1:].tolist())
        #Vol ultimo MAT
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-12-p,date_col_name].strftime('%b-%y') ] + [aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p,aux.shape[1]-1].sum() for col in range(1,len(aux.columns) )]
        #Vol MAT actual
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-p,date_col_name].strftime('%b-%y') ] + [aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-12-p:len(aux)-p,aux.shape[1]-1].values.sum() for col in range(1,len(aux.columns))]
        #Var
        apo.loc[len(apo)] = ["Var %"] +  [ aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum()-1 for col in range(1,len(aux.columns))]
        #Aporte
        apo.loc[len(apo)] = ["Aporte"] + [apo.iloc[-1, col] * apo.iloc[0, col] for col in range(1,len(aux.columns))]
        # Remove categories with invalid metrics (NaN/Inf) before formatting
        def _is_invalid_metric(value):
            if pd.isna(value):
                return True
            try:
                numeric = float(value)
            except (TypeError, ValueError):
                return False
            return not np.isfinite(numeric)
        invalid_columns = []
        for col_idx in range(1, len(apo.columns)):
            col_name = apo.columns[col_idx]
            var_value = apo.iloc[2, col_idx]
            aporte_value = apo.iloc[3, col_idx]
            if _is_invalid_metric(var_value) or _is_invalid_metric(aporte_value):
                invalid_columns.append(col_name)
        if invalid_columns:
            apo.drop(columns=invalid_columns, inplace=True)
        removed_all = list(dict.fromkeys(removed_non_numeric + invalid_columns))
        apo.attrs["removed_columns"] = removed_all
        share_period_values = []
        max_share_rows = min(2, len(apo.index))
        for row_idx in range(max_share_rows):
            period_label = normalize_period_label(apo.iloc[row_idx, 0])
            period_shares = OrderedDict()
            for col_idx in range(1, len(apo.columns)):
                column_name = str(apo.columns[col_idx])
                normalized_name = column_name.strip().lower()
                if not normalized_name or normalized_name == 'total':
                    continue
                value = apo.iloc[row_idx, col_idx]
                numeric_value = None
                try:
                    numeric_value = float(value)
                except (TypeError, ValueError):
                    try:
                        numeric_value = float(str(value).replace('%', '').replace(',', '.')) / 100.0
                    except Exception:
                        numeric_value = None
                if numeric_value is None or not np.isfinite(numeric_value) or numeric_value < 0:
                    numeric_value = 0.0
                period_shares[column_name] = numeric_value
            share_period_values.append((period_label, period_shares))
        apo.attrs["share_period_values"] = share_period_values
        share_mat_values = {}
        if len(share_period_values) > 1:
            _, current_period_shares = share_period_values[1]
            share_mat_values = {
                column_name: value
                for column_name, value in current_period_shares.items()
                if np.isfinite(value)
            }
        apo.attrs["share_mat_values"] = share_mat_values
        #Formato del volumen
        apo.iloc[:2, 1:] = apo.iloc[:2, 1:].applymap(lambda x: f"{round(x * 100, 1)}%")
        #Formato de la variacion y el aporte
        apo.iloc[2:, 1:] = apo.iloc[2:, 1:].applymap(lambda x: f"{round(x * 100, 2)}%")
        return apo
# Funcion que crea el grafico de la tabla de aporte
def graf_apo(apo, c_fig, column_color_mapping=None):
    """Renderiza la tabla de aporte como imagen PNG."""
    fig, ax = plt.subplots(num=c_fig, dpi=DEFAULT_EXPORT_DPI)
    row_height = 0.45
    col_width = 1.0
    n_rows, n_cols = apo.shape
    fig_width = col_width * n_cols
    fig_height = row_height * (n_rows + 0.2)
    fig.set_size_inches(fig_width, fig_height)
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax.axis('off')
    ax.set_facecolor('white')
    display_data = apo.copy()
    if display_data.shape[1] > 0:
        display_data.iloc[:, 0] = display_data.iloc[:, 0].map(wrap_table_text)
    display_col_labels = [wrap_table_text(str(label)) for label in apo.columns]
    table = ax.table(
        cellText=display_data.values,
        colLabels=display_col_labels,
        loc='center',
        cellLoc='center'
    )
    for i, _ in enumerate(apo.columns):
        table.auto_set_column_width(i)
    table.scale(1, 1.2)
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    header_main = HEADER_COLOR_PRIMARY
    header_secondary = HEADER_COLOR_SECONDARY
    total_fill = HEADER_TOTAL_FILL
    first_col_fill = HEADER_FIRST_COL_FILL
    positive_color = TABLE_POSITIVE_COLOR
    negative_color = TABLE_NEGATIVE_COLOR
    bar_padding_ratio = 0.08
    bar_height_ratio = 0.55
    column_color_mapping = column_color_mapping or {}
    exact_color_mapping = {
        str(key).strip(): value
        for key, value in column_color_mapping.items()
        if key is not None and value
    }
    normalized_color_mapping = {
        key.lower(): value
        for key, value in exact_color_mapping.items()
        if key
    }
    suffix_variants = ('.c', '.v', '_c', '_v', '-c', '-v')
    def resolve_column_color(column_name):
        if column_name is None:
            return None
        key = str(column_name).strip()
        if not key:
            return None
        color_value = exact_color_mapping.get(key)
        if color_value:
            return color_value
        lower_key = key.lower()
        color_value = normalized_color_mapping.get(lower_key)
        if color_value:
            return color_value
        for suffix in suffix_variants:
            if lower_key.endswith(suffix):
                base = key[:-len(suffix)].strip()
                if not base:
                    continue
                color_value = exact_color_mapping.get(base)
                if color_value:
                    return color_value
                color_value = normalized_color_mapping.get(base.lower())
                if color_value:
                    return color_value
        return None
    total_column_name = next((col for col in apo.columns if str(col).strip().lower() == 'total'), None)
    data_columns = [
        idx
        for idx in range(1, n_cols)
        if total_column_name is None or apo.columns[idx] != total_column_name
    ]
    start_rgb = np.array(mcolors.to_rgb(VOLUME_BAR_START))
    end_rgb = np.array(mcolors.to_rgb(VOLUME_BAR_END))
    def share_colors(count: int):
        if count <= 0:
            return []
        if count == 1:
            return [mcolors.to_hex(np.clip(start_rgb, 0, 1))]
        colors = []
        for i in range(count):
            ratio = i / (count - 1)
            rgb = np.clip(start_rgb + (end_rgb - start_rgb) * ratio, 0, 1)
            colors.append(mcolors.to_hex(rgb))
        return colors
    fallback_column_colors = share_colors(len(data_columns))
    column_colors = []
    for rel_idx, col_idx in enumerate(data_columns):
        resolved = resolve_column_color(apo.columns[col_idx])
        if resolved is None:
            resolved = fallback_column_colors[rel_idx] if rel_idx < len(fallback_column_colors) else VOLUME_BAR_END
        try:
            column_colors.append(mcolors.to_hex(resolved))
        except (ValueError, TypeError):
            column_colors.append(fallback_column_colors[rel_idx] if rel_idx < len(fallback_column_colors) else VOLUME_BAR_END)
    white_rgb = np.array([1.0, 1.0, 1.0])
    volume_row_indexes = [idx for idx, label in enumerate(apo.iloc[:, 0]) if str(label).lower().startswith('vol')]
    aporte_row_index = next((idx for idx, label in enumerate(apo.iloc[:, 0]) if str(label).strip().lower() == 'aporte'), None)
    def to_percentage(value):
        try:
            text_value = str(value).strip().replace('%', '').replace(',', '.')
            if not text_value:
                return None
            return max(0.0, float(text_value) / 100.0)
        except Exception:
            return None
    def to_signed_percentage(value):
        try:
            text_value = str(value).strip().replace('%', '').replace(',', '.')
            if not text_value:
                return None
            return float(text_value) / 100.0
        except Exception:
            return None
    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor(TABLE_GRID_COLOR)
        cell.set_linewidth(1)
        cell.get_text().set_zorder(3)
        if row == 0:
            text = cell.get_text()
            text.set_weight('bold')
            text.set_fontsize(TABLE_HEADER_FONT_SIZE)
            header_text_color = HEADER_FONT_COLOR
            if col == 0 or (total_column_name is not None and apo.columns[col] == total_column_name):
                cell.set_facecolor(header_main)
            else:
                header_color = resolve_column_color(apo.columns[col])
                if header_color is None:
                    cell.set_facecolor(header_secondary)
                else:
                    try:
                        cell.set_facecolor(header_color)
                        rgb = mcolors.to_rgb(header_color)
                        luminance = 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]
                        header_text_color = '#FFFFFF' if luminance < 0.6 else TABLE_TEXT_PRIMARY
                    except (ValueError, TypeError):
                        cell.set_facecolor(header_secondary)
            text.set_color(header_text_color)
        else:
            label = str(apo.iloc[row - 1, 0])
            text = cell.get_text()
            if col == 0:
                cell.set_facecolor(first_col_fill)
                text.set_weight('bold')
                text.set_color(TABLE_TEXT_PRIMARY)
            else:
                if total_column_name is not None and apo.columns[col] == total_column_name:
                    cell.set_facecolor(total_fill)
                else:
                    cell.set_facecolor('white')
                if label.lower().startswith('vol'):
                    text.set_color(TABLE_TEXT_PRIMARY)
                else:
                    value = apo.iloc[row - 1, col]
                    try:
                        numeric = float(str(value).replace('%', '').replace(',', '.'))
                        text.set_color(positive_color if numeric >= 0 else negative_color)
                    except Exception:
                        pass
    if aporte_row_index is not None:
        aporte_values = []
        for col in data_columns:
            numeric = to_signed_percentage(apo.iloc[aporte_row_index, col])
            if numeric is not None:
                aporte_values.append((col, numeric))
        if aporte_values:
            aporte_numbers = [value for _, value in aporte_values]
            min_val = min(aporte_numbers)
            max_val = max(aporte_numbers)
            cmap = None
            norm = None
            if not math.isclose(min_val, max_val, rel_tol=1e-9, abs_tol=1e-9):
                if min_val < 0 and max_val > 0:
                    cmap = mcolors.LinearSegmentedColormap.from_list(
                        'aporte_diverging',
                        [APORTE_NEGATIVE_COLOR, APORTE_NEUTRAL_COLOR, TABLE_POSITIVE_COLOR]
                    )
                    norm = mcolors.TwoSlopeNorm(vmin=min_val, vcenter=0, vmax=max_val)
                elif max_val <= 0:
                    cmap = mcolors.LinearSegmentedColormap.from_list(
                        'aporte_negative',
                        [APORTE_NEGATIVE_COLOR, APORTE_NEUTRAL_COLOR]
                    )
                    norm = mcolors.Normalize(vmin=min_val, vmax=0)
                else:
                    cmap = mcolors.LinearSegmentedColormap.from_list(
                        'aporte_positive',
                        [APORTE_NEUTRAL_COLOR, TABLE_POSITIVE_COLOR]
                    )
                    norm = mcolors.Normalize(vmin=0, vmax=max_val)
            for col, numeric in aporte_values:
                if cmap is not None and norm is not None:
                    rgba = cmap(norm(numeric))
                    rgb = tuple(rgba[:3])
                else:
                    if math.isclose(numeric, 0.0, abs_tol=1e-9):
                        rgb = mcolors.to_rgb(APORTE_NEUTRAL_COLOR)
                    elif numeric > 0:
                        rgb = mcolors.to_rgb(TABLE_POSITIVE_COLOR)
                    else:
                        rgb = mcolors.to_rgb(APORTE_NEGATIVE_COLOR)
                cell = table[(aporte_row_index + 1, col)]
                cell.set_facecolor(rgb)
                text = cell.get_text()
                luminance = 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]
                text.set_color(TABLE_TEXT_PRIMARY if luminance > 0.5 else '#FFFFFF')
    for df_row_idx in volume_row_indexes:
        table_row = df_row_idx + 1
        row_percentages = {}
        max_percent = 0.0
        for col in data_columns:
            percent = to_percentage(apo.iloc[df_row_idx, col])
            row_percentages[col] = percent
            if percent is not None and percent > max_percent:
                max_percent = percent
        for rel_idx, col in enumerate(data_columns):
            percent = row_percentages.get(col)
            if percent is None:
                continue
            cell = table[(table_row, col)]
            bar_color = column_colors[rel_idx] if rel_idx < len(column_colors) else VOLUME_BAR_END
            base_rgb = np.array(mcolors.to_rgb(bar_color))
            intensity = min(percent / max_percent, 1.0) if max_percent > 0 else 0.0
            blended_rgb = white_rgb * (1.0 - intensity) + base_rgb * intensity
            cell.set_facecolor(blended_rgb)
            cell.get_text().set_color(TABLE_TEXT_PRIMARY)
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    table_bbox = table.get_tightbbox(renderer).transformed(fig.dpi_scale_trans.inverted())
    target_height_in = TABLE_TARGET_HEIGHT_CM / 2.54
    if table_bbox.height > 0:
        scale = target_height_in / table_bbox.height
        if abs(scale - 1.0) > 1e-3:
            current_width, current_height = fig.get_size_inches()
            fig.set_size_inches(current_width * scale, current_height * scale)
            fig.canvas.draw()
            renderer = fig.canvas.get_renderer()
            table_bbox = table.get_tightbbox(renderer).transformed(fig.dpi_scale_trans.inverted())
    return figure_to_stream(fig, dpi=TABLE_EXPORT_DPI, bbox_inches=table_bbox, pad_inches=0)
def simplify_name_segment(value: str, max_len: int) -> str:
    """
    Reduce una cadena a un segmento seguro para archivos.
    Conserva alfanumericos, reemplaza otros con guiones y limita longitud.
    """
    if value is None:
        return 'NA'
    value_str = str(value).strip()
    if not value_str:
        return 'NA'
    cleaned = ''.join(ch if ch.isalnum() else '-' for ch in value_str)
    cleaned = cleaned.strip('-')
    if not cleaned:
        cleaned = 'NA'
    return cleaned[:max_len]
def select_excel_file(base_dir: Path) -> str:
    """Lista archivos Excel disponibles y devuelve el nombre seleccionado."""
    excel_files = sorted([p for p in base_dir.glob('*.xlsx') if not p.name.startswith('~$')])
    def metadata_for(file_path: Path) -> str:
        try:
            if file_path.stem.lower() == 'plantilla_entrada_5w1h':
                return 'Plantilla de entrada - usar como referencia'
            parts = file_path.stem.split('_')
            country_code = int(parts[0]) if parts else None
            category_code = parts[1].strip().upper() if len(parts) > 1 else None
            pais_name = pais.loc[pais['cod'] == country_code, 'pais'].iloc[0] if country_code is not None else 'Pais N/D'
            cesta = categ.loc[categ['cod'] == category_code, 'cest'].iloc[0] if category_code else 'Cesta N/D'
            categoria = categ.loc[categ['cod'] == category_code, 'cat'].iloc[0] if category_code else 'Categoria N/D'
            return f"{pais_name} | {cesta} | {categoria}"
        except Exception:
            return 'Metadata no disponible'
    if not excel_files:
        raise FileNotFoundError(colorize(f'No se encontraron archivos .xlsx en {base_dir}', COLOR_RED))
    if len(excel_files) == 1:
        meta = metadata_for(excel_files[0])
        meta_color = COLOR_RED if meta == 'Metadata no disponible' else COLOR_YELLOW
        print_colored(f"Archivo encontrado: {excel_files[0].name} {colorize('[' + meta + ']', meta_color)}", COLOR_BLUE)
        return excel_files[0].name
    print_colored('Archivos Excel disponibles:', COLOR_QUESTION)
    for idx, archivo in enumerate(excel_files, 1):
        meta = metadata_for(archivo)
        meta_color = COLOR_RED if meta == 'Metadata no disponible' else COLOR_YELLOW
        print(f"{colorize(f'  {idx}. {archivo.name}', COLOR_BLUE)} {colorize('[' + meta + ']', meta_color)}")
    prompt = colorize(f"Seleccione el numero del archivo (1-{len(excel_files)}). Enter para {excel_files[0].name}: ", COLOR_QUESTION)
    while True:
        choice = input(prompt).strip()
        if not choice:
            return excel_files[0].name
        if choice.isdigit():
            idx = int(choice)
            if 1 <= idx <= len(excel_files):
                return excel_files[idx - 1].name
        matching = [archivo.name for archivo in excel_files if archivo.name.lower() == choice.lower()]
        if matching:
            return matching[0]
        print_colored('Entrada invalida. Intente nuevamente.', COLOR_RED)
def configure_cover_slide(presentation, lang):
    """Configura el titulo de portada segun el idioma."""
    title_by_lang = {
        'P': 'Adicionais de Cobertura - 5W1H',
        'E': 'Adicionales de Cobertura - 5W1H'
    }
    cover_slide = presentation.slides[0]
    target_text = title_by_lang.get(lang, title_by_lang['E'])
    fallback_geometry = (Inches(0.42), Inches(2.0), Inches(10), Inches(0.8))
    geometry = None
    for shape in list(cover_slide.shapes):
        if not getattr(shape, 'has_text_frame', False):
            continue
        text_value = shape.text.strip()
        if 'Cobertura' in text_value and '5W1H' in text_value:
            geometry = (shape.left, shape.top, shape.width, shape.height)
            cover_slide.shapes._spTree.remove(shape._element)
    left, top, width, height = geometry if geometry else fallback_geometry
    text_box = cover_slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = target_text
    run = paragraph.runs[0]
    font = run.font
    font.name = 'Arial'
    font.size = Inches(0.45)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)
def plot_ven():
    """Solicita el modo de grafico para ventas y compras."""
    options = [
        ("1 - Plotear Ventas y Compras Juntas", "1"),
        ("2 - Plotear Ventas y Compras Separadas", "2"),
        ("3 - No hay W con Ventas en esa base", "3"),
    ]
    print_colored('\nOpciones de grafico disponibles:', COLOR_QUESTION)
    for idx, (texto, _) in enumerate(options, 1):
        print_colored(f"  {idx}. {texto}", COLOR_BLUE)
    prompt = colorize(f"Seleccione el numero de la opcion (1-{len(options)}). Enter para {options[0][1]}: ", COLOR_QUESTION)
    while True:
        choice = input(prompt).strip()
        if not choice:
            return options[0][1]
        if choice.isdigit():
            idx = int(choice)
            if 1 <= idx <= len(options):
                return options[idx - 1][1]
        for _, valor in options:
            if choice == valor:
                return valor
        print_colored('Entrada invalida. Intente nuevamente.', COLOR_RED)
#Codigos paises
pais = pd.DataFrame(
    {'cod': [10, 54, 91, 55, 12, 56, 57, 93, 52, 51, 66, 63, 62, 64, 65, 67, 69],
     'pais': ['LatAm', 'Argentina', 'Bolivia', 'Brasil', 'CAM', 'Chile', 'Colombia', 'Ecuador', 'Mexico', 'Peru', 
              'Costa Rica', 'El Salvador', 'Guatemala', 'Honduras', 'Nicaragua', 'Panama', 'Republica Dominicana']
    }
)
#Codigos categorias
CATEG_CSV_DATA = """
cod,cest,cat
ALCB,Bebidas,Bebidas Alcoholicas
BEER,Bebidas,Cervezas
CARB,Bebidas,Bebidas Gaseosas
CWAT,Bebidas,Agua Gasificada
COCW,Bebidas,Agua de Coco
COFF,Bebidas,Cafe-Consolidado de Cafe
CRBE,Bebidas,Cross Category (Bebidas)
ENDR,Bebidas,Bebidas Energeticas
FLBE,Bebidas,Bebidas Saborizadas Sin Gas
GCOF,Bebidas,Cafe Tostado y Molido
HJUI,Bebidas,Jugos Caseros
ITEA,Bebidas,Te Helado
ICOF,Bebidas,Cafe Instantaneo-Cafe Sucedaneo
JUNE,Bebidas,Jugos y Nectares
VEJU,Bebidas,Zumos de Vegetales
WATE,Bebidas,Agua Natural
CSDW,Bebidas,Gaseosas + Aguas
MXCM,Bebidas,Mixta Cafe+Malta
MXDG,Bebidas,Mixta Dolce Gusto-Mixta Te Helado + Cafe + Modificadores
MXJM,Bebidas,Mixta Jugos y Leches
MXJS,Bebidas,Mixta Jugos Liquidos + Bebidas de Soja
MXTC,Bebidas,Mixta Te+Cafe
JUIC,Bebidas,Jugos Liquidos-Jugos Polvo
PWDJ,Bebidas,Refrescos en Polvo-Jugos - Bebidas Instantaneas En Polvo - Jugos Polvo
RFDR,Bebidas,Bebidas Refrescantes
RTDJ,Bebidas,Refrescos Liquidos-Jugos Liquidos
RTEA,Bebidas,Te Liquido - Listo para Tomar
SOYB,Bebidas,Bebidas de Soja
SPDR,Bebidas,Bebidas Isotonicas
TEAA,Bebidas,Te e Infusiones-Te-Infusion Hierbas
YERB,Bebidas,Yerba Mate
BUTT,Lacteos,Manteca
CHEE,Lacteos,Queso Fresco y para Untar
CMLK,Lacteos,Leche Condensada
CRCH,Lacteos,Queso Untable
DYOG,Lacteos,Yoghurt p-beber
EMLK,Lacteos,Leche Culinaria-Leche Evaporada
FRMM,Lacteos,Leche Fermentada
FMLK,Lacteos,Leche Liquida Saborizada-Leche Liquida Con Sabor
FRMK,Lacteos,Formulas Infantiles
LQDM,Lacteos,Leche Liquida
LLFM,Lacteos,Leche Larga Vida
MARG,Lacteos,Margarina
MCHE,Lacteos,Queso Fundido
MKCR,Lacteos,Crema de Leche
MXDI,Lacteos,Mixta Lacteos-Postre+Leches+Yogurt
MXMI,Lacteos,Mixta Leches
MXYD,Lacteos,Mixta Yoghurt+Postres
PTSS,Lacteos,Petit Suisse
PWDM,Lacteos,Leche en Polvo
SYOG,Lacteos,Yoghurt p-comer
MILK,Lacteos,Leche-Leche Liquida Blanca - Leche Liq. Natural
YOGH,Lacteos,Yoghurt
CLOT,Ropas y Calzados,Ropas
FOOT,Ropas y Calzados,Calzados
SOCK,Ropas y Calzados,Medias-Calcetines
AREP,Alimentos,Arepas
BCER,Alimentos,Cereales Infantiles
BABF,Alimentos,Nutricion Infantil-Colados y Picados
BEAN,Alimentos,Frijoles
BISC,Alimentos,Galletas
BOUI,Alimentos,Caldos-Caldos y Sazonadores
BREA,Alimentos,Pan
BRCR,Alimentos,Apanados-Empanizadores
BRDC,Alimentos,Empanados
CERE,Alimentos,Cereales-Cereales Desayuno-Avenas y Cereales
BURG,Alimentos,Hamburguesas
CCMX,Alimentos,Mezclas Listas para Tortas-Preparados Base Harina Trigo
CAKE,Alimentos,Queques-Ponques Industrializados
FISH,Alimentos,Conservas De Pescado
CFAV,Alimentos,Conservas de Frutas y Verduras
CRML,Alimentos,Dulce de Leche-Manjar
CMLC,Alimentos,Alfajores
CBAR,Alimentos,Barras de Cereal
CHCK,Alimentos,Pollo
CHOC,Alimentos,Chocolate
COCO,Alimentos,Chocolate de Taza-Achocolatados - Cocoas
COLS,Alimentos,Salsas Frias
COMP,Alimentos,Compotas
SPIC,Alimentos,Condimentos y Especias
CKCH,Alimentos,Chocolate de Mesa
COIL,Alimentos,Aceite-Aceites Comestibles
CSAU,Alimentos,Salsas Listas-Salsas Caseras Envasadas
CNML,Alimentos,"Grano, Harina y Masa de Maiz"
CNST,Alimentos,Fecula de Maiz
CNFL,Alimentos,Harina De Maiz
CAID,Alimentos,Ayudantes Culinarios
DESS,Alimentos,Postres Preparados
DHAM,Alimentos,Jamon Endiablado
DFNS,Alimentos,Semillas y Frutos Secos
EBRE,Alimentos,Pan de Pascua
EEGG,Alimentos,Huevos de Pascua
EGGS,Alimentos,Huevos
FLSS,Alimentos,Flash Cecinas
FLOU,Alimentos,Harinas
MEAT,Alimentos,Carne Fresca
FRDS,Alimentos,Platos Listos Congelados
FRFO,Alimentos,Alimentos Congelados
HAMS,Alimentos,Jamones
HCER,Alimentos,Cereales Calientes-Cereales Precocidos
HOTS,Alimentos,Salsas Picantes
ICEC,Alimentos,Helados
IBRE,Alimentos,Pan Industrializado
IMPO,Alimentos,Pure Instantaneo
INOO,Alimentos,Fideos Instantaneos
JAMS,Alimentos,Mermeladas
KETC,Alimentos,Ketchup
LJDR,Alimentos,Jugo de Limon Adereso
MALT,Alimentos,Maltas
SEAS,Alimentos,Adobos - Sazonadores
MAYO,Alimentos,Mayonesa
MEAT,Alimentos,Carnicos
MLKM,Alimentos,Modificadores de Leche-Saborizadores p-leche
MXCO,Alimentos,Mixta Cereales Infantiles+Avenas
MXBS,Alimentos,Mixta Caldos + Saborizantes
MXSB,Alimentos,Mixta Caldos + Sopas
MXCH,Alimentos,Mixta Cereales + Cereales Calientes
MXCC,Alimentos,Mixta Chocolate + Manjar
MXSN,Alimentos,"Galletas, snacks y mini tostadas"
COBT,Alimentos,Aceites + Mantecas
COCF,Alimentos,Aceites + Conservas De Pescado
CABB,Alimentos,Ayudantes Culinarios + Bolsa de Hornear
MXEC,Alimentos,Mixta Huevos de Pascua + Chocolates
MXDP,Alimentos,Mixta Platos Listos Congelados + Pasta
MXFR,Alimentos,Mixta Platos Congelados y Listos para Comer
MXFM,Alimentos,Mixta Alimentos Congelados + Margarina
MXMC,Alimentos,Mixta Modificadores + Cocoa
MXPS,Alimentos,Mixta Pastas
MXSO,Alimentos,Mixta Sopas+Cremas+Ramen
MXSP,Alimentos,Mixta Margarina + Mayonesa + Queso Crema
MXSW,Alimentos,Mixta Azucar+Endulzantes
MUST,Alimentos,Mostaza
NDCR,Alimentos,Sustitutos de Crema
NOOD,Alimentos,Fideos
NUGG,Alimentos,Nuggets
OAFL,Alimentos,Avena en hojuelas-liquidas
OLIV,Alimentos,Aceitunas
PANC,Alimentos,Tortilla
PANE,Alimentos,Paneton
PAST,Alimentos,Pastas
PSAU,Alimentos,Salsas para Pasta
PNOU,Alimentos,Turron de mani
PORK,Alimentos,Carne Porcina
PPMX,Alimentos,Postres en Polvo-Postres para Preparar - Horneables-Gelificables
PWSM,Alimentos,Leche de Soya en Polvo
PCCE,Alimentos,Cereales Precocidos
DOUG,Alimentos,Masas Frescas-Tapas Empanadas y Tarta
PPIZ,Alimentos,Pre-Pizzas
REFR,Alimentos,Meriendas listas
RICE,Alimentos,Arroz
RBIS,Alimentos,Galletas de Arroz
RTEB,Alimentos,Frijoles Procesados
RTEM,Alimentos,Pratos Prontos - Comidas Listas
SDRE,Alimentos,Aderezos para Ensalada
SALT,Alimentos,Sal
SLTC,Alimentos,Galletas Saladas-Galletas No Dulce
SARD,Alimentos,Sardina Envasada
SAUS,Alimentos,Cecinas
SCHN,Alimentos,Milanesas
SNAC,Alimentos,Snacks
SNOO,Alimentos,Fideos Sopa
SOUP,Alimentos,Sopas-Sopas Cremas
SOYS,Alimentos,Siyau
SPAG,Alimentos,Tallarines-Spaguetti
SPCH,Alimentos,Chocolate para Untar
SUGA,Alimentos,Azucar
SWCO,Alimentos,Galletas Dulces
SWSP,Alimentos,Untables Dulces
SWEE,Alimentos,Endulzantes
TOAS,Alimentos,Torradas - Tostadas
TOMA,Alimentos,Salsas de Tomate
TUNA,Alimentos,Atun Envasado
VMLK,Alimentos,Leche Vegetal
WFLO,Alimentos,Harinas de trigo
AIRC,Cuidado del Hogar,Ambientadores-Desodorante Ambiental
BARS,Cuidado del Hogar,Jabon en Barra-Jabon de lavar
BLEA,Cuidado del Hogar,Cloro-Lavandinas-Lejias-Blanqueadores
CBLK,Cuidado del Hogar,Pastillas para Inodoro
CGLO,Cuidado del Hogar,Guantes de latex
CLSP,Cuidado del Hogar,Esponjas de Limpieza-Esponjas y panos
CLTO,Cuidado del Hogar,Utensilios de Limpieza
FILT,Cuidado del Hogar,Filtros de Cafe
CRHC,Cuidado del Hogar,Cross Category (Limpiadores Domesticos)
CRLA,Cuidado del Hogar,Cross Category (Lavanderia)
CRPA,Cuidado del Hogar,Cross Category (Productos de Papel)
DISH,Cuidado del Hogar,Lavavajillas-Lavaplatos - Lavalozas mano
DPAC,Cuidado del Hogar,Empaques domesticos-Bolsas plasticas-Plastico Adherente-Papel encerado-Papel aluminio
DRUB,Cuidado del Hogar,Destapacanerias
FBRF,Cuidado del Hogar,Perfumantes para Ropa-Perfumes para Ropa
FWAX,Cuidado del Hogar,Cera p-pisos
FDEO,Cuidado del Hogar,Desodorante para Pies
FRNP,Cuidado del Hogar,Lustramuebles
GBBG,Cuidado del Hogar,Bolsas de Basura
GCLE,Cuidado del Hogar,Limpiadores verdes
CLEA,Cuidado del Hogar,Limpiadores-Limpiadores y Desinfectantes
INSE,Cuidado del Hogar,Insecticidas-Raticidas
KITT,Cuidado del Hogar,Toallas de papel-Papel Toalla - Toallas de Cocina - Rollos Absorbentes de Papel
LAUN,Cuidado del Hogar,Detergentes para ropa
LSTA,Cuidado del Hogar,Apresto
MXBC,Cuidado del Hogar,Mixta Pastillas para Inodoro + Limpiadores
MXHC,Cuidado del Hogar,Mixta Home Care-Cloro-Limpiadores-Ceras-Ambientadores
MXCB,Cuidado del Hogar,Mixta Limpiadores + Cloro
MXLB,Cuidado del Hogar,Mixta Detergentes + Cloro
MXLD,Cuidado del Hogar,Mixta Detergentes + Lavavajillas
CRTO,Cuidado del Hogar,Panitos + Papel Higienico
NAPK,Cuidado del Hogar,Servilletas
PLWF,Cuidado del Hogar,Film plastico e papel aluminio
SCOU,Cuidado del Hogar,Esponjas de Acero
SOFT,Cuidado del Hogar,Suavizantes de Ropa
STRM,Cuidado del Hogar,Quitamanchas-Desmanchadores
TOIP,Cuidado del Hogar,Papel Higienico
WIPE,Cuidado del Hogar,Panos de Limpieza
ANLG,OTC,Analgesicos-Painkillers
FSUP,OTC,Suplementos alimentares
GMED,OTC,Gastrointestinales-Efervescentes
VITA,OTC,Vitaminas y Calcio
nan,Otros,Categoria Desconocida
BATT,Otros,Pilas-Baterias
CGAS,Otros,Combustible Gas
PFHH,Otros,Panel Financiero de Hogares
PFIN,Otros,Panel Financiero de Hogares
INKC,Otros,Cartuchos de Tintas
PETF,Otros,Alimento para Mascota-Alim.p - perro - gato
TELE,Otros,Telecomunicaciones - Convergencia
TILL,Otros,Tickets - Till Rolls
TOBA,Otros,Tabaco - Cigarrillos
ADIP,Cuidado Personal,Incontinencia de Adultos
BSHM,Cuidado Personal,Shampoo Infantil
RAZO,Cuidado Personal,Maquinas de Afeitar
BDCR,Cuidado Personal,Cremas Corporales
CWIP,Cuidado Personal,Panos Humedos
COMB,Cuidado Personal,Cremas para Peinar
COND,Cuidado Personal,Acondicionador-Balsamo
CRHY,Cuidado Personal,Cross Category (Higiene)
CRPC,Cuidado Personal,Cross Category (Personal Care)
DEOD,Cuidado Personal,Desodorantes
DIAP,Cuidado Personal,Panales-Panales Desechables
FCCR,Cuidado Personal,Cremas Faciales
FTIS,Cuidado Personal,Panuelos Faciales
FEMI,Cuidado Personal,Proteccion Femenina-Toallas Femeninas
FRAG,Cuidado Personal,Fragancias
HAIR,Cuidado Personal,Cuidado del Cabello-Hair Care
HRCO,Cuidado Personal,Tintes para el Cabello-Tintes - Tintura - Tintes y Coloracion para el cabello
HREM,Cuidado Personal,Depilacion
HRST,Cuidado Personal,Alisadores para el Cabello
HSTY,Cuidado Personal,Fijadores para el Cabello-Modeladores-Gel-Fijadores para el cabello
HRTR,Cuidado Personal,Tratamientos para el Cabello
LINI,Cuidado Personal,leo Calcareo
MAKE,Cuidado Personal,Maquillaje-Cosmeticos
MEDS,Cuidado Personal,Jabon Medicinal
CRDT,Cuidado Personal,Panitos + Panales
MXMH,Cuidado Personal,Mixta Make Up+Tinturas
MOWA,Cuidado Personal,Enjuague Bucal-Refrescante Bucal
ORAL,Cuidado Personal,Cuidado Bucal
SPAD,Cuidado Personal,Protectores Femeninos
STOW,Cuidado Personal,Toallas Femininas
SHAM,Cuidado Personal,Shampoo
SHAV,Cuidado Personal,Afeitado-Crema afeitar-Locion de afeitar-Pord. Antes del afeitado
SKCR,Cuidado Personal,Cremas Faciales y Corporales-Cremas de Belleza - Cremas Cuerp y Faciales
SUNP,Cuidado Personal,Proteccion Solar
TALC,Cuidado Personal,Talcos-Talco para pies
TAMP,Cuidado Personal,Tampones Femeninos
TOIL,Cuidado Personal,Jabon de Tocador
TOOB,Cuidado Personal,Cepillos Dentales
TOOT,Cuidado Personal,Pastas Dentales
BAGS,Material Escolar,Morrales y MAletas Escoalres
CLPC,Material Escolar,Lapices de Colores
GRPC,Material Escolar,Lapices De Grafito
MRKR,Material Escolar,Marcadores
NTBK,Material Escolar,Cuadernos
SCHS,Material Escolar,tiles Escolares
CSTD,Diversos,Estudio de Categorias
CORP,Diversos,Corporativa
CROS,Diversos,Cross Category
CRBA,Diversos,Cross Category (Bebes)
CRBR,Diversos,"Cross Category (Desayuno)-Yogurt, Cereal, Pan y Queso"
CRDT,Diversos,Cross Category (Diet y Light)
CRDF,Diversos,Cross Category (Alimentos Secos)
CRFO,Diversos,Cross Category (Alimentos)
CRSA,Diversos,Cross Category (Salsas)-Mayonesas-Ketchup - Salsas Frias
CRSN,Diversos,Cross Category (Snacks)
DEMO,Diversos,Demo
FLSH,Diversos,Flash
HLVW,Diversos,Holistic View
COCP,Diversos,Mezcla para cafe instantaneo y crema no lactea
CRSN,Diversos,Mezclas nutricionales y suplementos
MULT,Diversos,Consolidado-Multicategory
PCHK,Diversos,Pantry Check
STCK,Diversos,Inventario
MIHC,Diversos,Leche y Cereales Calientes-Cereales Precocidos y Leche Liquida Blanca
FLWT,Alimentos,Agua Saborizada"""
categ = pd.read_csv(io.StringIO(CATEG_CSV_DATA), dtype={'cod': str, 'cest': str, 'cat': str})
categ['cod'] = categ['cod'].str.strip().str.upper()
CLIENT_NAME_SUFFIX_PATTERN = re.compile(r'[\s_-]*5w1h$', re.IGNORECASE)
#obtém o país,categoria,cesta e fabricante para template e ppt
base_dir = Path(__file__).resolve().parent
os.chdir(base_dir)
excel = select_excel_file(base_dir)
file = pd.ExcelFile(str(base_dir / excel))
W = file.sheet_names
#Obtém o pais cesta categoria fabricante marca e idioma para o qual se fará o estudo
excel_parts = excel.split('_')
land = pais.loc[pais.cod==int(excel_parts[0]),'pais'].iloc[0]
category_code = excel_parts[1].strip().upper()
cesta = categ.loc[categ.cod==category_code,'cest'].iloc[0]
cat = categ.loc[categ.cod==category_code,'cat'].iloc[0]
raw_client = excel_parts[2].rsplit('.', 1)[0]
sanitized_client = CLIENT_NAME_SUFFIX_PATTERN.sub('', raw_client).strip()
client = sanitized_client if sanitized_client else raw_client.strip()
lang= "P" if land=='Brasil' else "E"
modelo_path = base_dir / 'Modelo_5W1H.pptx'
if not modelo_path.exists():
    raise FileNotFoundError(colorize(f'No se encontro el template {modelo_path.name} en {base_dir}', COLOR_RED))
ppt= Presentation(str(modelo_path))
configure_cover_slide(ppt, lang)
brand = W[0][2:]
# Diccionario con las correspondencias de los números y su W respectivo
c_w = { 
    ('P','1')   : '1W - Quando?',
    ('P','2')   : '2W - Por quê?',
    ('P','3-1') : '3W - O quê? Tamanhos',
    ('P','3-2') : '3W - O quê? Marcas',
    ('P','3-3') : '3W - O quê? Sabores',
    ('P','4')   : '4W - Quem? NSE (Nível Socioeconômico)',
    ('P','5-1') : '5W - Onde? Regiões',
    ('P','5-2') : '5W - Onde? Canais',
    ('P','6')   : 'Players',
    ('P','6-1') : 'Players - Preço indexado',
    ('P','7-R') : '7W - Distribuição por regiões',
    ('P','7-NSE'): '7W - Distribuição por NSE',
    ('P','8')   : 'Intervalos de confiança',

    ('E','1')   : '1W - ¿Cuándo?',
    ('E','2')   : '2W - ¿Por qué?',
    ('E','3-1') : '3W - ¿Qué? Tamaños',
    ('E','3-2') : '3W - ¿Qué? Marcas',
    ('E','3-3') : '3W - ¿Qué? Sabores',
    ('E','4')   : '4W - ¿Quiénes? NSE (Nivel socioeconómico)',
    ('E','5-1') : '5W - ¿Dónde? Regiones',
    ('E','5-2') : '5W - ¿Dónde? Canales',
    ('E','6')   : 'Players',
    ('E','6-1') : 'Players - Precio indexado',
    ('E','7-R') : '7W - Distribución por regiones',
    ('E','7-NSE'): '7W - Distribución por NSE',
    ('E','8')   : 'Intervalos de confianza'
}

# Etiquetas de los slides
labels = {
    ('P','Data')   : 'Data',
    ('P','MAT')    : 'Avaliação em ano móvel acumulado',
    ('P','Var MAT'): 'Variação em ano móvel',
    ('P','comp')   : 'Concorrência no mercado de: ',
    ('P','dist')   : 'Distribuição',

    ('E','Data')   : 'Fecha',
    ('E','MAT')    : 'Evaluación en año móvil acumulado',
    ('E','Var MAT'): 'Variación en año móvil',
    ('E','comp')   : 'Competencia en el mercado de: ',
    ('E','dist')   : 'Distribución'
}

class SeriesConfig(NamedTuple):
    data: pd.DataFrame
    raw_tipo: str
    display_tipo: str
    pipeline: int
COLOR_TOKEN_STOPWORDS = {
    'gr', 'gr.', 'grupo', 'grup', 'compras', 'compra', 'ventas', 'venta',
    'canal', 'canales', 'channel', 'channels', 'marca', 'marcas', 'brand',
    'brands'
}
COLOR_SUFFIXES = ('.c', '.v', '_c', '_v', '-c', '-v')
COLOR_TOKEN_SPLIT_PATTERN = re.compile(r'[^0-9a-z]+')
COLOR_UNIT_SUFFIX_PATTERN = re.compile(
    r'(?:[_-](?:k|m|g|kg|grs?|gr|l|lt|lts|ml))+$',
    re.IGNORECASE
)
def normalize_color_text(value) -> str:
    """Normaliza texto para comparar etiquetas de color."""
    if value is None:
        return ''
    if not isinstance(value, str):
        value = str(value)
    normalized = unicodedata.normalize('NFKD', value)
    normalized = ''.join(ch for ch in normalized if not unicodedata.combining(ch))
    return normalized.lower().strip()
def strip_color_suffixes(text: str) -> str:
    """Elimina sufijos de compras/ventas y unidades antes de buscar color."""
    base = text.strip()
    for suffix in COLOR_SUFFIXES:
        if base.endswith(suffix):
            base = base[:-len(suffix)]
            break
    # Elimina sufijos de unidad (p. ej. _K, _M) para reutilizar el color de marca
    base = COLOR_UNIT_SUFFIX_PATTERN.sub('', base)
    return base.strip()
def generate_color_lookup_keys(label: str) -> list[str]:
    """Genera variantes de claves para mapear etiquetas a colores."""
    normalized = normalize_color_text(label)
    if not normalized:
        return []
    base = strip_color_suffixes(normalized)
    seen = set()
    keys: list[str] = []
    def _push(candidate: str):
        candidate = candidate.strip()
        if candidate and candidate not in seen:
            seen.add(candidate)
            keys.append(candidate)
    _push(base)
    collapsed = ' '.join(base.split())
    _push(collapsed)
    alnum = COLOR_TOKEN_SPLIT_PATTERN.sub('', base)
    _push(alnum)
    digit_key = ''.join(ch for ch in base if ch.isdigit())
    if digit_key:
        digit_key = digit_key.lstrip('0') or '0'
        _push(digit_key)
    raw_tokens = [tok for tok in COLOR_TOKEN_SPLIT_PATTERN.split(base) if tok]
    filtered_tokens = [tok for tok in raw_tokens if tok not in COLOR_TOKEN_STOPWORDS]
    if filtered_tokens:
        _push(' '.join(filtered_tokens))
        _push(''.join(filtered_tokens))
        if len(filtered_tokens) > 1:
            _push(' '.join(sorted(filtered_tokens)))
    if raw_tokens:
        _push(' '.join(raw_tokens))
    return keys
def build_color_lookup_dict(color_mapping: dict[str, str]) -> dict[str, str]:
    """Construye un diccionario de busqueda de colores."""
    lookup: dict[str, str] = {}
    if not color_mapping:
        return lookup
    for original_key, color_value in color_mapping.items():
        if color_value is None:
            continue
        normalized_key = str(original_key).strip()
        if not normalized_key or normalized_key.lower() == 'total':
            continue
        for lookup_key in generate_color_lookup_keys(normalized_key):
            if lookup_key and lookup_key not in lookup:
                lookup[lookup_key] = color_value
    return lookup
def lookup_color_for_label(label, lookup_dict: dict[str, str]) -> Optional[str]:
    """Busca el color asociado a una etiqueta normalizada."""
    if not lookup_dict or label is None:
        return None
    raw_key = str(label).strip()
    if not raw_key or raw_key.lower() == 'total':
        return None
    for lookup_key in generate_color_lookup_keys(raw_key):
        color_value = lookup_dict.get(lookup_key)
        if color_value:
            return color_value
    return None
def register_color_lookup(label, color_value: str, lookup_dict: dict[str, str], overwrite: bool = False) -> None:
    """Registra una etiqueta y su color en el lookup."""
    if label is None or not color_value:
        return
    raw_key = str(label).strip()
    if not raw_key or raw_key.lower() == 'total':
        return
    for lookup_key in generate_color_lookup_keys(raw_key):
        if not lookup_key:
            continue
        if overwrite or lookup_key not in lookup_dict:
            lookup_dict[lookup_key] = color_value

def hex_to_rgb_color(color_value: str) -> Optional[RGBColor]:
    """Convierte un color HEX a RGBColor."""
    try:
        rgb_float = mcolors.to_rgb(color_value)
    except (ValueError, TypeError):
        return None
    rgb_int = tuple(int(round(val * 255)) for val in rgb_float)
    return RGBColor(*rgb_int)

def hex_to_ansi_color(color_value: str) -> Optional[str]:
    """
    Convierte un color HEX a un codigo ANSI de 24 bits para la terminal.
    """
    try:
        r, g, b = (int(round(val * 255)) for val in mcolors.to_rgb(color_value))
    except (ValueError, TypeError):
        return None
    return f"\033[38;2;{r};{g};{b}m"

def assign_brand_palette_color(label: str, lookup_dict: dict[str, str], palette_sequence: Optional[list[str]] = None) -> Optional[str]:
    """Asigna un color de paleta a una marca si no existe en el lookup."""
    if not label:
        return None
    color_value = lookup_color_for_label(label, lookup_dict)
    if color_value:
        return color_value
    palette = palette_sequence if palette_sequence else TREND_COLOR_SEQUENCE
    if not palette:
        return None
    used_colors = {value for value in lookup_dict.values() if value}
    for candidate in palette:
        if candidate not in used_colors:
            register_color_lookup(label, candidate, lookup_dict, overwrite=False)
            return candidate
    fallback = palette[len(used_colors) % len(palette)]
    register_color_lookup(label, fallback, lookup_dict, overwrite=False)
    return fallback

def set_title_with_brand_color(
    text_frame,
    prefix_text: str,
    brand_text: Optional[str],
    suffix_text: str,
    font_size_inches: float,
    brand_lookup: dict[str, str],
    palette_sequence: Optional[list[str]] = None
) -> None:
    """Escribe un titulo con el segmento de marca coloreado."""
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = ''
    base_size = Inches(font_size_inches)
    def _add_run(text: str, color_hex: Optional[str] = None):
        if not text:
            return
        run = paragraph.add_run()
        run.text = text
        font = run.font
        font.bold = True
        font.size = base_size
        if color_hex:
            rgb_color = hex_to_rgb_color(color_hex)
            if rgb_color:
                font.color.rgb = rgb_color
    _add_run(prefix_text)
    brand_segment = (brand_text or '').strip()
    if brand_segment:
        brand_color = assign_brand_palette_color(brand_segment, brand_lookup, palette_sequence)
        _add_run(brand_segment, brand_color)
    _add_run(suffix_text)

BRAND_TITLE_COLOR_LOOKUP: dict[str, str] = {}
# Lookup independiente para categorías (segmento 8).
CATEGORY_TITLE_COLOR_LOOKUP: dict[str, str] = {}

def colorize_brand_for_terminal(
    label: str,
    brand_lookup: Optional[dict[str, str]] = None,
    palette_sequence: Optional[list[str]] = None
) -> str:
    """
    Devuelve el nombre de la marca coloreado con su color asignado en la paleta.
    """
    if not label:
        return ''
    lookup = brand_lookup if brand_lookup is not None else BRAND_TITLE_COLOR_LOOKUP
    brand_color = assign_brand_palette_color(label, lookup, palette_sequence)
    ansi_color = hex_to_ansi_color(brand_color) if brand_color else None
    if ansi_color:
        return colorize(label, ansi_color)
    return label

def build_terminal_label(sheet_name: str, lang: str, category_name: str) -> Optional[str]:
    """
    Devuelve un label corto y legible (ej: '3W Marcas - X').
    """
    if not sheet_name:
        return None
    sheet_clean = sheet_name.strip()
    if not sheet_clean:
        return None
    distribution_match = DISTRIBUTION_SHEET_PATTERN.match(sheet_clean)
    first_char = sheet_clean[0]
    suffix_char = sheet_clean[-1] if len(sheet_clean) >= 1 else ''
    brand_label = ''
    cw_key = None
    # Distribución (segmento 7) via patrón 7_*_* (ej.: R, NSE, WIFEAGE, FAMILY_SIZE)
    if distribution_match:
        dist_kind_raw = distribution_match.group(2)
        dist_kind = str(dist_kind_raw).upper().strip()
        pretty_dist = dist_kind.replace('_', ' ').replace('-', ' ')
        pretty_dist = ' '.join(word.capitalize() for word in pretty_dist.split())
        if dist_kind == 'R':
            cw_key = '7-R'
            step_label = '7W Distribución Regiones'
        elif dist_kind == 'NSE':
            cw_key = '7-NSE'
            step_label = '7W Distribución NSE'
        else:
            step_label = f'7W Distribución {pretty_dist or dist_kind}'
        brand_label = distribution_match.group(1).replace('.', ' ').strip()
        brand_label = distribution_match.group(1).replace('.', ' ').strip()
    elif sheet_clean.startswith('6-1') or sheet_clean.startswith('6_1'):
        cw_key = '6-1'
        brand_label = category_name
    elif first_char == '6':
        cw_key = '6'
        brand_label = category_name
    elif first_char == '8':
        cw_key = '8'
        brand_label = sheet_clean[2:].replace('_', ' ').strip()
    elif first_char == '5':
        cw_key = f"5-{suffix_char}" if suffix_char in {'1', '2'} else '5'
        brand_label = sheet_clean[2:-2].strip() if len(sheet_clean) > 3 else ''
    elif first_char == '4':
        cw_key = '4'
        brand_label = sheet_clean[2:].strip()
    elif first_char == '3':
        cw_key = f"3-{suffix_char}" if suffix_char in {'1', '2', '3'} else '3'
        brand_label = sheet_clean[2:-2].strip() if len(sheet_clean) > 3 else sheet_clean[2:].strip()
    elif first_char in {'2', '1'}:
        cw_key = first_char
        brand_label = sheet_clean[2:].strip()
    step_label = c_w.get((lang, cw_key), step_label if cw_key is None else cw_key if cw_key else step_label)
    if not step_label:
        return None
    # Para mantener el estilo corto en terminal, eliminamos el " - " del texto base
    if ' - ' in step_label:
        code, desc = step_label.split(' - ', 1)
        step_label = f"{code} {desc}"
    brand_terminal = colorize_brand_for_terminal(brand_label) if brand_label else ''
    detail = f" - {brand_terminal}" if brand_terminal else ''
    return f"{step_label}{detail}"

def build_terminal_progress_message(sheet_name: str, lang: str, category_name: str) -> Optional[str]:
    """
    Mensaje de progreso claro (ej: 'Generando 3W Marcas - X (hoja 3_Marca_P)').
    """
    label = build_terminal_label(sheet_name, lang, category_name)
    if not label:
        return None
    return f"Generando {label} (hoja {sheet_name})"

def build_terminal_done_message(sheet_name: str, lang: str, category_name: str) -> Optional[str]:
    """Devuelve el mensaje final de progreso en terminal."""
    label = build_terminal_label(sheet_name, lang, category_name)
    if not label:
        return None
    return f"Listo {label}"

def _extract_pipeline(col_name: str) -> int:
    """Extrae el numero de pipeline desde el encabezado."""
    if not isinstance(col_name, str):
        return 0
    parts = [part for part in col_name.split('_') if part.isdigit()]
    if parts:
        return int(parts[0])
    digits = ''.join(ch for ch in col_name if ch.isdigit())
    return int(digits) if digits else 0
def _is_separator_column(col) -> bool:
    """Detecta columnas separadoras vacias o 'Unnamed'."""
    if col is None:
        return True
    if not isinstance(col, str):
        return False
    stripped = col.strip()
    return not stripped or stripped.lower().startswith('unnamed')
def _normalize_header_text(value: str) -> str:
    normalized = unicodedata.normalize('NFKD', str(value))
    normalized = ''.join(ch for ch in normalized if not unicodedata.combining(ch))
    return normalized.lower().strip()
def _find_compras_header_row(df: pd.DataFrame, include_table: bool = True) -> Optional[int]:
    """Busca la fila con encabezado de Compras (y opcionalmente table)."""
    if df is None or df.empty:
        return None
    compras_idx = None
    table_idx = None
    for idx, row in df.iterrows():
        for value in row:
            if isinstance(value, str) and 'compras' in _normalize_header_text(value):
                compras_idx = idx
                break
            if include_table and isinstance(value, str) and 'table' in _normalize_header_text(value):
                if table_idx is None:
                    table_idx = idx
        if compras_idx is not None:
            break
    if compras_idx is not None:
        return compras_idx
    return table_idx


def _find_table_anchor_row(df: pd.DataFrame) -> Optional[int]:
    if df is None or df.empty or df.shape[1] == 0:
        return None
    first_col = df.iloc[:, 0]
    for idx, value in first_col.items():
        if isinstance(value, str) and 'table' in _normalize_header_text(value):
            return idx
    return None


def _find_first_date_row(
    df: pd.DataFrame,
    start_idx: int,
    date_format: str
) -> Optional[int]:
    if df is None or df.empty or df.shape[1] == 0:
        return None
    first_col = df.iloc[start_idx:, 0]
    parsed = pd.to_datetime(first_col, format=date_format, errors='coerce')
    non_empty = pd.Series(True, index=first_col.index)
    if first_col.dtype == object:
        non_empty = first_col.astype(str).str.strip().ne('')
    valid = parsed.notna() & non_empty
    if valid.any():
        return valid[valid].index[0]
    return None


def _find_header_row_near_date(df: pd.DataFrame, date_row_idx: int) -> Optional[int]:
    if df is None or df.empty:
        return None
    for idx in range(date_row_idx, max(date_row_idx - 6, -1), -1):
        row = df.iloc[idx]
        for value in row:
            if isinstance(value, str):
                text = _normalize_header_text(value)
                if 'compras' in text or 'ventas' in text:
                    return idx
    if date_row_idx > 0:
        return date_row_idx - 1
    return None


def _row_has_mat_headers(row: pd.Series) -> bool:
    for value in row:
        if isinstance(value, str) and 'mat' in _normalize_header_text(value):
            return True
    return False
# Segmentos 3-6: parsing de hojas con Compras/Ventas y series.
def parse_sheet_with_compras_header(
    excel_file: pd.ExcelFile,
    sheet_name: str,
    include_table: bool = True
) -> pd.DataFrame:
    """
    Carga la hoja buscando la fila con 'Compras' como encabezado para tolerar filas vacias o
    textos previos; si no se encuentra, usa el encabezado por defecto.
    """
    raw_df = excel_file.parse(sheet_name, header=None)
    header_row = _find_compras_header_row(raw_df, include_table=include_table)
    if header_row is None:
        return excel_file.parse(sheet_name)
    return excel_file.parse(sheet_name, header=header_row)


def parse_sheet_with_date_fallback(
    excel_file: pd.ExcelFile,
    sheet_name: str,
    date_format: str = '%b-%y  '
) -> pd.DataFrame:
    """
    Intenta leer la hoja con encabezado Compras/Ventas.
    Si falla, busca la primera fecha debajo de un ancla "table" y reintenta.
    """
    raw_df = excel_file.parse(sheet_name, header=None)
    try:
        base_df = parse_sheet_with_compras_header(excel_file, sheet_name, include_table=False)
        return ensure_date_column(base_df, sheet_name=sheet_name, date_format=date_format)
    except ValueError:
        table_idx = _find_table_anchor_row(raw_df)
        start_idx = table_idx + 1 if table_idx is not None else 0
        date_row_idx = _find_first_date_row(raw_df, start_idx, date_format)
        if date_row_idx is None:
            if table_idx is not None and _row_has_mat_headers(raw_df.iloc[table_idx]):
                raise ValueError(
                    f"La hoja {sheet_name} contiene un encabezado 'table' con periodos MAT, "
                    "pero no se encontraron fechas en la primera columna."
                )
            raise
        header_idx = _find_header_row_near_date(raw_df, date_row_idx)
        if header_idx is None:
            raise ValueError(
                f"No se encontro encabezado de Compras/Ventas alrededor de las fechas en la hoja {sheet_name}."
            )
        df = excel_file.parse(sheet_name, header=header_idx)
        return ensure_date_column(df, sheet_name=sheet_name, date_format=date_format)
def ensure_date_column(df: pd.DataFrame, sheet_name: Optional[str] = None, date_format: str = '%b-%y  ') -> pd.DataFrame:
    """
    Valida y convierte la primera columna a fecha (formato MMM-YY).
    Lanza un ValueError si no hay fechas legibles para evitar graficos con datos corruptos.
    """
    if df is None or df.empty or df.shape[1] == 0:
        return df
    first_col = df.iloc[:, 0]
    if np.issubdtype(first_col.dtype, np.datetime64):
        return df
    context = f" en la hoja {sheet_name}" if sheet_name else ""
    parsed = pd.to_datetime(first_col, format=date_format, errors='coerce')
    non_empty = pd.Series(True, index=first_col.index)
    if first_col.dtype == object:
        non_empty = first_col.astype(str).str.strip().ne('')
    invalid_count = int((parsed.isna() & non_empty).sum())
    if invalid_count > 0 or parsed.notna().sum() == 0:
        raise ValueError(
            f"No se pudieron leer fechas validas en la primera columna{context}. "
            "Use el formato MMM-YY (ej.: Ene-24)."
        )
    df_copy = df.copy()
    df_copy.iloc[:, 0] = parsed
    return df_copy
# Segmento 7: distribuciones (R/NSE/otros cortes).
DISTRIBUTION_SHEET_PATTERN = re.compile(r'^7[-_](.+?)[-_](.+)$', re.IGNORECASE)
def _normalize_simple(text: str) -> str:
    """Normaliza texto a minusculas y sin tildes."""
    normalized = unicodedata.normalize('NFKD', str(text))
    normalized = ''.join(ch for ch in normalized if not unicodedata.combining(ch))
    return normalized.lower().strip()
def _is_distribution_total_row(cat_value, row_values: pd.Series) -> bool:
    """
    Detecta filas de totales/resumen que no deben graficarse en distribuciones.
    Conserva filas como "Total Guatemala" o similares siempre que los valores
    no sean uniformes.
    """
    if not isinstance(cat_value, str):
        return False
    normalized = _normalize_simple(cat_value)
    if not normalized or 'total' not in normalized:
        return False
    numeric_vals = pd.to_numeric(row_values, errors='coerce').dropna()
    if numeric_vals.empty:
        return True
    unique_vals = pd.Series(np.round(numeric_vals, 3)).unique()
    if len(unique_vals) == 1:
        return True
    if np.allclose(numeric_vals, 100, atol=0.5):
        return True
    return False
def _find_distribution_header_indices(raw_df: pd.DataFrame) -> tuple[Optional[int], Optional[int]]:
    """
    Busca una fila de encabezados cuando no hay ancla "table/Compras".
    Heuristica: una fila con al menos dos celdas no vacías (col 1+) y la fila
    siguiente con al menos un número, asumiendo que es cabecera + datos.
    """
    if raw_df is None or raw_df.empty:
        return None, None
    max_scan = min(raw_df.shape[0] - 1, 12)
    for idx in range(max_scan):
        row_rest = raw_df.iloc[idx, 1:]
        non_empty = row_rest.dropna().astype(str).str.strip()
        if (non_empty != '').sum() < 2:
            continue
        next_row = raw_df.iloc[idx + 1, 1:] if idx + 1 < raw_df.shape[0] else None
        if next_row is None:
            continue
        numeric_count = pd.to_numeric(next_row, errors='coerce').count()
        if numeric_count >= 1:
            return idx, idx + 1
    return None, None
def parse_distribution_sheet(excel_file: pd.ExcelFile, sheet_name: str, client_name: str) -> tuple[pd.DataFrame, str]:
    """
    Extrae la distribucion (R/NSE) desde las hojas 7_*_*.
    Mantiene todas las series con datos y usa los encabezados de la fila 2 como nombres de marca.
    """
    raw_df = excel_file.parse(sheet_name, header=None)
    if raw_df.shape[0] < 3 or raw_df.shape[1] < 2:
        return pd.DataFrame(columns=['Categoria']), ''
    # Buscar la fila "table" (en la primera columna) para alinear la lectura, similar al segmento 2
    table_anchor_idx = None
    def _cell_has_anchor(value) -> bool:
        if value is None:
            return False
        text = str(value)
        if not text:
            return False
        normalized = unicodedata.normalize('NFKD', text)
        normalized = ''.join(ch for ch in normalized if not unicodedata.combining(ch))
        normalized = normalized.lower()
        return 'table' in normalized or 'compras' in normalized
    for idx in range(raw_df.shape[0]):
        first_cell = raw_df.iat[idx, 0] if raw_df.shape[1] > 0 else None
        if _cell_has_anchor(first_cell):
            table_anchor_idx = idx
            break
    header_idx = None
    data_start_idx = None
    if table_anchor_idx is not None and table_anchor_idx + 1 < raw_df.shape[0]:
        header_idx = table_anchor_idx
        data_start_idx = table_anchor_idx + 1
    if header_idx is None:
        header_idx, data_start_idx = _find_distribution_header_indices(raw_df)
    if header_idx is None:
        header_idx, data_start_idx = 1, 2
    if data_start_idx is None or data_start_idx >= raw_df.shape[0]:
        return pd.DataFrame(columns=['Categoria']), ''
    header_row = raw_df.iloc[header_idx, 1:]
    data_block = raw_df.iloc[data_start_idx:, :]
    categories = data_block.iloc[:, 0].astype(str).str.strip()
    df = pd.DataFrame({'Categoria': categories})
    series_headers: list[str] = []
    for idx in range(1, data_block.shape[1]):
        serie = pd.to_numeric(data_block.iloc[:, idx], errors='coerce')
        if serie.dropna().empty:
            continue
        header_value = header_row.iloc[idx - 1] if (idx - 1) < len(header_row) else ''
        header_text = str(header_value).strip() if pd.notna(header_value) else ''
        if not header_text:
            fallback_name = None
            for candidate in data_block.iloc[:, idx]:
                if isinstance(candidate, str) and candidate.strip():
                    fallback_name = candidate.strip()
                    break
            header_text = fallback_name if fallback_name else f"Serie {idx}"
        df[header_text] = serie
        series_headers.append(header_text)
    df = df[df['Categoria'].str.len() > 0]
    df = df[~df['Categoria'].str.lower().eq('nan')]
    value_columns = [col for col in df.columns if col != 'Categoria']
    if not value_columns:
        return pd.DataFrame(columns=['Categoria']), ''
    df = df.dropna(subset=value_columns, how='all')
    df = df[~df.apply(lambda row: _is_distribution_total_row(row['Categoria'], row[value_columns]), axis=1)]
    def _is_total_label(label: str) -> bool:
        return isinstance(label, str) and 'total' in label.lower()
    total_columns = [col for col in value_columns if _is_total_label(col)]
    non_total_columns = [col for col in value_columns if col not in total_columns]
    ordered_columns = total_columns + non_total_columns
    df = df[['Categoria'] + ordered_columns]
    ordered_headers = [col for col in ordered_columns if col]
    return df, ', '.join(ordered_headers)
def plot_distribution_chart(
    df: pd.DataFrame,
    c_fig: int,
    color_lookup: Optional[dict[str, str]] = None,
    title: Optional[str] = None,  # se conserva la firma para no romper llamadas previas
    adaptive_label_size: bool = True,
    target_width_in: Optional[float] = None,
    max_height_in: Optional[float] = None
) -> tuple[io.BytesIO, tuple[float, float]]:
    """
    Crea un grafico de barras agrupadas con etiquetas numericas para distribuciones R/NSE.
    Permite escalar el tamano de las etiquetas de barra cuando hay muchas categorias/series.
    Devuelve el stream PNG y el tamano de la figura en pulgadas para insertarlo sin reescalar.
    """
    if df.empty or df.shape[1] <= 1:
        fig, ax = plt.subplots(num=c_fig)
        ax.axis('off')
        fig_size = fig.get_size_inches()
        return figure_to_stream(fig), (float(fig_size[0]), float(fig_size[1]))
    categories = df['Categoria'].tolist()
    series_names = [col for col in df.columns if col != 'Categoria']
    filtered_series: list[tuple[str, pd.Series]] = []
    for serie_name in series_names:
        values_raw = pd.to_numeric(df[serie_name], errors='coerce')
        if values_raw.dropna().empty:
            continue
        values = values_raw.fillna(0)
        if values.replace(0, np.nan).dropna().empty:
            continue
        filtered_series.append((serie_name, values))
    if not categories or not filtered_series:
        fig, ax = plt.subplots(num=c_fig, figsize=(8, 4), dpi=DEFAULT_EXPORT_DPI)
        ax.axis('off')
        fig_size = fig.get_size_inches()
        return figure_to_stream(fig), (float(fig_size[0]), float(fig_size[1]))
    # Descarta categorias cuyo valor sea ~100 para todas las series (evita barras "planas" de referencia)
    keep_mask = []
    tol_100 = 0.1  # tolerancia para considerar un valor como 100
    for cat_idx in range(len(categories)):
        vals_at_cat = []
        for _, serie_vals in filtered_series:
            try:
                val = float(serie_vals.iloc[cat_idx])
            except Exception:
                continue
            if np.isfinite(val):
                vals_at_cat.append(val)
        if vals_at_cat:
            max_v = max(vals_at_cat)
            min_v = min(vals_at_cat)
            if abs(max_v - 100.0) <= tol_100 and abs(min_v - 100.0) <= tol_100:
                keep_mask.append(False)
                continue
        else:
            keep_mask.append(False)
            continue
        keep_mask.append(True)
    if not any(keep_mask):
        fig, ax = plt.subplots(num=c_fig, figsize=(8, 4), dpi=DEFAULT_EXPORT_DPI)
        ax.axis('off')
        fig_size = fig.get_size_inches()
        return figure_to_stream(fig), (float(fig_size[0]), float(fig_size[1]))
    keep_indices = [idx for idx, keep in enumerate(keep_mask) if keep]
    categories = [categories[idx] for idx in keep_indices]
    filtered_series = [
        (name, pd.Series([vals.iloc[idx] for idx in keep_indices])) for name, vals in filtered_series
    ]
    n_series = len(filtered_series)
    x = np.arange(len(categories))
    width = 0.8 / max(n_series, 1)
    base_fig_width = max(10.0, min(16.0, 1.3 * len(categories)))
    base_fig_height = max(6.5, min(12.0, 2.5 + 0.9 * len(categories)))
    fig_width = target_width_in if target_width_in else base_fig_width
    if fig_width <= 0:
        fig_width = base_fig_width
    fig_height = base_fig_height
    if max_height_in is not None and max_height_in > 0:
        fig_height = min(fig_height, max_height_in)
    default_label_font_size = 9
    label_font_size = default_label_font_size
    if adaptive_label_size:
        total_bars = max(len(categories) * n_series, 1)
        density = total_bars / max(fig_width, 1.0)  # barras por pulgada aproximada
        if density > 5.5:
            label_font_size = 6
        elif density > 4.0:
            label_font_size = 7
        elif density > 3.0:
            label_font_size = 8
        elif density < 1.5:
            label_font_size = min(11, default_label_font_size + 1)
    label_font_size = max(6, min(11, int(round(label_font_size))))
    fig, ax = plt.subplots(num=c_fig, figsize=(fig_width, fig_height), dpi=DEFAULT_EXPORT_DPI)
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F9FAFB')
    ax.set_axisbelow(True)
    palette_lookup = {}
    def _is_total_label(label: str) -> bool:
        return isinstance(label, str) and 'total' in label.lower()
    def _is_reference_total_series(label: str, serie_values: pd.Series, tol: float = 0.5) -> bool:
        """
        Considera serie de referencia solo si contiene 'total' y sus valores ~100.
        Evita pintar de negro series con 'total' en el nombre que no son totales reales.
        """
        if not _is_total_label(label):
            return False
        numeric_vals = pd.to_numeric(serie_values, errors='coerce')
        numeric_vals = numeric_vals[np.isfinite(numeric_vals)]
        if numeric_vals.empty:
            return False
        return np.allclose(numeric_vals, 100.0, atol=tol)
    if TREND_COLOR_PALETTE:
        palette_values = [
            TREND_COLOR_PALETTE[key]
            for key in sorted(
                TREND_COLOR_PALETTE.keys(),
                key=lambda k: int(k.split('_')[1]) if '_' in k and k.split('_')[1].isdigit() else k
            )
        ]
    else:
        palette_values = TREND_COLOR_SEQUENCE
    if not palette_values:
        palette_values = [mcolors.to_hex(c) for c in plt.get_cmap('tab20').colors]
    palette_index = 0
    bars_by_series = []
    legend_handles = []
    legend_labels = []
    max_val = 0.0
    for idx, (serie_name, values) in enumerate(filtered_series):
        is_total_series = _is_reference_total_series(serie_name, values)
        if idx == 0 or is_total_series:
            color_val = REFERENCE_SERIES_COLOR
        else:
            color_val = palette_values[palette_index % len(palette_values)]
            palette_index += 1
        register_color_lookup(serie_name, color_val, palette_lookup, overwrite=True)
        offset = (idx - (n_series - 1) / 2) * width
        bars = ax.bar(
            x + offset,
            values,
            width,
            label=serie_name,
            color=color_val,
            edgecolor='none',
            alpha=0.9
        )
        bars_by_series.append((bars, values))
        legend_handles.append(bars)
        legend_labels.append(serie_name)
        serie_max = values.max() if not values.empty else 0
        if np.isfinite(serie_max):
            max_val = max(max_val, float(serie_max))
    ax.set_xticks(x)
    tick_font_size = 10
    ax.set_xticklabels(categories, rotation=30, ha='right', va='top', rotation_mode='anchor', fontweight='normal', fontsize=tick_font_size)
    ax.tick_params(axis='x', pad=2)
    ax.tick_params(axis='y', labelsize=tick_font_size, pad=2)
    ax.set_ylabel('%', fontweight='bold', fontsize=11, labelpad=10)
    ax.grid(axis='y', linestyle='--', alpha=0.45, color='#D9D9D9')
    ax.margins(x=0.02)
    for spine in ax.spines.values():
        spine.set_visible(False)
    limit = max_val * 1.15 if max_val > 0 else None
    if limit:
        ax.set_ylim(0, limit)
    ax.set_xlim(-0.5, len(categories) - 0.5)
    ref_idx = 0 if filtered_series else None
    highlight_map: dict[tuple[int, int], str] = {}
    if ref_idx is not None and len(filtered_series) > 1:
        for cat_idx in range(len(categories)):
            try:
                ref_val = float(bars_by_series[ref_idx][1].iloc[cat_idx])
            except Exception:
                ref_val = None
            if ref_val is None or not np.isfinite(ref_val):
                continue
            diffs = []
            for s_idx, (_, vals) in enumerate(bars_by_series):
                if s_idx == ref_idx:
                    continue
                try:
                    val = float(vals.iloc[cat_idx])
                except Exception:
                    continue
                if not np.isfinite(val):
                    continue
                diffs.append((s_idx, val - ref_val))
            if not diffs:
                continue
            max_entry = max(diffs, key=lambda t: t[1])
            min_entry = min(diffs, key=lambda t: t[1])
            if max_entry[1] > 0:
                highlight_map[(max_entry[0], cat_idx)] = BAR_LABEL_COLOR_POS_ALT
            if min_entry[1] < 0:
                highlight_map[(min_entry[0], cat_idx)] = BAR_LABEL_COLOR_NEG_ALT
    label_padding = max(max_val * 0.01, 0.05 * (label_font_size / default_label_font_size))
    max_label_y = 0.0
    for serie_idx, (bars, values) in enumerate(bars_by_series):
        for cat_idx, (bar, value) in enumerate(zip(bars, values)):
            label_val = f"{value:.1f}" if abs(value) >= 10 else f"{value:.2f}"
            color_val = highlight_map.get((serie_idx, cat_idx), TABLE_TEXT_PRIMARY)
            proposed_y = bar.get_height() + label_padding
            if proposed_y > max_label_y:
                max_label_y = proposed_y
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                proposed_y,
                label_val,
                ha='center',
                va='bottom',
                rotation=0,
                fontsize=label_font_size,
                fontweight='normal',
                color=color_val,
                clip_on=False
            )
    if max_label_y > 0:
        ylim_bottom, ylim_top = ax.get_ylim()
        extra_top = max(max_val * 0.03, label_padding)
        target_top = max(ylim_top, max_label_y + extra_top)
        ax.set_ylim(ylim_bottom, target_top)
    legend = ax.legend(
        legend_handles,
        legend_labels,
        frameon=True,
        loc='upper center',
        bbox_to_anchor=(0.5, 1.02),
        ncol=max(1, len(legend_labels)),
        borderaxespad=0.2,
        columnspacing=1.0,
        handlelength=1.4,
        handletextpad=0.4
    )
    if legend:
        legend.get_frame().set_facecolor('#FFFFFF')
        legend.get_frame().set_edgecolor('#CCCCCC')
        legend.get_frame().set_alpha(0.95)
    # Ajuste dinamico de margenes: conserva la ubicacion de la leyenda arriba
    # y expande el area util del eje para evitar que el contenido quede comprimido.
    fig.canvas.draw()
    fig_height_in = fig.get_size_inches()[1] or 1.0
    legend_height_frac = 0.0
    if legend:
        renderer = fig.canvas.get_renderer()
        if renderer:
            bbox = legend.get_window_extent(renderer=renderer)
            legend_height_frac = bbox.height / (fig_height_in * fig.dpi)
    # Margen inferior estimado segun el alto de las etiquetas del eje X (rotadas)
    xtick_labels = ax.get_xticklabels()
    bottom_margin = 0.1
    if xtick_labels:
        max_label_pts = max(label.get_size() for label in xtick_labels if label.get_text())
        estimated_bottom = (max_label_pts * 1.6) / (fig_height_in * 72.0)
        bottom_margin = max(0.06, min(0.16, estimated_bottom))
    # Margen superior reducido dinamicamente para dejar espacio a la leyenda
    top_margin = 0.93 - legend_height_frac * 1.2
    top_margin = max(0.84, min(0.95, top_margin))
    # Aseguramos consistencia con los bar charts de otras secciones (3-6)
    fig.subplots_adjust(left=0.06, right=0.98, top=top_margin, bottom=bottom_margin)
    try:
        # Afina el layout sin mover la leyenda
        fig.tight_layout(rect=[0.05, bottom_margin + 0.01, 0.97, top_margin - 0.01])
    except Exception:
        pass
    return figure_to_stream(fig), (fig_width, fig_height)

def split_compras_ventas(df: pd.DataFrame, sheet_name: Optional[str] = None) -> tuple[list[pd.DataFrame], int]:
    """Divide la hoja en bloques de compras/ventas y detecta pipeline."""
    warning_context = f" en la hoja {sheet_name}" if sheet_name else ""
    normalized_columns = [
        str(col).strip().lower() if isinstance(col, str) else ''
        for col in df.columns
    ]
    has_compras_header = any('compra' in col for col in normalized_columns)
    has_ventas_header = any('venta' in col for col in normalized_columns)
    if not has_compras_header:
        print_colored(
            f"No se encontro el encabezado de \"Compras\"{warning_context}. Se usaran los datos disponibles.",
            COLOR_RED
        )
    ventas_idx = None
    pipeline_header = None
    for idx, col in enumerate(df.columns):
        if isinstance(col, str) and 'ventas' in col.lower():
            ventas_idx = idx
            pipeline_header = col
            break
    fallback_separator_idx = None
    for idx in range(len(df.columns) - 1):
        if _is_separator_column(df.columns[idx]):
            next_idx = idx + 1
            if next_idx < len(df.columns) and not _is_separator_column(df.columns[next_idx]):
                fallback_separator_idx = idx
                break
    suffix_hint = any(
        isinstance(col, str) and col.endswith('.1')
        for col in df.columns
    )
    fallback_used = False
    if ventas_idx is None and fallback_separator_idx is not None:
        candidate_idx = fallback_separator_idx + 1
        if candidate_idx < len(df.columns):
            ventas_idx = candidate_idx
            pipeline_header = df.columns[candidate_idx]
            fallback_used = True
    pipeline = _extract_pipeline(pipeline_header) if pipeline_header is not None else 0
    has_pipeline_digits = isinstance(pipeline_header, str) and any(ch.isdigit() for ch in pipeline_header)
    if ventas_idx is None:
        if (fallback_separator_idx is not None or suffix_hint) and not has_ventas_header:
            print_colored(
                f"No se encontro el encabezado de \"Ventas\"{warning_context}. No se pudo dividir la hoja.",
                COLOR_RED
            )
        return [df], 0
    if fallback_used and not has_ventas_header:
        print_colored(
            f"No se encontro el encabezado de \"Ventas\"{warning_context}. Se detecto el bloque por la columna en blanco y se graficara con pipeline = 0.",
            COLOR_RED
        )
        pipeline = 0
    if pipeline == 0 and pipeline_header is not None and not has_pipeline_digits:
        print_colored(
            f"No se definio un pipeline numerico en la columna '{pipeline_header}'{warning_context}. Se usara pipeline = 0.",
            COLOR_RED
        )
        pipeline = 0
    separator_idx = ventas_idx - 1
    has_separator = separator_idx >= 0 and _is_separator_column(df.columns[separator_idx])
    compras_end = separator_idx if has_separator else ventas_idx
    if compras_end <= 0:
        compras_end = ventas_idx
    compras = df.iloc[:, :compras_end].copy()
    ventas = df.iloc[:, ventas_idx:].copy()
    if not compras.empty:
        compras = compras.loc[:, ~pd.Series(compras.columns).apply(_is_separator_column).to_numpy()]
        if len(compras.columns) > 1:
            compras.columns = [compras.columns[0]] + [f"{col}.C" for col in compras.columns[1:]]
    if not ventas.empty:
        ventas = ventas.loc[:, ~pd.Series(ventas.columns).apply(_is_separator_column).to_numpy()]
        ventas_cols = list(ventas.columns)
        if ventas_cols:
            first = ventas_cols[0]
            base_name = first[:-2] if isinstance(first, str) and len(first) > 2 and first[-2] == '_' else first
            ventas_cols[0] = base_name
            ventas_cols[1:] = [str(col).replace('.1', '.V') for col in ventas_cols[1:]]
            ventas.columns = ventas_cols
    if ventas.shape[1] <= 1 or ventas.iloc[:, 1:].isna().all().all():
        return [compras if not compras.empty else df], 0
    df_list: list[pd.DataFrame] = []
    if not compras.empty:
        df_list.append(compras)
    df_list.append(ventas)
    return df_list, pipeline


def _detect_ventas_pipeline(columns: pd.Index) -> int:
    for col in columns:
        if isinstance(col, str) and 'ventas' in col.lower():
            return _extract_pipeline(col)
    return 0
def prepare_series_configs(df_list, lang, p_ventas, sheet_name: Optional[str] = None):
    """Prepara configuraciones de series para graficar."""
    configs = []
    for original_df in df_list:
        if original_df is None or original_df.empty or original_df.shape[1] == 0:
            context = f" en la hoja {sheet_name}" if sheet_name else ""
            print_colored(f"No se encontraron datos graficables{context}. Se omite esta seccion.", COLOR_RED)
            continue
        df_local = ensure_date_column(original_df.copy(), sheet_name=sheet_name)
        if df_local.shape[1] == 0:
            context = f" en la hoja {sheet_name}" if sheet_name else ""
            print_colored(f"No hay columnas disponibles{context}. Se omite esta seccion.", COLOR_RED)
            continue
        raw_tipo = str(df_local.columns[0])
        df_local.rename(columns={df_local.columns[0]: labels[(lang,'Data')]}, inplace=True)
        is_compras = 'compras' in raw_tipo.lower()
        if not is_compras and len(df_local.columns) > 1:
            is_compras = '.c' in str(df_local.columns[1]).lower()
        display_tipo = 'Compras' if is_compras else 'Ventas'
        pipeline = 0 if is_compras else p_ventas
        configs.append(SeriesConfig(df_local, raw_tipo, display_tipo, pipeline))
    return configs
def format_origin_label(display_tipo: Optional[str], lang: str) -> str:
    """
    Normaliza la etiqueta de origen de datos (Compras/Ventas) para textos como 'Corte a:'.
    """
    if not display_tipo:
        return ""
    base = str(display_tipo).strip()
    if not base:
        return ""
    normalized = base.lower()
    if normalized.startswith('comp') or 'compra' in normalized or normalized.endswith('.c'):
        return 'Compras'
    if normalized.startswith('vent') or 'venta' in normalized or normalized.endswith('.v') or 'sell-out' in normalized:
        return 'Vendas' if lang == 'P' else 'Ventas'
    if 'sell-in' in normalized:
        return 'Compras'
    return base
def extract_players_base_key(sheet_name: str) -> Optional[str]:
    """Deriva la clave base para hojas de jugadores (segmento 6)."""
    if not isinstance(sheet_name, str):
        return None
    if not sheet_name or not sheet_name.startswith('6'):
        return None
    parts = [part for part in sheet_name.split('_') if part]
    if not parts:
        return None
    if parts[-1].isdigit():
        if len(parts) > 1:
            return '_'.join(parts[:-1])
        return parts[0]
    return '_'.join(parts)
# Precio indexado: identifica hojas y arma la tabla indexada con contexto Players.
def is_price_index_sheet(sheet_name: str) -> bool:
    """Indica si la hoja corresponde a precio indexado."""
    if not isinstance(sheet_name, str):
        return False
    if not sheet_name.startswith('6'):
        return False
    parts = [part for part in sheet_name.split('_') if part]
    if len(parts) < 3:
        return False
    last_part = parts[-1]
    return last_part == '1'
def prepare_price_index_dataframe(series_config: SeriesConfig, top10_columns: list[str]) -> tuple[pd.DataFrame, list[str], pd.DataFrame]:
    """Prepara dataframes para grafico y tabla de precio indexado."""
    df_price = series_config.data.copy()
    if df_price.empty or df_price.shape[1] <= 1:
        return pd.DataFrame(), [], pd.DataFrame()
    date_col = df_price.columns[0]
    data_cols = list(df_price.columns[1:])
    if not data_cols:
        return pd.DataFrame(), [], pd.DataFrame()
    total_col = next(
        (col for col in data_cols if 'total' in str(col).strip().lower()),
        data_cols[0]
    )
    df_price[total_col] = pd.to_numeric(df_price[total_col], errors='coerce')
    df_price = df_price[df_price[total_col].notna()]
    df_price = df_price[df_price[total_col] > 0]
    if df_price.empty:
        return pd.DataFrame(columns=[date_col]), [], pd.DataFrame(columns=[date_col])
    normalized_top10 = {str(col).strip().lower() for col in top10_columns if str(col).strip()}
    candidate_columns = []
    filtered_columns = []
    for col in data_cols:
        if col == total_col:
            continue
        numeric_series = pd.to_numeric(df_price[col], errors='coerce')
        if numeric_series.dropna().empty:
            continue
        if (numeric_series == 0).any():
            continue
        candidate_columns.append(col)
        col_lower = str(col).strip().lower()
        if not normalized_top10 or col_lower in normalized_top10:
            filtered_columns.append(col)
    if normalized_top10 and not filtered_columns:
        return pd.DataFrame(columns=[date_col]), [], pd.DataFrame(columns=[date_col])
    if not filtered_columns:
        filtered_columns = candidate_columns[:10]
    if not filtered_columns:
        return pd.DataFrame(columns=[date_col]), [], pd.DataFrame(columns=[date_col])
    total_numeric = pd.to_numeric(df_price[total_col], errors='coerce')
    ratio_df = pd.DataFrame({date_col: df_price[date_col]})
    ratio_df['Total'] = np.where(total_numeric.notna(), 100.0, np.nan)
    original_df = pd.DataFrame({date_col: df_price[date_col], 'Total': total_numeric})
    for col in filtered_columns:
        numeric_series = pd.to_numeric(df_price[col], errors='coerce')
        original_df[col] = numeric_series
        ratio_df[col] = numeric_series
    total_series = total_numeric.astype(float)
    with np.errstate(divide='ignore', invalid='ignore'):
        for col in filtered_columns:
            ratio_df[col] = ratio_df[col].astype(float) / total_series * 100.0
    ratio_df = ratio_df.replace([np.inf, -np.inf], np.nan)
    ratio_df.dropna(subset=filtered_columns, how='all', inplace=True)
    return ratio_df, ['Total'] + filtered_columns, original_df
def graf_price_index_table(
    columns: list[str],
    averages: dict[str, float],
    color_mapping: dict[str, str],
    metric_label: str,
    c_fig: int
) -> tuple[io.BytesIO, tuple[float, float]]:
    """Genera la tabla de indicadores de precio indexado."""
    column_count = max(len(columns), 1)
    fig_width = max(6.0, 1.25 * column_count)
    fig_height = 1.5
    fig, ax = plt.subplots(num=c_fig, figsize=(fig_width, fig_height), dpi=DEFAULT_EXPORT_DPI)
    ax.axis('off')
    ax.set_facecolor('white')
    header_labels = ['Indicador'] + [wrap_table_text(str(col)) for col in columns]
    display_columns = ['Indicador'] + columns
    row_values = [metric_label]
    for col in columns:
        value = averages.get(col)
        if value is None or not np.isfinite(value):
            row_values.append('-')
        else:
            row_values.append(f"{value:.2f}%")
    row_values[0] = wrap_table_text(row_values[0])
    table = ax.table(
        cellText=[row_values],
        colLabels=header_labels,
        cellLoc='center',
        loc='center'
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1.05, 1.4)
    default_header_color = HEADER_COLOR_PRIMARY
    header_cell = table[(0, 0)]
    header_cell.set_facecolor(default_header_color)
    header_cell.set_text_props(color='white', weight='bold')
    header_cell.set_height(header_cell.get_height() * 1.5)
    for idx, col in enumerate(columns, start=1):
        header_cell = table[(0, idx)]
        header_color = color_mapping.get(col, '#777777')
        header_cell.set_facecolor(header_color)
        header_cell.set_text_props(color='white', weight='bold')
        header_cell.set_height(header_cell.get_height() * 1.5)
    for idx in range(len(display_columns)):
        data_cell = table[(1, idx)]
        data_cell.set_height(data_cell.get_height() * 1.35)
        data_cell.set_facecolor('#FFFFFF')
        data_cell.set_edgecolor('#DDDDDD')
    fig.tight_layout(pad=0.2)
    fig_size = fig.get_size_inches()
    return figure_to_stream(fig, dpi=TABLE_EXPORT_DPI), fig_size
def build_price_index_slide(
    ppt: Presentation,
    lang: str,
    cat: str,
    apo_entries: list[dict],
    removed_headers: list[str],
    price_df: pd.DataFrame,
    price_original_df: pd.DataFrame,
    chart_pipeline: int,
    chart_share_lookup: dict,
    c_fig: int
) -> int:
    """Construye la diapositiva de precio indexado."""
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    titulo_base = c_w.get((lang, '6-1'), 'Precio indexado')
    titulo_completo = f"{titulo_base} | {labels[(lang,'comp')]}{cat}"
    title_prefix = f"{titulo_base} | {labels[(lang,'comp')]}"
    title_box = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
    title_tf = title_box.text_frame
    competition_title_lookup: dict[str, str] = {}
    set_title_with_brand_color(
        title_tf,
        title_prefix,
        cat,
        '',
        0.33,
        competition_title_lookup,
        COMPETITION_TITLE_PALETTE
    )
    comment_box = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
    comment_tf = comment_box.text_frame
    comment_tf.clear()
    comment_paragraph = comment_tf.paragraphs[0]
    comment_paragraph.text = "Comentario" if lang == 'P' else "Comentario"
    comment_paragraph.font.size = Inches(0.25)
    chart_share_lookup = chart_share_lookup or {}
    if not price_df.empty and price_df.shape[1] > 1:
        left_margin_chart = Inches(0.33)
        right_margin_chart = Inches(0.33)
        available_line_width = available_width(ppt, left_margin_chart, right_margin_chart)
        c_fig += 1
        chart_height = Cm(12)
        chart_colors = {}
        slide.shapes.add_picture(
            line_graf(
                price_df,
                chart_pipeline,
                titulo_completo,
                c_fig,
                1,
                width_emu=available_line_width,
                height_emu=chart_height,
                share_lookup=chart_share_lookup,
                y_axis_percent=True,
                top_value_annotations=2,
                color_collector=chart_colors,
                color_overrides={'Total': '#000000'},
                linestyle_overrides={'Total': '--'},
                show_title=False,
            ),
            left_margin_chart,
            Inches(CHART_TOP_INCH),
            width=available_line_width,
            height=chart_height
        )
        plt.clf()
        chart_colors.setdefault('Total', '#000000')
        if not price_original_df.empty and chart_colors:
            data_cols = [col for col in price_df.columns[1:] if col in chart_colors]
            if 'Total' in data_cols:
                data_cols = ['Total'] + [col for col in data_cols if col != 'Total']
            if data_cols:
                tail_df = price_original_df.iloc[-12:] if len(price_original_df) >= 12 else price_original_df
                averages_raw: dict[str, float] = {}
                for col in data_cols:
                    series = pd.to_numeric(tail_df[col], errors='coerce')
                    avg_value = series.mean() if not series.dropna().empty else np.nan
                    averages_raw[col] = avg_value
                total_key = next((col for col in data_cols if str(col).strip().lower() == 'total'), None)
                total_avg = averages_raw.get(total_key) if total_key else np.nan
                averages = {}
                if total_key and np.isfinite(total_avg) and not np.isclose(total_avg, 0.0):
                    for col, avg_value in averages_raw.items():
                        if avg_value is None or not np.isfinite(avg_value):
                            averages[col] = np.nan
                        elif col == total_key:
                            averages[col] = 100.0
                        else:
                            averages[col] = (avg_value / total_avg) * 100.0
                else:
                    for col, avg_value in averages_raw.items():
                        averages[col] = avg_value if avg_value is not None and np.isfinite(avg_value) else np.nan
                if any(np.isfinite(val) for val in averages.values()):
                    c_fig += 1
                    metric_label = "Preco indexado\nmedio 12 meses" if lang == 'P' else "Precio indexado 12m"
                    table_stream, fig_size = graf_price_index_table(
                        data_cols,
                        averages,
                        chart_colors,
                        metric_label,
                        c_fig
                    )
                    target_table_height = Cm(TABLE_TARGET_HEIGHT_CM)
                    bottom_margin = Cm(1.5)
                    slide_height_emu = int(ppt.slide_height)
                    chart_bottom_emu = int(Inches(CHART_TOP_INCH) + chart_height)
                    min_table_top = chart_bottom_emu + int(Cm(0.2))
                    table_top_emu = slide_height_emu - int(target_table_height) - int(bottom_margin)
                    if table_top_emu < min_table_top:
                        table_top_emu = min_table_top
                    pic_table = slide.shapes.add_picture(
                        table_stream,
                        left_margin_chart,
                        table_top_emu,
                        height=target_table_height
                    )
                    constrain_picture_width(pic_table, available_line_width)
                    plt.clf()
    return c_fig
# Variable de control del numero de graficos
def _try_parse_reference_date(value):
    """Intenta parsear una fecha desde distintos formatos."""
    month_token_map = {
        'jan': 1, 'ene': 1,
        'feb': 2,
        'mar': 3,
        'apr': 4, 'abr': 4,
        'may': 5,
        'jun': 6,
        'jul': 7,
        'aug': 8, 'ago': 8,
        'sep': 9, 'set': 9,
        'oct': 10,
        'nov': 11,
        'dec': 12, 'dic': 12
    }
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return None
    if isinstance(value, str):
        value = value.strip()
        if not value:
            return None
        match = re.search(r'([A-Za-z]{3})[-/](\d{2,4})', value)
        if match:
            month_token, year_token = match.groups()
            month = month_token_map.get(month_token.lower())
            if month:
                year = int(year_token)
                if year < 100:
                    year += 2000 if year < 50 else 1900
                return dt(year, month, 1)
        # Intentar formatos cortos antes de delegar en pandas
        for fmt in ('%b-%y', '%b-%y  '):
            try:
                return dt.strptime(value, fmt)
            except ValueError:
                continue
    try:
        parsed = pd.to_datetime(value, errors='coerce')
    except Exception:
        return None
    if pd.isna(parsed):
        return None
    return parsed.to_pydatetime() if hasattr(parsed, 'to_pydatetime') else parsed

c_fig=0
plot=plot_ven()
last_reference_source = None
last_reference_origin = None
last_tree_period_dt = None
players_share_context = {}
#---------------------------------------------------------------------------------------------------------------------
chart_generation_start = dt.now()
for w in W:
    progress_message = build_terminal_progress_message(w, lang, cat)
    if progress_message:
        print_colored(progress_message, COLOR_QUESTION)
    distribution_match = DISTRIBUTION_SHEET_PATTERN.match(w.strip())
    # Segmento 7: distribuciones (regiones/NSE/otros cortes) desde hojas 7_*_*.
    if distribution_match:
        dist_kind_raw = distribution_match.group(2)
        dist_kind = str(dist_kind_raw).upper().strip()
        category_segment = distribution_match.group(1).replace('.', ' ').strip()
        if dist_kind == 'R':
            dist_label = 'Regiones'
        elif dist_kind == 'NSE':
            dist_label = 'NSE'
        else:
            pretty_dist = dist_kind.replace('_', ' ').replace('-', ' ')
            dist_label = ' '.join(word.capitalize() for word in pretty_dist.split()) or dist_kind
        dist_df, dist_series = parse_distribution_sheet(file, w, client)
        if dist_df.empty:
            print_colored(f"No se encontraron datos para graficar en la hoja {w}.", COLOR_RED)
            continue
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title_box = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        title_tf = title_box.text_frame
        suffix_text = ''
        set_title_with_brand_color(
            title_tf,
            f"Distribucion {dist_label} | ",
            category_segment,
            suffix_text,
            0.35,
            BRAND_TITLE_COLOR_LOOKUP
        )
        left_margin = Inches(0.33)
        right_margin = Inches(0.33)
        available_width_emu = available_width(ppt, left_margin, right_margin)
        available_width_in = emu_to_inches(int(available_width_emu))
        slide_height_in = emu_to_inches(int(ppt.slide_height))
        comment_top_in = emu_to_inches(int(Inches(6.33)))
        max_height_in = min(
            slide_height_in - CHART_TOP_INCH - 0.3,
            comment_top_in - CHART_TOP_INCH - 0.2
        )
        if max_height_in <= 0:
            max_height_in = None
        c_fig += 1
        chart_stream, fig_size = plot_distribution_chart(
            dist_df,
            c_fig,
            BRAND_TITLE_COLOR_LOOKUP,
            target_width_in=available_width_in,
            max_height_in=max_height_in
        )
        fig_width_in, fig_height_in = fig_size
        target_height_in = fig_height_in
        if fig_width_in and fig_width_in > 0 and available_width_in:
            target_height_in = fig_height_in * (available_width_in / fig_width_in)
        if max_height_in is not None and target_height_in > max_height_in:
            target_height_in = max_height_in
        slide.shapes.add_picture(
            chart_stream,
            left_margin,
            Inches(CHART_TOP_INCH),
            width=available_width_emu,
            height=Inches(target_height_in)
        )
        comment_box = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
        comment_tf = comment_box.text_frame
        comment_tf.clear()
        comment_para = comment_tf.paragraphs[0]
        comment_para.text = "Comentario"
        comment_para.font.size = Inches(0.25)
        plt.clf()
        continue
    sheet_clean = str(w).strip()
    # Segmento 8: intervalos y error muestral con bloques mensuales + agregados.
    if sheet_clean.startswith('8'):
        try:
            raw_df = file.parse(w, header=None)
        except Exception as exc:
            print_colored(f"No se pudo leer la hoja {w}: {exc}", COLOR_RED)
            continue
        monthly_blocks = ppt8_parse_monthly_blocks(raw_df)
        agg_blocks = ppt8_parse_agg_blocks(raw_df)
        if not monthly_blocks or not agg_blocks:
            print_colored(f"Segmento 8: no se detectaron bloques validos en la hoja {w}.", COLOR_RED)
            continue
        agg_blocks, agg_dropped = ppt8_filter_agg_blocks(monthly_blocks, agg_blocks)
        if agg_dropped:
            unique_labels: list[str] = []
            for label in agg_dropped:
                cleaned = str(label).strip()
                if cleaned and cleaned not in unique_labels:
                    unique_labels.append(cleaned)
            preview = ", ".join(unique_labels[:3])
            extra = len(unique_labels) - 3
            extra_note = f", +{extra} mas" if extra > 0 else ""
            detail = f": {preview}{extra_note}" if preview else ""
            print_colored(
                f"Segmento 8: se omitieron {len(agg_dropped)} bloque(s) agregados sin datos{detail}.",
                COLOR_YELLOW
            )
        if not agg_blocks:
            print_colored(f"Segmento 8: no se detectaron bloques agregados validos en la hoja {w}.", COLOR_RED)
            continue
        if len(monthly_blocks) != len(agg_blocks):
            print_colored(
                f"Segmento 8: se usara el minimo comun de bloques en {w} (mensuales={len(monthly_blocks)}, agregados={len(agg_blocks)}).",
                COLOR_YELLOW
            )
        mat_base_label, mat_curr_label = ppt8_find_mat_labels(raw_df)
        ppt_rows, ppt_errors = ppt8_compute_rows(monthly_blocks, agg_blocks)
        for err in ppt_errors:
            print_colored(err, COLOR_YELLOW)
        if not ppt_rows:
            print_colored(f"Segmento 8: no hay filas calculables en la hoja {w}.", COLOR_RED)
            continue
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        title_box = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        title_tf = title_box.text_frame
        title_prefix = c_w.get((lang, '8'), 'Intervalos') + ' | '
        brand_label = sheet_clean[2:].replace('_', ' ').strip() or cat
        set_title_with_brand_color(
            title_tf,
            title_prefix,
            brand_label,
            '',
            0.35,
            CATEGORY_TITLE_COLOR_LOOKUP,
            CATEGORY_TITLE_PALETTE
        )
        left_margin = Inches(0.33)
        right_margin = Inches(0.33)
        available_width_emu = available_width(ppt, left_margin, right_margin)
        available_width_in = emu_to_inches(int(available_width_emu))
        slide_height_in = emu_to_inches(int(ppt.slide_height))
        comment_top_in = emu_to_inches(int(Inches(6.33)))
        max_height_in = max(2.5, slide_height_in - CHART_TOP_INCH - 0.3)
        if comment_top_in > 0:
            max_height_cap = max(1.5, comment_top_in - CHART_TOP_INCH - 0.2)
            max_height_in = min(max_height_in, max_height_cap)
        c_fig += 1
        # Colores de marcas en tabla PPT8: iniciar desde el primer color de la paleta,
        # sin heredar asignaciones previas de otros segmentos/títulos.
        ppt8_brand_color_lookup: dict[str, str] = {}
        table_stream, fig_size = render_ppt8_table(ppt_rows, c_fig, ppt8_brand_color_lookup, mat_curr_label)
        fig_width_in, fig_height_in = fig_size
        target_height_in = fig_height_in
        if fig_width_in and fig_width_in > 0 and available_width_in and available_width_in > 0:
            target_height_in = fig_height_in * (available_width_in / fig_width_in)
        if max_height_in and target_height_in > max_height_in:
            target_height_in = max_height_in
        slide.shapes.add_picture(
            table_stream,
            left_margin,
            Inches(CHART_TOP_INCH),
            width=available_width_emu,
            height=Inches(target_height_in)
        )
        comment_box = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
        comment_tf = comment_box.text_frame
        comment_tf.clear()
        comment_para = comment_tf.paragraphs[0]
        comment_para.text = "Comentario" if lang == 'P' else "Comentario"
        comment_para.font.size = Inches(0.25)
        plt.clf()
        last_reference_source = raw_df
        last_reference_origin = 'Segmento 8'
        continue
    # Precio indexado: usa el contexto de Players (segmento 6) para construir la tabla.
    if is_price_index_sheet(w):
        players_base_key = extract_players_base_key(w)
        context = players_share_context.get(players_base_key)
        if context is None:
            print_colored(f"No se encontro contexto de Players para la hoja {w}. Se omite Precio indexado.", COLOR_YELLOW)
            continue
        df_start = parse_sheet_with_compras_header(file, w)
        df_list, p_ventas = split_compras_ventas(df_start, sheet_name=w)
        try:
            series_configs = prepare_series_configs(df_list, lang, p_ventas, sheet_name=w)
        except ValueError as exc:
            msg = str(exc)
            color = COLOR_YELLOW if "encabezado 'table' con periodos MAT" in msg else COLOR_RED
            print_colored(msg, color)
            continue
        price_series = None
        for cfg in series_configs:
            if cfg.display_tipo.lower() == 'compras':
                price_series = cfg
                break
        if price_series is None:
            price_series = series_configs[0] if series_configs else None
        if price_series is None:
            print_colored(f"No se detecto informacion de Compras en la hoja {w}.", COLOR_YELLOW)
            continue
        price_df, selected_columns, price_original_df = prepare_price_index_dataframe(price_series, context.get('compras_top10', []))
        share_lookup_context = context.get('share_lookup', {})
        chart_share_lookup = {}
        if share_lookup_context and selected_columns:
            share_lookup_lower = {
                str(key).strip().lower(): value
                for key, value in share_lookup_context.items()
            }
            for col in selected_columns:
                share_val = share_lookup_lower.get(str(col).strip().lower())
                if share_val is not None:
                    chart_share_lookup[col] = share_val
        c_fig = build_price_index_slide(
            ppt,
            lang,
            cat,
            context.get('apo_entries', []),
            context.get('removed_headers', []),
            price_df,
            price_original_df,
            price_series.pipeline,
            chart_share_lookup,
            c_fig
        )
        last_reference_source = price_series.data
        last_reference_origin = price_series.display_tipo
        continue
    # Segmento 2: arbol de medidas (por que) con comparacion MAT.
    if w.startswith('2_'):
        raw_tree_df = file.parse(w, header=None)
        sheet_name_clean = w.strip()
        tree_table, tree_periods = _limpiar_tabla_excel(raw_tree_df)
        if tree_table is None or not tree_periods:
            print_colored(f"No se encontro la tabla de variables en la hoja {w}.", COLOR_YELLOW)
            continue
        if len(tree_periods) < 2:
            print_colored(f"La hoja {w} no tiene al menos dos periodos MAT para comparar.", COLOR_YELLOW)
            continue
        unidad_detectada = _unidad_desde_nombre_hoja(w) or 'Units'
        periodo_inicial, periodo_final = tree_periods[-2], tree_periods[-1]
        parsed_tree_dt = _try_parse_reference_date(periodo_final)
        if parsed_tree_dt is not None:
            last_tree_period_dt = parsed_tree_dt
        try:
            metrics_calculated = calcular_cambios(tree_table, periodo_inicial, periodo_final, unidad_detectada)
        except KeyError as exc:
            print_colored(f"No se pudo calcular el arbol para {w}: {exc}", COLOR_RED)
            continue
        c_fig += 1
        tree_stream, fig_size = graficar_arbol(
            metrics_calculated,
            volumen_unidad=unidad_detectada,
            hoja=sheet_name_clean.replace(" ", "_"),
            output_dir=None
        )
        if fig_size:
            fig_width_in, fig_height_in = fig_size
        else:
            fig_width_in, fig_height_in = (21, 9)
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        titulo_arbol = c_w.get((lang, '2'), '2W - Arbol de Medidas')
        match_nombre = re.match(r'^\d+_(.+?)(?:_[A-Za-z])?$', sheet_name_clean)
        display_target = match_nombre.group(1) if match_nombre else sheet_name_clean[2:]
        title_box = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        title_tf = title_box.text_frame
        set_title_with_brand_color(
            title_tf,
            f"{titulo_arbol} | ",
            display_target,
            '',
            0.35,
            BRAND_TITLE_COLOR_LOOKUP
        )
        left_margin = Inches(0.33)
        right_margin = Inches(0.33)
        chart_top = Inches(CHART_TOP_INCH)
        available_width_emu = available_width(ppt, left_margin, right_margin)
        slide_width_in = emu_to_inches(ppt.slide_width)
        slide_height_in = emu_to_inches(ppt.slide_height)
        left_margin_in = emu_to_inches(left_margin)
        right_margin_in = emu_to_inches(right_margin)
        available_width_in = max(1.0, slide_width_in - left_margin_in - right_margin_in)
        aspect_ratio = fig_height_in / fig_width_in if fig_width_in else 0.52
        max_chart_height_in = max(3.0, slide_height_in - emu_to_inches(chart_top) - 0.4)
        target_width_in = available_width_in
        target_height_in = target_width_in * aspect_ratio if aspect_ratio else 5.5
        scale_factor = 1.12
        target_width_in *= scale_factor
        target_height_in *= scale_factor
        if target_width_in > available_width_in:
            target_width_in = available_width_in
            target_height_in = target_width_in * aspect_ratio if aspect_ratio else target_height_in
        if target_height_in > max_chart_height_in:
            target_height_in = max_chart_height_in
            target_width_in = target_height_in / aspect_ratio if aspect_ratio else target_width_in
        chart_shape = slide.shapes.add_picture(
            tree_stream,
            left_margin,
            chart_top,
            width=Inches(target_width_in),
            height=Inches(target_height_in)
        )
        # Caja de comentario en posicion fija aun si sobrepone el grafico
        comment_box = slide.shapes.add_textbox(Inches(0.33), Inches(5.8), Inches(10), Inches(0.5))
        comment_tf = comment_box.text_frame
        comment_tf.clear()
        comment_para = comment_tf.paragraphs[0]
        comment_para.text = "Comentario"
        comment_para.font.size = Inches(0.25)
        # No imprimimos mensaje de finalización para mantener la terminal concisa
        continue
    # Segmento 1: grafico de variaciones MAT (cuando).
    if w[0]=='1':
        try:
            sheet_df = parse_sheet_with_date_fallback(file, w)
        except ValueError as exc:
            print_colored(str(exc), COLOR_RED)
            continue
        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        title_prefix = c_w[(lang,w[0])]+' '+ labels[(lang,'MAT')]+' | '
        title_brand = w[2:].strip()
        set_title_with_brand_color(tf, title_prefix, title_brand, '', 0.35, BRAND_TITLE_COLOR_LOOKUP)
        #Obtém pipeline das vendas
        p = _detect_ventas_pipeline(sheet_df.columns)
        #Cria a base
        mat=df_mat(sheet_df,p)
        last_reference_source = mat
        last_reference_origin = None
        #Elimina linhas com divisão com zero devido ao pipeline
        mat=mat[~np.isinf(mat.iloc[:, 3])]
        #Incrementa contador do gráfico
        c_fig+=1
        #Insere o gráfico do MAT
        left_margin = Inches(0.33)
        right_margin = Inches(0.33)
        available = available_width(ppt, left_margin, right_margin)
        pic=slide.shapes.add_picture(graf_mat(mat,c_fig,p), left_margin, Inches(CHART_TOP_INCH),width=available)
        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(5.8), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = "Comentario"
        t.font.size = Inches(0.28)
        #Limpa área de plotagem
        plt.clf()
        #mensaje de conclusion por cada slide
    # Segmentos 3-6: comparativos y competencia con grafico + tabla de aportes.
    else: 
        #Carrega a base
        df_start=parse_sheet_with_compras_header(file, w)
        df_list, p_ventas = split_compras_ventas(df_start, sheet_name=w)
        try:
            series_configs = prepare_series_configs(df_list, lang, p_ventas, sheet_name=w)
        except ValueError as exc:
            print_colored(str(exc), COLOR_RED)
            continue
        last_reference_source = series_configs[0].data if series_configs else df_start
        last_reference_origin = series_configs[0].display_tipo if series_configs else None
        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        title_prefix = ''
        title_brand_label = None
        title_font_inches = 0.35
        title_color_lookup = BRAND_TITLE_COLOR_LOOKUP
        title_palette_sequence = None
        if w[0] in ['3','5']:
            title_prefix = c_w[(lang,w[0]+'-'+w[-1])]+' | '
            title_brand_label = w[2:-2].strip()
            titulo = f"{title_prefix}{title_brand_label}"
        elif w[0] == '4':
            title_prefix = c_w[(lang,w[0])] + ' | '
            title_brand_label = w[2:].strip()
            titulo = f"{title_prefix}{title_brand_label}"
        elif w[0]=='6':
            title_prefix =  c_w[(lang,w[0])] + ' | ' + labels[(lang,'comp')]
            title_brand_label = cat
            titulo = f"{title_prefix}{title_brand_label}"
            title_font_inches = 0.33
            title_color_lookup = {}
            title_palette_sequence = COMPETITION_TITLE_PALETTE
        else:
            title_prefix = c_w[(lang,w[0])]+' | '+ w[2:] 
            titulo = title_prefix
        set_title_with_brand_color(tf, title_prefix, title_brand_label, '', title_font_inches, title_color_lookup, title_palette_sequence)
        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text ="Comentario"
        t.font.size= Inches(0.25)
        # Control para la tabla de aporte
        target_table_height = Cm(TABLE_TARGET_HEIGHT_CM)
        bottom_margin = Cm(1.5)
        left_margin = int(Cm(TABLE_SIDE_MARGIN_CM))
        right_margin = left_margin
        vertical_spacing = Cm(0.2)
        target_table_height_emu = int(target_table_height)
        bottom_margin_emu = int(bottom_margin)
        slide_width = int(ppt.slide_width)
        slide_height = int(ppt.slide_height)
        removed_headers = []
        series_share_lookup = {}
        stacked_share_sources = []
        should_collect_share = False
        apo_entries = []
        chart_color_mappings = {}
        chart_top_emu = int(Inches(CHART_TOP_INCH))
        chart_height_emu = int(Cm(10))
        min_table_top = chart_top_emu + chart_height_emu + int(vertical_spacing)
        if w:
            first_char = w[0]
            last_char = w[-1]
            if first_char == '6':
                should_collect_share = True
            elif first_char == '3' and last_char in {'1', '2'}:
                should_collect_share = True
            elif first_char == '5' and last_char == '2':
                should_collect_share = True
        for idx, serie in enumerate(series_configs):
            apo = aporte(serie.data.copy(), serie.pipeline, lang, serie.raw_tipo)
            apo_entries.append({
                "apo": apo,
                "display_tipo": serie.display_tipo,
            })
            removed_headers.extend(apo.attrs.get("removed_columns", []))
            share_values = apo.attrs.get("share_mat_values", {})
            if share_values:
                series_share_lookup.update(share_values)
            if should_collect_share:
                share_periods = apo.attrs.get("share_period_values", [])
                if share_periods:
                    stacked_share_sources.append(
                        {
                            "display_tipo": serie.display_tipo,
                            "share_periods": share_periods,
                            "columns": [
                                str(col)
                                for col in apo.columns[1:]
                                if str(col).strip() and str(col).strip().lower() != 'total'
                            ],
                        }
                    )
        #Cuadro de texto que indica los encabezados eliminados
        unique_removed_headers = [header for header in dict.fromkeys(removed_headers) if header]
        if unique_removed_headers:
            footer_left = Cm(4.6)
            footer_right_margin = Cm(1.2)
            footer_width = ppt.slide_width - footer_left - footer_right_margin
            if footer_width < Cm(2):
                footer_width = Cm(2)
            footer_top = ppt.slide_height - bottom_margin + Cm(0.1)
            footer_height = Cm(0.9)
            footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
            footer_tf = footer_box.text_frame
            footer_tf.clear()
            footer_tf.word_wrap = True
            footer_tf.text = "Se eliminaron los encabezados sin información: " + ", ".join(unique_removed_headers)
            footer_run = footer_tf.paragraphs[0].runs[0]
            footer_font = footer_run.font
            footer_font.name = 'Arial'
            footer_font.size = Pt(10)
            footer_font.color.rgb = RGBColor(120, 120, 120)
        compras_top10 = []
        if series_share_lookup:
            compras_items = []
            for col_name, share_value in series_share_lookup.items():
                if col_name is None:
                    continue
                col_str = str(col_name).strip()
                if not col_str:
                    continue
                if col_str.lower().find('total') != -1:
                    continue
                if not col_str.lower().endswith('.c'):
                    continue
                try:
                    numeric_share = float(share_value)
                except (TypeError, ValueError):
                    continue
                if not np.isfinite(numeric_share):
                    continue
                compras_items.append((col_name, numeric_share))
            if compras_items:
                compras_items.sort(key=lambda item: item[1], reverse=True)
                compras_top10 = [item[0] for item in compras_items[:10]]
        players_base_key = extract_players_base_key(w)
        if players_base_key and apo_entries:
            players_share_context[players_base_key] = {
                "apo_entries": apo_entries,
                "removed_headers": unique_removed_headers,
                "share_lookup": dict(series_share_lookup),
                "titulo": titulo,
                "compras_top10": compras_top10,
            }
        if plot=="1" and len(series_configs)>1:
            df_full = series_configs[0].data.copy()
            for extra in series_configs[1:]:
                df_full = pd.concat([df_full, extra.data.iloc[:,1:]], axis=1)
            pipeline_combined = max((cfg.pipeline for cfg in series_configs), default=0)
            c_fig+=1
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available_line_width = available_width(ppt, left_margin, right_margin)
            chart_colors = {}
            chart_stream = line_graf(
                df_full,
                pipeline_combined,
                titulo,
                c_fig,
                len(series_configs),
                width_emu=available_line_width,
                height_emu=Cm(10),
                share_lookup=series_share_lookup,
                color_collector=chart_colors,
                color_overrides=chart_color_mappings,
                force_right_labels=True,
                show_title=False,
            )
            if isinstance(chart_colors, dict):
                chart_color_mappings.update(chart_colors)
            pic=slide.shapes.add_picture(chart_stream, left_margin, Inches(CHART_TOP_INCH),width=available_line_width,height=Cm(10))
        elif plot=="2" and len(series_configs)>1:
            ven_param = max(len(series_configs)-1, 1)
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available = available_width(ppt, left_margin, right_margin)
            count = len(series_configs)
            gap = Inches(0.1) if count > 1 else Inches(0)
            total_gap = int(gap) * (count - 1)
            effective_width = available - total_gap
            if count:
                if effective_width <= 0:
                    chart_width = available // count
                else:
                    chart_width = effective_width // count
            else:
                chart_width = available
            chart_width = max(chart_width, 1)
            for idx, serie in enumerate(series_configs):
                c_fig+=1
                left_position = left_margin + idx * (chart_width + int(gap))
                chart_colors = {}
                chart_title = format_title_with_bold_suffix(titulo, serie.display_tipo)
                chart_stream = line_graf(
                    serie.data,
                    serie.pipeline,
                    chart_title,
                    c_fig,
                    ven_param,
                    width_emu=chart_width,
                    height_emu=Cm(10),
                    multi_chart=True,
                    share_lookup=series_share_lookup,
                    color_collector=chart_colors,
                    color_overrides=chart_color_mappings,
                    force_right_labels=True,
                    include_origin_suffix=False
                )
                if isinstance(chart_colors, dict):
                    chart_color_mappings.update(chart_colors)
                pic=slide.shapes.add_picture(chart_stream, left_position, Inches(CHART_TOP_INCH),width=chart_width,height=Cm(10))
                plt.clf()
        elif series_configs:
            c_fig+=1
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available_single_width = available_width(ppt, left_margin, right_margin)
            chart_colors = {}
            chart_stream = line_graf(
                series_configs[0].data,
                series_configs[0].pipeline,
                titulo,
                c_fig,
                len(series_configs),
                width_emu=available_single_width,
                height_emu=Cm(10),
                share_lookup=series_share_lookup,
                color_collector=chart_colors,
                color_overrides=chart_color_mappings,
                show_title=False,
            )
            if isinstance(chart_colors, dict):
                chart_color_mappings.update(chart_colors)
            pic=slide.shapes.add_picture(chart_stream, left_margin, Inches(CHART_TOP_INCH),width=available_single_width,height=Cm(10))
            plt.clf()
        color_lookup_keys = build_color_lookup_dict(chart_color_mappings)
        table_entries = [
            entry for entry in apo_entries
            if not entry["apo"].attrs.get("skip_table")
            and not entry["apo"].empty
            and entry["apo"].shape[1] > 1
        ]
        if table_entries:
            def column_color_for_table(column_name):
                return lookup_color_for_label(column_name, color_lookup_keys)
            def build_table_color_mapping(apo_df):
                mapping = {}
                for column in apo_df.columns[1:]:
                    color_value = column_color_for_table(column)
                    if color_value:
                        mapping[column] = color_value
                return mapping
            table_top = slide_height - target_table_height_emu - bottom_margin_emu
            if table_top > min_table_top:
                table_top = min_table_top
            if len(table_entries) == 2:
                gap = Cm(TABLE_PAIR_GAP_CM)
                half_slide_width = slide_width // 2
                gap_half = int(gap) // 2
                paired_shapes = []
                for entry in table_entries:
                    apo_df = entry["apo"]
                    display_tipo = str(entry.get("display_tipo", "")).strip().lower()
                    if display_tipo == 'ventas':
                        left_position = half_slide_width + gap_half
                        max_width = slide_width - left_position - right_margin
                    else:
                        left_position = left_margin
                        max_width = half_slide_width - left_position - gap_half
                    c_fig += 1
                    table_colors = build_table_color_mapping(apo_df)
                    pic = slide.shapes.add_picture(
                        graf_apo(apo_df, c_fig, column_color_mapping=table_colors),
                        left_position,
                        table_top,
                        height=target_table_height
                    )
                    constrain_picture_width(pic, max_width)
                    paired_shapes.append((pic, max_width))
                    plt.clf()
                if paired_shapes:
                    min_height = min(shape.height for shape, _ in paired_shapes if shape is not None)
                    for shape, max_width in paired_shapes:
                        if shape is None or min_height <= 0:
                            continue
                        # Escalamos de forma proporcional para igualar alturas sin deformar
                        current_h = float(shape.height)
                        current_w = float(shape.width)
                        if current_h <= 0 or current_w <= 0:
                            continue
                        scale_h = min_height / current_h
                        # Evitar cualquier upscale involuntario
                        scale_h = min(scale_h, 1.0)
                        if max_width:
                            scale_w_limit = float(max_width) / current_w
                            if scale_w_limit <= 0:
                                scale_w_limit = scale_h
                            scale_h = min(scale_h, scale_w_limit)
                        new_h = int(current_h * scale_h)
                        new_w = int(current_w * scale_h)
                        shape.height = new_h
                        shape.width = new_w
            else:
                for entry in table_entries:
                    apo_df = entry["apo"]
                    c_fig += 1
                    table_colors = build_table_color_mapping(apo_df)
                    top_position = slide_height - target_table_height_emu - bottom_margin_emu
                    if top_position > min_table_top:
                        top_position = min_table_top
                    left_position = left_margin
                    max_width = slide_width - left_position - right_margin
                    pic = slide.shapes.add_picture(
                        graf_apo(apo_df, c_fig, column_color_mapping=table_colors),
                        left_position,
                        top_position,
                        height=target_table_height
                    )
                    constrain_picture_width(pic, max_width)
                    plt.clf()
        render_share_items = []
        if should_collect_share and stacked_share_sources:
            preferred_order = ['compras', 'ventas']
            ordered_sources = []
            consumed_indexes = set()
            for desired in preferred_order:
                for idx_source, entry in enumerate(stacked_share_sources):
                    if idx_source in consumed_indexes:
                        continue
                    display_tipo = (entry.get("display_tipo") or "").strip().lower()
                    if display_tipo == desired:
                        ordered_sources.append(entry)
                        consumed_indexes.add(idx_source)
                        break
            for idx_source, entry in enumerate(stacked_share_sources):
                if idx_source not in consumed_indexes:
                    ordered_sources.append(entry)
                    consumed_indexes.add(idx_source)
            def _sanitize_share_periods(entry_dict):
                valid_periods = []
                for period_label, shares in entry_dict.get("share_periods", []):
                    if not isinstance(shares, dict):
                        continue
                    sanitized = {}
                    has_positive = False
                    for column_name, value in shares.items():
                        if column_name is None:
                            continue
                        try:
                            numeric_value = float(value)
                        except (TypeError, ValueError):
                            continue
                        if not np.isfinite(numeric_value):
                            continue
                        if numeric_value < 0:
                            numeric_value = 0.0
                        sanitized[column_name] = numeric_value
                        if numeric_value > 0:
                            has_positive = True
                    if sanitized and has_positive:
                        valid_periods.append((period_label, sanitized))
                if len(valid_periods) > 2:
                    valid_periods = valid_periods[-2:]
                return valid_periods
            share_chart_entries = []
            for entry in ordered_sources:
                sanitized_periods = _sanitize_share_periods(entry)
                if sanitized_periods:
                    share_chart_entries.append(
                        {
                            "display_tipo": entry.get("display_tipo") or "",
                            "columns": entry.get("columns") or [],
                            "periods": sanitized_periods,
                        }
                    )
            if share_chart_entries:
                max_sections = 2
                share_chart_entries = share_chart_entries[:max_sections]
                palette_values = TREND_COLOR_SEQUENCE or [mcolors.to_hex(c) for c in plt.get_cmap('tab20').colors]
                share_color_lookup = dict(color_lookup_keys)
                used_color_values = {value for value in share_color_lookup.values() if value}
                palette_index = [0]
                palette_length = len(palette_values)
                def _next_palette_color():
                    if not palette_values:
                        return '#888888'
                    attempts = 0
                    while attempts < max(1, palette_length):
                        candidate = palette_values[palette_index[0] % palette_length]
                        palette_index[0] += 1
                        attempts += 1
                        if candidate not in used_color_values:
                            used_color_values.add(candidate)
                            return candidate
                    candidate = palette_values[palette_index[0] % palette_length]
                    palette_index[0] += 1
                    return candidate
                label_color_cache: dict[str, str] = {}
                def _color_for_share_label(label):
                    if label is None:
                        return _next_palette_color()
                    normalized_label = str(label).strip()
                    if not normalized_label or normalized_label.lower() == 'total':
                        return _next_palette_color()
                    cached_color = label_color_cache.get(normalized_label)
                    if cached_color:
                        return cached_color
                    color_value = lookup_color_for_label(normalized_label, share_color_lookup)
                    if color_value:
                        label_color_cache[normalized_label] = color_value
                        register_color_lookup(normalized_label, color_value, share_color_lookup, overwrite=False)
                        return color_value
                    color_value = _next_palette_color()
                    label_color_cache[normalized_label] = color_value
                    register_color_lookup(normalized_label, color_value, share_color_lookup, overwrite=False)
                    return color_value
                if len(share_chart_entries) == 1:
                    entry_info = share_chart_entries[0]
                    column_candidates = [col for col in entry_info["columns"] if col]
                    for _, shares in entry_info["periods"]:
                        for column_name in shares.keys():
                            if column_name and column_name not in column_candidates:
                                column_candidates.append(column_name)
                    color_mapping = {}
                    for column_name in column_candidates:
                        color_mapping[column_name] = _color_for_share_label(column_name)
                    for idx_period, (period_label, shares) in enumerate(entry_info["periods"]):
                        previous_shares = entry_info["periods"][idx_period - 1][1] if idx_period > 0 else None
                        render_share_items.append(
                            {
                                "display_tipo": entry_info["display_tipo"],
                                "period_label": period_label,
                                "shares": shares,
                                "previous_shares": previous_shares,
                                "color_mapping": color_mapping,
                                "show_delta": previous_shares is not None,
                            }
                        )
                else:
                    for entry_info in share_chart_entries:
                        column_candidates = [col for col in entry_info["columns"] if col]
                        for _, shares in entry_info["periods"]:
                            for column_name in shares.keys():
                                if column_name and column_name not in column_candidates:
                                    column_candidates.append(column_name)
                        color_mapping = {}
                        for column_name in column_candidates:
                            color_mapping[column_name] = _color_for_share_label(column_name)
                        latest_label, latest_shares = entry_info["periods"][-1]
                        previous_shares = entry_info["periods"][-2][1] if len(entry_info["periods"]) > 1 else None
                        render_share_items.append(
                            {
                                "display_tipo": entry_info["display_tipo"],
                                "period_label": latest_label,
                                "shares": latest_shares,
                                "previous_shares": previous_shares,
                                "color_mapping": color_mapping,
                                "show_delta": False,
                            }
                        )
        if render_share_items:
            delta_header = 'Crescimento vs ano anterior' if lang == 'P' else 'Crecimiento vs año anterior'
            share_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            share_title_box = share_slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
            share_tf = share_title_box.text_frame
            share_suffix = " - Share 100% Apilado"
            set_title_with_brand_color(
                share_tf,
                title_prefix,
                title_brand_label,
                share_suffix,
                0.33,
                title_color_lookup,
                title_palette_sequence
            )
            comment_box = share_slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
            comment_tf = comment_box.text_frame
            comment_tf.clear()
            comment_paragraph = comment_tf.paragraphs[0]
            comment_paragraph.text = "Comentario"
            comment_paragraph.font.size = Inches(0.25)
            chart_top = Inches(0.55)
            multiple_sections = len(render_share_items) > 1
            left_start = Inches(1.25) if multiple_sections else Inches(2.0)
            right_margin = Inches(0.6)
            horizontal_gap = Inches(0.4) if multiple_sections else Inches(0.5)
            caption_gap = Inches(0.12)
            caption_height = Inches(0.35)
            current_left = left_start
            reference_label = 'Corte em' if lang == 'P' else 'Corte a'
            for idx_chart, item in enumerate(render_share_items):
                shares = item["shares"]
                total_share = sum(max(value, 0) for value in shares.values())
                if total_share <= 0:
                    continue
                c_fig += 1
                chart_stream, fig_size = stacked_share_chart(
                    item["period_label"],
                    shares,
                    item["color_mapping"],
                    c_fig,
                    title=f"{item['display_tipo']} - {item['period_label']}"
                )
                fig_width_in, fig_height_in = fig_size
                if fig_height_in <= 0:
                    continue
                aspect_ratio = fig_width_in / fig_height_in if fig_height_in else 1.0
                target_height_in = max(6.0, fig_height_in * 1.05)
                slide_height_in = emu_to_inches(ppt.slide_height)
                chart_top_in = emu_to_inches(chart_top)
                caption_gap_in = emu_to_inches(caption_gap)
                caption_height_in = emu_to_inches(caption_height)
                available_height_in = max(3.0, slide_height_in - chart_top_in - caption_gap_in - caption_height_in - 0.2)
                target_height_in = min(available_height_in, target_height_in)
                target_width_in = max(3.8, target_height_in * aspect_ratio)
                slide_width_in = emu_to_inches(ppt.slide_width)
                available_width_in = slide_width_in - emu_to_inches(current_left) - emu_to_inches(right_margin)
                if available_width_in <= 0:
                    continue
                if target_width_in > available_width_in:
                    target_width_in = available_width_in
                    target_height_in = target_width_in / aspect_ratio if aspect_ratio else target_height_in
                target_height = Inches(target_height_in)
                target_width = Inches(target_width_in)
                chart_shape = share_slide.shapes.add_picture(
                    chart_stream,
                    current_left,
                    chart_top,
                    width=target_width,
                    height=target_height
                )
                caption_left = chart_shape.left
                caption_width = target_width
                caption_top = chart_top + target_height + caption_gap
                caption_box = share_slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
                caption_tf = caption_box.text_frame
                caption_tf.clear()
                caption_paragraph = caption_tf.paragraphs[0]
                origin_label = format_origin_label(item.get("display_tipo"), lang)
                origin_suffix = f" - {origin_label}" if origin_label else ""
                caption_paragraph.text = f"{reference_label}: {item['period_label']}{origin_suffix}"
                caption_paragraph.font.size = Pt(11)
                caption_paragraph.font.bold = True
                caption_paragraph.font.color.rgb = RGBColor(80, 80, 80)
                should_show_delta = item["show_delta"] and item["previous_shares"] is not None and idx_chart == len(render_share_items) - 1
                if should_show_delta:
                    delta_box_width = Inches(2.6)
                    delta_box_left = chart_shape.left + chart_shape.width + Inches(0.2)
                    slide_right_limit = ppt.slide_width - Inches(0.2)
                    if delta_box_left + delta_box_width > slide_right_limit:
                        delta_box_left = slide_right_limit - delta_box_width
                    delta_box = share_slide.shapes.add_textbox(
                        delta_box_left,
                        chart_top,
                        delta_box_width,
                        chart_shape.height + caption_gap + caption_height
                    )
                    delta_tf = delta_box.text_frame
                    delta_tf.clear()
                    delta_tf.word_wrap = True
                    header_para = delta_tf.paragraphs[0]
                    header_para.text = delta_header
                    header_para.font.bold = True
                    header_para.font.size = Pt(12)
                    header_para.font.color.rgb = RGBColor(70, 70, 70)
                    growth_entries = []
                    previous_shares = item["previous_shares"] or {}
                    for brand, current_value in shares.items():
                        prev_value = previous_shares.get(brand, 0.0)
                        if current_value is None or not np.isfinite(current_value):
                            current_value = 0.0
                        if prev_value is None or not np.isfinite(prev_value):
                            prev_value = 0.0
                        delta_value = (current_value - prev_value) * 100
                        growth_entries.append((brand, delta_value, current_value))
                    growth_entries.sort(key=lambda val: val[1], reverse=True)
                    for brand, delta_value, _ in growth_entries:
                        para = delta_tf.add_paragraph()
                        para.level = 1
                        para.text = f"{brand}: {delta_value:+.1f} pp"
                        para.font.size = Pt(11)
                        color_hex = item["color_mapping"].get(brand)
                        if color_hex:
                            rgb_tuple = mcolors.to_rgb(color_hex)
                            color_rgb = tuple(int(round(c * 255)) for c in rgb_tuple)
                            para.font.color.rgb = RGBColor(*color_rgb)
                        else:
                            para.font.color.rgb = RGBColor(90, 90, 90)
                current_left += target_width + horizontal_gap
                plt.clf()
        # Mensaje final omitido para evitar ruido en consola
chart_generation_end = dt.now()
# Referencia de la base
ref_source = last_reference_source if last_reference_source is not None else parse_sheet_with_compras_header(file, W[0])
last_dt = None
if ref_source is not None and not ref_source.empty:
    for val in reversed(ref_source.iloc[:, 0].tolist()):
        last_dt = _try_parse_reference_date(val)
        if last_dt is not None:
            break
if last_dt is None and last_tree_period_dt is not None:
    last_dt = last_tree_period_dt

if last_dt is None:
    ref = 'NA'
    ref_display = 'NA'
else:
    ref = last_dt.strftime('%m-%y')
    month_list = MONTH_NAMES.get(lang, MONTH_NAMES['E'])
    month_name = month_list[last_dt.month - 1]
    ref_display = f"{month_name}-{last_dt.year}"
slide=ppt.slides[0]
#Fabricante
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(3.77), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
t.text = "Numerator vs " + client
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.44), True
font.color.rgb = RGBColor(255, 255, 255)
#Referencia
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(5.98), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
reference_label = 'Corte em' if lang == 'P' else 'Corte a'
display_ref = ref_display if ref_display != 'NA' else ref
origin_label = format_origin_label(last_reference_origin, lang)
reference_text = f"{reference_label} {display_ref}"
if origin_label:
    reference_text = f"{reference_text} - {origin_label}"
t.text = reference_text
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.32), True
font.color.rgb = RGBColor(255, 255, 255)
output_filename = "-".join([
    simplify_name_segment(land, 30),
    simplify_name_segment(client, 30),
    simplify_name_segment(cat, 6),
    simplify_name_segment(brand, 8),
    '5W1H',
    simplify_name_segment(ref, 5)
]) + '.pptx'
ppt.save(output_filename)
chart_elapsed = int((chart_generation_end - chart_generation_start).total_seconds())
print_colored(
    f'Tiempo de generacion de graficos : {chart_elapsed//60} min {chart_elapsed%60} s'
    if chart_elapsed >= 60
    else f'Tiempo de generacion de graficos : {chart_elapsed} s',
    COLOR_BLUE
)
