#Bibliotecas necessarias
#---------------------------------------------------------------------------------------------------------------------
import pandas as pd
import numpy as np
import warnings
from datetime import datetime as dt
import os
import io
import tkinter as tk
from matplotlib import pyplot as plt
from datetime import datetime as dt
from matplotlib.ticker import FuncFormatter
from decimal import Decimal, ROUND_CEILING, ROUND_DOWN, ROUND_FLOOR, ROUND_HALF_DOWN, ROUND_HALF_EVEN, ROUND_HALF_UP, ROUND_UP, ROUND_05UP
from pptx import Presentation 
from pptx.util import Inches
from pptx.dml.color import RGBColor
from tkinter import messagebox

pd.set_option('future.no_silent_downcasting', True)
pd.set_option('mode.chained_assignment', None)
warnings.filterwarnings('ignore')

agora = dt.now()

#Funcao que prepara os dados para criação do gráfico MAT
def df_mat(df,p):
    
    v1 = pd.DataFrame([(df.iloc[i-12-p:i-p,1].sum()/df.iloc[i-24-p:i-12-p,1].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-in'])
    v2 = pd.DataFrame([(df.iloc[i-12:i,2].sum()/df.iloc[i-24:i-12,2].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-out'])

    ac1, ac2 = [pd.DataFrame(df[col].rolling(window=12).sum().dropna()) for col in df.columns[1:]]
    v1.index,v2.index=v1.index+11,v2.index+11
    v1 ,v2 = v1.dropna(), v2.dropna()

    mat =pd.DataFrame(df.iloc[:,0].copy())
    mat= mat.join(ac1.join(ac2)[-len(v1):].join([v1]).join([v2]))
    mat= mat[mat.notna().all(axis=1).idxmax():]

    return mat

#Funcao que gera o gráfico de linhas e barras para o MAT
def graf_mat (mat,c_fig,p):

    #Cria figura e área plotável   
        fig = plt.figure(c_fig, (15, 5))
        ax = fig.add_subplot(1, 1, 1)
        
        #Cria o eixo do tempo e os acumulados juntamente com as variações
        ran=mat.iloc[:,0].copy()
        ran=[x.strftime('%m-%y') for x in ran]
        ac1 = mat.iloc[:,1].copy()
        ac2 = mat.iloc[:,2].copy()
        v1= mat.iloc[:,3].copy()
        v2= mat.iloc[:,4].copy()

        #Plota os acumulados em linhas    
        l1=ax.plot(ran,ac1, color='#7F7F7F', linewidth=4.2, label='Acumulado Cliente')
        l2=ax.plot(ran,ac2, color='#000000', linewidth=4.2, label='Acumulado Numerator') 

        #Plota as variações em barras
        ax2= ax.twinx()
        b1 = ax2.bar(np.arange(len(ran)) - 0.3, v1.values, 0.3, color='#FFFFFF', edgecolor='black', label='Var % Cliente Pipeline: '+ str(p) +' '+labels[(lang,'Var MAT')])
        b2 = ax2.bar(np.arange(len(ran)) + 0.3, v2.values, 0.3, color='#000000', edgecolor='black', label='Var % Numerator '+labels[(lang,'Var MAT')])

        ax.set_xticks(np.arange(len(ran)))
        ax.set_xticklabels(ran, rotation=30)
        ax2.tick_params(left=False, labelleft=False, top=False, labeltop=False,
            right=False, labelright=False, bottom=False, labelbottom=False)

        for v in [v1,v2]:
            for x,y in zip(np.arange(len(v1))+0.2,v):
                label = "{:.1f}%".format(y)
                bbox_props_white=dict(facecolor='#F2F2F2',edgecolor='black')
                plt.annotate(f"{y*100:.1f}%", (x, y), textcoords="offset points", xytext=(0, 10), ha='center', color='red' if y < 0 else 'green', size=9, bbox=bbox_props_white)

        
        #Cria a legenda 
        lns = l1 + l2 + [b1, b2]
        labs = [l.get_label() for l in l1 + l2] + [b1.get_label(), b2.get_label()]
        ax.legend(lns, labs, loc='upper center', bbox_to_anchor=(0.475,-0.115), borderaxespad=0.05, frameon=True, prop={'size': 12}, ncol=2)
        img_stream = io.BytesIO()


        #Cria o titulo
        ax.set_title(labels[(lang,'MAT')]+' | ' + w[2:], size=18, pad=10)

        #Salva o grafico
        fig.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.01, transparent=True)

        #Insere o grafico
        img_stream.seek(0) 

        return img_stream

#Função que gera o gráfico de linhas
def line_graf(df, p, title, c_fig,ven):
    # Cria figura
    fig = plt.figure(c_fig, (15, 5))
    ax = fig.add_subplot(1, 1, 1)
    
    aux = df.copy()
    ran = [x.strftime('%m-%y') for x in aux.iloc[:, 0]]

    if ven>1:
        colunas = aux.columns[1:]  
        comp=colunas[:len(colunas)//2]
        vent=colunas[len(colunas)//2:]
        cmap = plt.get_cmap('tab20')
        cor_map = {base: cmap(i % 20) for i, base in enumerate(comp)}
        df_cor = pd.DataFrame({
        'Compras': comp,
        'Ventas': vent,
        'Cor': [cor_map[base] for base in comp]
        })
    else:
        colunas = aux.columns[1:]  
        comp=colunas
        cmap = plt.get_cmap('tab20')
        cor_map = {base: cmap(i % 20) for i, base in enumerate(comp)}
        df_cor = pd.DataFrame({
        'Compras': comp,
        'Ventas': comp,
        'Cor': [cor_map[base] for base in comp]
        })
    lns = []

    for col in colunas:
        if ven>1: 
            estilo = '-' if '.v' in col.lower() else '--' 
        else: 
            estilo='-'

        tipo = 'Ventas' if '.v' in col.lower() else 'Compras'
        cor = df_cor['Cor'][df_cor[tipo]==col].iloc[0]

        y = aux[col].values
        line, = ax.plot(ran[p:], y[p:], color=cor, linewidth=2, linestyle=estilo, label=col)
        lns.append(line)

    # Eixo X,Y
    plt.xticks(rotation=30)
    plt.ylim(0)

    # Legenda
    labs = [l.get_label() for l in lns]
    ax.legend(lns, labs, loc='upper center', bbox_to_anchor=(0.5, -0.15), borderaxespad=0, frameon=False,
              prop={'size': 10}, ncol=max(1, len(colunas)/2))
    plt.tight_layout()

    # Título
    plt.title(title, size=18, pad=10)

    # Salva imagem
    img_stream = io.BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight', pad_inches=0.01, transparent=True)
    img_stream.seek(0)

    return img_stream

#Função que cria a tabela de aporte
def aporte(df,p,lang,tipo):

        aux = df.copy()

        aux['Total'] = aux.iloc[:len(df),1:].sum(axis=1)

        apo=pd.DataFrame(columns=[tipo] + aux.columns[1:].tolist())
        #Vol ultimo MAT
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-12-p,labels[(lang,'Data')]].strftime('%b-%y') ] + [aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p,aux.shape[1]-1].sum() for col in range(1,len(aux.columns) )]
        #Vol MAT atual
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-p,labels[(lang,'Data')]].strftime('%b-%y') ] + [aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-12-p:len(aux)-p,aux.shape[1]-1].values.sum() for col in range(1,len(aux.columns))]
        #Var
        apo.loc[len(apo)] = ["Var %"] +  [ aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum()-1 for col in range(1,len(aux.columns))]
        #Aporte
        apo.loc[len(apo)] = ["Aporte"] + [apo.iloc[-1, col] * apo.iloc[0, col] for col in range(1,len(aux.columns))]

        #Formatação do volume
        apo.iloc[:2, 1:] = apo.iloc[:2, 1:].applymap(lambda x: f"{round(x * 100, 1)}%")
        #Formatação da variação e aporte
        apo.iloc[2:, 1:] = apo.iloc[2:, 1:].applymap(lambda x: f"{round(x * 100, 2)}%")

        return apo

#Função que cria o gráfico tabela de aporte
def graf_apo(apo,c_fig):

        fig, ax = plt.subplots(num=c_fig)
        row_height = 0.4  # polegadas por linha, ajuste conforme gosto
        col_width = 1.0   # polegadas por coluna, ajuste conforme gosto

        n_rows, n_cols = apo.shape
        fig_width = col_width * n_cols
        fig_height = row_height * (n_rows-1)  # +1 para cabeçalho

        fig.set_size_inches(fig_width, fig_height)
        ax.axis('off')

        table = ax.table(
        cellText=apo.values,
        colLabels=apo.columns,
        loc='center',
        cellLoc='center'
        )

        for i,col in enumerate(apo.columns):
            table.auto_set_column_width(i)

        n = apo.shape[1]

        for (row, col), cell in table.get_celld().items():
            cell.set_edgecolor('black')
            cell.set_linewidth(1)

            if row == 0:
                if col == 0 or col == n - 1:
                    cell.set_facecolor('navy')
                    cell.get_text().set_color('white')
                    cell.get_text().set_weight('bold')
                else:
                    cell.set_facecolor('gray')
                    cell.get_text().set_color('white')
                    cell.get_text().set_weight('bold')

            elif col >= 1 and  'Vol' not in str(apo.iloc[row-1,0]):
                val = apo.iloc[row-1, col]
                try:
                    num = float(str(val).replace('%', '').replace(',', '.'))
                    if num >= 0:
                        cell.get_text().set_color('green')
                    elif num < 0:
                        cell.get_text().set_color('red')
                except:
                    pass

        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=400, bbox_inches='tight')
        buf.seek(0)   
        return buf

def plot_ven():

    def selecionar(valor):
        plot[0] = valor
        messagebox.showinfo("Seleção", f"Tu elegiste: {valor}")
        root.destroy()
        
    root =tk.Tk()
    root.title("Elija una opción")

    plot = [None]
    opcoes = [
    ("1 - Plotear Ventas y Compras Juntas", "1"),
    ("2 - Plotear Ventas y Compras Separadas", "2"),
    ("3 - No hay W con Ventas en esa base", "3"),
            ]

    tk.Label(root, text="Si hay ventas en la base, elija cómo desea graficarlas:", font=("Arial", 12)).pack(anchor="w", pady=(10, 5), padx=10)

    for texto, valor in opcoes:
        tk.Button(root, text=texto, command=lambda v=valor: selecionar(v)).pack(anchor="w",pady=5,padx=10)

    root.mainloop()

    return plot[0]

#Codigos paises
pais = pd.DataFrame(
    {'cod': [10, 54, 91, 55, 12, 56, 57, 93, 52, 51, 66, 63, 62, 64, 65, 67, 69],
     'pais': ['LatAm', 'Argentina', 'Bolivia', 'Brasil', 'CAM', 'Chile', 'Colombia', 'Ecuador', 'Mexico', 'Peru', 
              'Costa Rica', 'El Salvador', 'Guatemala', 'Honduras', 'Nicaragua', 'Panama', 'Republica Dominicana']
    }
)

#Codigos categorias
categ = pd.DataFrame({
    'cest': ['Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 
               'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Lacteos', 
               'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 
               'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Ropas y Calzados', 'Ropas y Calzados', 'Ropas y Calzados', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'OTC', 'OTC', 'OTC', 'OTC', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 
               'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 
               'Diversos', 'Diversos', 'Diversos', 'Diversos','Alimentos'],
    'cat': ['Bebidas Alcohólicas', 'Cervezas', 'Bebidas Gaseosas', 'Agua Gasificada', 'Água de Coco', 'Café-Consolidado de Café', 'Cross Category (Bebidas)', 'Bebidas Energéticas', 
                  'Bebidas Saborizadas Sin Gas', 'Café Tostado y Molido', 'Jugos Caseros', 'Té Helado', 'Café Instantáneo-Café Sucedáneo', 'Jugos y Nectares', 'Zumos de Vegetales', 'Agua Natural', 
                  'Gaseosas + Aguas', 'Mixta Café+Malta', 'Mixta Dolce Gusto-Mixta Té Helado + Café + Modificadores', 'Mixta Jugos y Leches', 'Mixta Jugos Líquidos + Bebidas de Soja', 'Mixta Té+Café', 
                  'Jugos Liquidos-Jugos Polvo', 'Refrescos en Polvo-Jugos - Bebidas Instantáneas En Polvo - Jugos Polvo', 'Bebidas Refrescantes', 'Refrescos Líquidos-Jugos Líquidos', 'Té Líquido - Listo para Tomar',
                  'Bebidas de Soja', 'Bebidas Isotónicas', 'Té e Infusiones-Te-Infusión Hierbas', 'Yerba Mate', 'Manteca', 'Queso Fresco y para Untar', 'Leche Condensada', 'Queso Untable', 'Yoghurt p-beber', 
                  'Leche Culinaria-Leche Evaporada', 'Leche Fermentada', 'Leche Líquida Saborizada-Leche Líquida Con Sabor', 'Fórmulas Infantiles', 'Leche Líquida', 'Leche Larga Vida', 'Margarina', 'Queso Fundido', 
                  'Crema de Leche', 'Mixta Lácteos-Postre+Leches+Yogurt', 'Mixta Leches', 'Mixta Yoghurt+Postres', 'Petit Suisse', 'Leche en Polvo', 'Yoghurt p-comer', 'Leche-Leche Líquida Blanca - Leche Liq. Natural', 
                  'Yoghurt', 'Ropas', 'Calzados', 'Medias-Calcetines', 'Arepas', 'Cereales Infantiles', 'Nutrición Infantil-Colados y Picados', 'Frijoles', 'Galletas', 'Caldos-Caldos y Sazonadores', 'Pan', 
                  'Apanados-Empanizadores', 'Empanados', 'Cereales-Cereales Desayuno-Avenas y Cereales', 'Hamburguesas', 'Mezclas Listas para Tortas-Preparados Base Harina Trigo', 'Queques-Ponques Industrializados', 
                  'Conservas De Pescado', 'Conservas de Frutas y Verduras', 'Dulce de Leche-Manjar', 'Alfajores', 'Barras de Cereal', 'Pollo', 'Chocolate', 'Chocolate de Taza-Achocolatados - Cocoas', 'Salsas Frías', 
                  'Compotas', 'Condimentos y Especias', 'Chocolate de Mesa', 'Aceite-Aceites Comestibles', 'Salsas Listas-Salsas Caseras Envasadas', 'Grano, Harina y Masa de Maíz', 'Fécula de Maíz', 'Harina De Maíz', 
                  'Ayudantes Culinarios', 'Postres Preparados', 'Jamón Endiablado', 'Semillas y Frutos Secos', 'Pan de Pascua', 'Huevos de Páscua', 'Huevos', 'Flash Cecinas', 'Harinas', 'Carne Fresca', 
                  'Platos Listos Congelados','Alimentos Congelados', 'Jamones', 'Cereales Calientes-Cereales Precocidos', 'Salsas Picantes', 'Helados', 'Pan Industrializado', 'Puré Instantáneo', 
                  'Fideos Instantáneos', 'Mermeladas', 'Ketchup', 'Jugo de Limon Adereso', 'Maltas', 'Adobos - Sazonadores', 'Mayonesa', 'Cárnicos', 'Modificadores de Leche-Saborizadores p-leche', 
                  'Mixta Cereales Infantiles+Avenas', 'Mixta Caldos + Saborizantes', 'Mixta Caldos + Sopas', 'Mixta Cereales + Cereales Calientes', 'Mixta Chocolate + Manjar', 'Galletas, snacks y mini tostadas', 
                  'Aceites + Mantecas', 'Aceites + Conservas De Pescado', 'Ayudantes Culinarios + Bolsa de Hornear', 'Mixta Huevos de Páscua + Chocolates', 'Mixta Platos Listos Congelados + Pasta', 
                  'Mixta Platos Congelados y Listos para Comer', 'Mixta Alimentos Congelados + Margarina', 'Mixta Modificadores + Cocoa', 'Mixta Pastas', 'Mixta Sopas+Cremas+Ramen', 
                  'Mixta Margarina + Mayonesa + Queso Crema', 'Mixta Azúcar+Endulzantes', 'Mostaza', 'Sustitutos de Crema', 'Fideos', 'Nuggets', 'Avena en hojuelas-liquidas', 'Aceitunas', 'Tortilla', 
                  'Panetón', 'Pastas', 'Salsas para Pasta', 'Turrón de maní', 'Carne Porcina', 'Postres en Polvo-Postres para Preparar - Horneables-Gelificables', 'Leche de Soya en Polvo', 'Cereales Precocidos', 
                  'Masas Frescas-Tapas Empanadas y Tarta', 'Pre-Pizzas', 'Meriendas listas', 'Arroz', 'Galletas de Arroz', 'Frijoles Procesados', 'Pratos Prontos - Comidas Listas', 'Aderezos para Ensalada', 
                  'Sal', 'Galletas Saladas-Galletas No Dulce', 'Sardina Envasada', 'Cecinas', 'Milanesas', 'Snacks', 'Fideos Sopa', 'Sopas-Sopas Cremas', 'Siyau', 'Tallarines-Spaguetti', 'Chocolate para Untar', 
                  'Azucar', 'Galletas Dulces', 'Untables Dulces', 'Endulzantes', 'Torradas - Tostadas', 'Salsas de Tomate', 'Atún Envasado', 'Leche Vegetal', 'Harinas de trigo', 
                  'Ambientadores-Desodorante Ambiental', 'Jabón en Barra-Jabón de lavar', 'Cloro-Lavandinas-Lejías-Blanqueadores', 'Pastillas para Inodoro', 'Guantes de látex', 
                  'Esponjas de Limpieza-Esponjas y paños', 'Utensilios de Limpieza', 'Filtros de Café', 'Cross Category (Limpiadores Domesticos)', 'Cross Category (Lavandería)', 
                  'Cross Category (Productos de Papel)', 'Lavavajillas-Lavaplatos - Lavalozas mano', 'Empaques domésticos-Bolsas plásticas-Plástico Adherente-Papel encerado-Papel aluminio', 
                  'Destapacañerias', 'Perfumantes para Ropa-Perfumes para Ropa', 'Cera p-pisos', 'Desodorante para Pies', 'Lustramuebles', 'Bolsas de Basura', 'Limpiadores verdes', 
                  'Limpiadores-Limpiadores y Desinfectantes', 'Insecticidas-Raticidas', 'Toallas de papel-Papel Toalla - Toallas de Cocina - Rollos Absorbentes de Papel', 'Detergentes para ropa', 
                  'Apresto', 'Mixta Pastillas para Inodoro + Limpiadores', 'Mixta Home Care-Cloro-Limpiadores-Ceras-Ambientadores', 'Mixta Limpiadores + Cloro', 'Mixta Detergentes + Cloro', 
                  'Mixta Detergentes + Lavavajillas', 'Pañitos + Papel Higienico', 'Servilletas', 'Film plastico e papel aluminio', 'Esponjas de Acero', 'Suavizantes de Ropa', 'Quitamanchas-Desmanchadores', 
                  'Papel Higiénico', 'Paños de Limpieza', 'Analgésicos-Painkillers', 'Suplementos alimentares', 'Gastrointestinales-Efervescentes', 'Vitaminas y Calcio', 'Categoría Desconocida', 
                  'Pilas-Baterías', 'Combustible Gas', 'Panel Financiero de Hogares', 'Panel Financiero de Hogares', 'Cartuchos de Tintas', 'Alimento para Mascota-Alim.p - perro - gato', 
                  'Telecomunicaciones - Convergencia', 'Tickets - Till Rolls', 'Tabaco - Cigarrillos', 'Incontinencia de Adultos', 'Shampoo Infantil', 'Maquinas de Afeitar', 'Cremas Corporales', 
                  'Paños Húmedos', 'Cremas para Peinar', 'Acondicionador-Bálsamo', 'Cross Category (Higiene)', 'Cross Category (Personal Care)', 'Desodorantes', 'Pañales-Pañales Desechables', 
                  'Cremas Faciales', 'Pañuelos Faciales', 'Protección Femenina-Toallas Femeninas', 'Fragancias', 'Cuidado del Cabello-Hair Care', 
                  'Tintes para el Cabello-Tintes - Tintura - Tintes y Coloración para el cabello', 'Depilación', 'Alisadores para el Cabello', 
                  'Fijadores para el Cabello-Modeladores-Gel-Fijadores para el cabello', 'Tratamientos para el Cabello', 'Óleo Calcáreo', 'Maquillaje-Cosméticos', 'Jabón Medicinal', 
                  'Pañitos + Pañales', 'Mixta Make Up+Tinturas', 'Enjuague Bucal-Refrescante Bucal', 'Cuidado Bucal', 'Protectores Femeninos', 'Toallas Femininas', 'Shampoo', 
                  'Afeitado-Crema afeitar-Loción de afeitar-Pord. Antes del afeitado', 'Cremas Faciales y Corporales-Cremas de Belleza - Cremas Cuerp y Faciales', 'Protección Solar', 
                  'Talcos-Talco para pies', 'Tampones Femeninos', 'Jabón de Tocador', 'Cepillos Dentales', 'Pastas Dentales', 'Morrales y MAletas Escoalres', 'Lapices de Colores', 'Lapices De Grafito', 
                  'Marcadores', 'Cuadernos', 'Útiles Escolares', 'Estudio de Categorías', 'Corporativa', 'Cross Category', 'Cross Category (Bebés)', 'Cross Category (Desayuno)-Yogurt, Cereal, Pan y Queso', 
                  'Cross Category (Diet y Light)', 'Cross Category (Alimentos Secos)', 'Cross Category (Alimentos)', 'Cross Category (Salsas)-Mayonesas-Ketchup - Salsas Frías', 'Cross Category (Snacks)', 
                  'Demo', 'Flash', 'Holistic View', 'Mezcla para café instantaneo y crema no láctea', 'Mezclas nutricionales y suplementos', 'Consolidado-Multicategory', 'Pantry Check', 'Inventario', 
                  'Leche y Cereales Calientes-Cereales Precocidos y Leche Líquida Blanca','Agua Saborizada'],
    'cod': ['ALCB', 'BEER', 'CARB', 'CWAT', 'COCW', 'COFF', 'CRBE', 'ENDR', 'FLBE', 'GCOF', 'HJUI', 'ITEA', 'ICOF', 'JUNE', 'VEJU', 'WATE', 'CSDW', 'MXCM', 'MXDG', 'MXJM', 'MXJS', 'MXTC', 'JUIC', 'PWDJ', 
               'RFDR', 'RTDJ', 'RTEA', 'SOYB', 'SPDR', 'TEAA', 'YERB', 'BUTT', 'CHEE', 'CMLK', 'CRCH', 'DYOG', 'EMLK', 'FRMM', 'FMLK', 'FRMK', 'LQDM', 'LLFM', 'MARG', 'MCHE', 'MKCR', 'MXDI', 'MXMI', 'MXYD', 
               'PTSS', 'PWDM', 'SYOG', 'MILK', 'YOGH', 'CLOT', 'FOOT', 'SOCK', 'AREP', 'BCER', 'BABF', 'BEAN', 'BISC', 'BOUI', 'BREA', 'BRCR', 'BRDC', 'CERE', 'BURG', 'CCMX', 'CAKE', 'FISH', 'CFAV', 'CRML', 
               'CMLC', 'CBAR', 'CHCK', 'CHOC', 'COCO', 'COLS', 'COMP', 'SPIC', 'CKCH', 'COIL', 'CSAU', 'CNML', 'CNST', 'CNFL', 'CAID', 'DESS', 'DHAM', 'DFNS', 'EBRE', 'EEGG', 'EGGS', 'FLSS', 'FLOU', 'MEAT', 
               'FRDS', 'FRFO', 'HAMS', 'HCER', 'HOTS', 'ICEC', 'IBRE', 'IMPO', 'INOO', 'JAMS', 'KETC', 'LJDR', 'MALT', 'SEAS', 'MAYO', 'MEAT', 'MLKM', 'MXCO', 'MXBS', 'MXSB', 'MXCH', 'MXCC', 'MXSN', 'COBT', 
               'COCF', 'CABB', 'MXEC', 'MXDP', 'MXFR', 'MXFM', 'MXMC', 'MXPS', 'MXSO', 'MXSP', 'MXSW', 'MUST', 'NDCR', 'NOOD', 'NUGG', 'OAFL', 'OLIV', 'PANC', 'PANE', 'PAST', 'PSAU', 'PNOU', 'PORK', 'PPMX', 
               'PWSM', 'PCCE', 'DOUG', 'PPIZ', 'REFR', 'RICE', 'RBIS', 'RTEB', 'RTEM', 'SDRE', 'SALT', 'SLTC', 'SARD', 'SAUS', 'SCHN', 'SNAC', 'SNOO', 'SOUP', 'SOYS', 'SPAG', 'SPCH', 'SUGA', 'SWCO', 'SWSP', 
               'SWEE', 'TOAS', 'TOMA', 'TUNA', 'VMLK', 'WFLO', 'AIRC', 'BARS', 'BLEA', 'CBLK', 'CGLO', 'CLSP', 'CLTO', 'FILT', 'CRHC', 'CRLA', 'CRPA', 'DISH', 'DPAC', 'DRUB', 'FBRF', 'FWAX', 'FDEO', 'FRNP', 
               'GBBG', 'GCLE', 'CLEA', 'INSE', 'KITT', 'LAUN', 'LSTA', 'MXBC', 'MXHC', 'MXCB', 'MXLB', 'MXLD', 'CRTO', 'NAPK', 'PLWF', 'SCOU', 'SOFT', 'STRM', 'TOIP', 'WIPE', 'ANLG', 'FSUP', 'GMED', 'VITA', 'nan', 
               'BATT', 'CGAS', 'PFHH', 'PFIN', 'INKC', 'PETF', 'TELE', 'TILL', 'TOBA', 'ADIP', 'BSHM', 'RAZO', 'BDCR', 'CWIP', 'COMB', 'COND', 'CRHY', 'CRPC', 'DEOD', 'DIAP', 'FCCR', 'FTIS', 'FEMI', 'FRAG', 'HAIR', 
               'HRCO', 'HREM', 'HRST', 'HSTY', 'HRTR', 'LINI', 'MAKE', 'MEDS', 'CRDT', 'MXMH', 'MOWA', 'ORAL', 'SPAD', 'STOW', 'SHAM', 'SHAV', 'SKCR', 'SUNP', 'TALC', 'TAMP', 'TOIL', 'TOOB', 'TOOT', 'BAGS', 'CLPC', 
               'GRPC', 'MRKR', 'NTBK', 'SCHS', 'CSTD', 'CORP', 'CROS', 'CRBA', 'CRBR', 'CRDT', 'CRDF', 'CRFO', 'CRSA', 'CRSN', 'DEMO', 'FLSH', 'HLVW', 'COCP', 'CRSN', 'MULT', 'PCHK', 'STCK', 'MIHC','FLWT'],
})

    #obtém o país,categoria,cesta e fabricante para template e ppt

root=r"C:\Users\70085757\OneDrive - Kantar\Desktop\Fim\Adicionais"

os.chdir(root)

excel='55_MILK_Tetra Pak.xlsx'

file = pd.ExcelFile(os.path.join(root,excel))

W = file.sheet_names

#Obtém o pais cesta categoria fabricante marca e idioma para o qual se fará o estudo
land, cesta, cat, client = pais.loc[pais.cod==int(excel.split('_')[0]),'pais'].iloc[0], categ.loc[categ.cod==excel.split('_')[1],'cest'].iloc[0] ,categ.loc[categ.cod==excel.split('_')[1],'cat'].iloc[0] ,excel.split('_')[2].rsplit('.', 1)[0]

lang= "P" if land=='Brasil' else "E"

modelo= "Modelo_ADD_PT.pptx" if lang =="P" else  "Modelo_ADD_ES.pptx" 

ppt= Presentation(modelo)

brand = W[0][2:]
#Dicionario com as correspondencias dos numeros e o respectivo W
c_w={
('P','1'):'1W - Quando?',
('P','3-1'):'3W - O Que? Tamanhos',
('P','3-2'):'3W - O Que? Marcas',
('P','3-3'):'3W - O Que? Sabores',
('P','4'):'4W - Quem? NSE',
('P','5-1'):'5W - Onde? Regiões',
('P','5-2'):'5W - Onde? Canais',
('P','6'):'Players',

('E','1'):'1W - ¿Cuando?',
('E','3-1'):'3W - ¿Lo Que? Tamaños',
('E','3-2'):'3W - ¿Lo Que? Marcas',
('E','3-3'):'3W - ¿Lo Que? Sabores',
('E','4'):'4W - ¿Quién? NSE',
('E','5-1'):'5W - ¿Dónde? Regiones',
('E','5-2'):'5W - ¿Dónde? Canales',
('E','6'):'Players'
}

#Etiquetas dos slides
labels  ={
('P','Data'):'Data',
('P','MAT'):'Avaliação em Ano Móvel Acumulado',
('P','Var MAT'):"em ano móvel",
('P','comp'):"Concorrência do mercado de: ",

('E','Data'):'Fecha',
('E','MAT'):'Evaluación en Año Móvil Acumulado',
('E','Var MAT'):"en año móvil",
('E','comp'):"Competencia para el mercado de: "}

#Variavel de controle do numero de graficos
c_fig=0

plot=plot_ven()
#---------------------------------------------------------------------------------------------------------------------
for w in W:

    #1- Quando, grafico de variações MAT
    if w[0]=='1':

        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[35])

        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = c_w[(lang,w[0])]+' '+ labels[(lang,'MAT')]+' | ' + w[2:] 
        t.font.bold = True
        t.font.size = Inches(0.35)

        #Obtém pipeline das vendas
        p=int(file.parse(w).columns[1].split("_")[1])

        #Cria a base
        mat=df_mat(file.parse(w),p)
        
        #Elimina linhas com divisão com zero devido ao pipeline
        mat=mat[~np.isinf(mat.iloc[:, 3])]
      
        #Incrementa contador do gráfico
        c_fig+=1
        
        #Insere o gráfico do MAT
        pic=slide.shapes.add_picture(graf_mat(mat,c_fig,p), Inches(0.33), Inches(1.15),width=Inches(10))

        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(5.8), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = "Comentário"
        t.font.size = Inches(0.28)

        #Limpa área de plotagem
        plt.clf()

        print(c_w[(lang,w[0])]+' '+' realizado para' +' '+ w[2:])

    #Outros
    else: 

        #Carrega a base
        df_start=file.parse(w)

        #Lista para as bases
        df_list=[]

        #Verifica se há dados de vendas para obter as bases e o pipeline
        try:
            start,p=[(i-1,col.split("_")[1]) for i,col in enumerate(df_start.columns) if "ventas" in col.lower()][0]
            p_ventas=int(p)
            compras=df_start.iloc[:,:start]
            ventas=df_start.iloc[:,start+1:]
            compras.columns= [compras.columns[0]]+[x+'.C' for x in compras.columns[1:]]
            ventas.columns= [ventas.columns[0][:-2]]+[x.replace('.1','.V') for x in ventas.columns[1:]]
            df_list.append(compras)
            df_list.append(ventas)
            
        except:
            p=0
            df_list.append(df_start)

        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[35])

        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        if w[0] in ['3','5']:
            titulo = c_w[(lang,w[0]+'-'+w[-1])]+' | ' + w[2:-2] 
            t.text = titulo
            t.font.size = Inches(0.35)
        elif w[0]=='6':
            titulo =  c_w[(lang,w[0])] + ' | ' + labels[(lang,'comp')] + cat
            t.text =  titulo
            t.font.size = Inches(0.33)
        else:
            titulo = c_w[(lang,w[0])]+' | '+ w[2:] 
            t.text = titulo 
            t.font.size = Inches(0.35)

        t.font.bold = True

        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text ="Comentário"
        t.font.size= Inches(0.25)

        #Controle posicao tabela aporte
        k=0

        for df in df_list:
            
            #Ajusta as datas
            for i in range(0,len(df)):
                if isinstance(df.iloc[:,0][i],str):
                    df.iloc[i,0]=dt.strptime(df.iloc[i,0], '%b-%y  ')       

            #Tipo do dado
            tipo= df.columns[0]    

            df.rename(columns={df.columns[0]: labels[(lang,'Data')]}, inplace=True)

            #Cria a tabela de aporte
            p=0 if 'compras' in tipo.lower() else p_ventas 

            apo=aporte(df,p,lang,tipo)
            
            #Incrementa contador do gráfico
            c_fig+=1

            #Insere o gráfico da tabela de aporte
            pic = slide.shapes.add_picture(graf_apo(apo,c_fig), Inches(0.33), Inches(4.61+k))

            #Limpa área de plotagem
            plt.clf()

            k=+1

        #Plota vendas e compras juntas
        if plot=="1" and len(df_list)>1:

            #Cria a base para o gráfico de linhas
            df_full = pd.concat([df_list[0],df_list[1].iloc[:,1:]],axis=1)

            #Incrementa contador do gráfico
            c_fig+=1

            #Insere o gráfico de linhas
            pic=slide.shapes.add_picture(line_graf(df_full,p,titulo,c_fig,len(df_list)), Inches(0.33), Inches(1.15),width=Inches(10),height=Inches(10/3))

            #Limpa área de plotagem
            plt.clf()
        #Plota vendas e compras separadas
        elif plot=="2" and len(df_list)>1:

            m=0
            for df in df_list:
                #Incrementa contador do gráfico
                c_fig+=1

                tipo = 'Compras' if '.C' in df.columns[1] else 'Ventas'
                
                #Insere o gráfico de linhas
                pic=slide.shapes.add_picture(line_graf(df,p,titulo+' '+tipo,c_fig,len(df_list)-1), Inches(0.15+m), Inches(1.15),width=Inches(10)*0.63,height=Inches(10/3)*0.63)

                #Limpa área de plotagem
                plt.clf()

                m=+6.75
        #Plota somente compras
        else:
            #Incrementa contador do gráfico
            c_fig+=1

            #Insere o gráfico de linhas
            pic=slide.shapes.add_picture(line_graf(df_list[0],p,titulo,c_fig,len(df_list)), Inches(0.33), Inches(1.15),width=Inches(10),height=Inches(10/3))

            #Limpa área de plotagem
            plt.clf()

        if w[0] in ['3','5']:
            print(c_w[(lang,w[0]+'-'+w[-1])]+' realizado para ' + w[2:-2]) 
        elif w[0]=='6':
            print(c_w[(lang,w[0])] + ' realizado para ' + labels[(lang,'comp')] + cat)
        else:
            print(c_w[(lang,w[0])]+' realizado para '+ w[2:])


#Referencia da base
ref =  dt.strptime(df.iloc[:,0][-1:].iloc[0], '%b-%y  ').strftime('%m-%y')  if isinstance(df.iloc[:,0][-1:].iloc[0],str)  else df.iloc[:,0][-1:].iloc[0].strftime('%m-%y')   

slide=ppt.slides[0]

#Fabricante
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(3.77), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
t.text = "Numerator vs " + client
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.44), True
font.color.rgb = RGBColor(255, 255, 255)

#Referencia
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(5.98), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
t.text = "Corte a "+ref
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.32), True
font.color.rgb = RGBColor(255, 255, 255)

ppt.save(land+'-'+cat+'-'+client+'-'+brand+'-5W1H-'+ref+'.pptx')

fim = dt.now()

t = int((fim - agora).total_seconds())
print(f"Tiempo de ejecución : {t//60} min {t%60} s" if t >= 60 else f"Tiempo de ejecución : {t} s")