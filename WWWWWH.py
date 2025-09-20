# -*- coding: utf-8 -*-

#Bibliotecas necessarias
#---------------------------------------------------------------------------------------------------------------------
import pandas as pd
import numpy as np
import warnings
from datetime import datetime as dt
import os
import io
from matplotlib import pyplot as plt
from datetime import datetime as dt
from matplotlib.ticker import FuncFormatter
from matplotlib import colors as mcolors
from matplotlib.patches import Rectangle
from decimal import Decimal, ROUND_CEILING, ROUND_DOWN, ROUND_FLOOR, ROUND_HALF_DOWN, ROUND_HALF_EVEN, ROUND_HALF_UP, ROUND_UP, ROUND_05UP
from pptx import Presentation 
from pptx.util import Inches, Cm
from pptx.dml.color import RGBColor
from typing import NamedTuple
from pathlib import Path

COLOR_BLUE = '\033[94m'
COLOR_YELLOW = '\033[93m'
COLOR_GREEN = '\033[92m'
COLOR_RED = '\033[91m'
COLOR_RESET = '\033[0m'

COLOR_QUESTION = '\033[38;5;37m'

HEADER_COLOR_PRIMARY = '#286B72'
HEADER_COLOR_SECONDARY = '#3EBBC7'
HEADER_FONT_COLOR = '#FFFFFF'
HEADER_TOTAL_FILL = '#E8F0FF'
HEADER_FIRST_COL_FILL = '#F5F5F5'
TABLE_TEXT_PRIMARY = '#1F1F1F'
TABLE_POSITIVE_COLOR = '#27AE60'
TABLE_NEGATIVE_COLOR = '#C0392B'
TABLE_GRID_COLOR = '#1A1A1A'
VOLUME_BAR_START = '#1452B8'
VOLUME_BAR_END = '#8CB8FF'
VOLUME_BAR_ALPHA = 0.85
VOLUME_BAR_BASE_ALPHA = 0.18
LINE_COLOR_CLIENT = '#2C3E50'
LINE_COLOR_NUMERATOR = "#D4AC0D"
BAR_COLOR_VARIATION_CLIENT = '#7F8C8D'
BAR_COLOR_VARIATION_NUMERATOR = '#F1C40F'
BAR_EDGE_COLOR = 'black'
ANNOTATION_BOX_FACE = '#F2F2F2'
ANNOTATION_BOX_EDGE = 'black'

MONTH_NAMES = {
    'P': ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho', 'Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro'],
    'E': ['Enero', 'Febrero', 'Marzo', 'Abril', 'Mayo', 'Junio', 'Julio', 'Agosto', 'Septiembre', 'Octubre', 'Noviembre', 'Diciembre']
}


def colorize(text: str, color: str = COLOR_BLUE) -> str:
    return f"{color}{text}{COLOR_RESET}"


def print_colored(text: str, color: str = COLOR_BLUE) -> None:
    print(colorize(text, color))
def available_width(presentation, left=Inches(0), right=Inches(0)):
    return presentation.slide_width - int(left) - int(right)


EMU_PER_INCH = 914400
DEFAULT_LINE_CHART_RATIO = 3
TABLE_TARGET_HEIGHT_CM = 4.0


def emu_to_inches(value: int) -> float:
    return float(value) / EMU_PER_INCH


pd.set_option('future.no_silent_downcasting', True)
pd.set_option('mode.chained_assignment', None)
warnings.filterwarnings('ignore')

agora = dt.now()

#Funcao que prepara os dados para criação do gráfico MAT
def df_mat(df,p):
    
    v1 = pd.DataFrame([(df.iloc[i-12-p:i-p,1].sum()/df.iloc[i-24-p:i-12-p,1].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-in'])
    v2 = pd.DataFrame([(df.iloc[i-12:i,2].sum()/df.iloc[i-24:i-12,2].sum()) - 1 if i >= 24 else np.nan for i in range(12, len(df)+1)],columns=['Var Sell-out'])

    ac1, ac2 = [pd.DataFrame(df[col].rolling(window=12).sum().dropna()) for col in df.columns[1:]]
    v1.index,v2.index=v1.index+11,v2.index+11
    v1 ,v2 = v1.dropna(), v2.dropna()

    mat =pd.DataFrame(df.iloc[:,0].copy())
    mat= mat.join(ac1.join(ac2)[-len(v1):].join([v1]).join([v2]))
    mat= mat[mat.notna().all(axis=1).idxmax():]

    return mat

#Grafico de MAT 
def graf_mat (mat,c_fig,p):

    #Cria figura e área plotável   
        altura = 4 # Puedes ajustar este valor para cambiar la altura
        # Cerrar cualquier figura existente con el mismo número para que figsize sea aplicado
        try:
            plt.close(c_fig)
        except Exception:
            pass
        fig = plt.figure(num=c_fig, figsize=(15, altura))
        ax = fig.add_subplot(1, 1, 1)
        
        #Cria o eixo do tempo e os acumulados juntamente com as variações
        ran=mat.iloc[:,0].copy()
        ran=[x.strftime('%m-%y') for x in ran]
        ac1 = mat.iloc[:,1].copy()
        ac2 = mat.iloc[:,2].copy()
        v1= mat.iloc[:,3].copy()
        v2= mat.iloc[:,4].copy()

        #Plota os acumulados em linhas    
        l1=ax.plot(ran,ac1, color=LINE_COLOR_CLIENT, linewidth=4.2, label='Acumulado Cliente')
        l2=ax.plot(ran,ac2, color=LINE_COLOR_NUMERATOR, linewidth=4.2, label='Acumulado Numerator') 

        #Plota as variações em barras
        ax2= ax.twinx()
        b1 = ax2.bar(np.arange(len(ran)) - 0.3, v1.values, 0.3, color=BAR_COLOR_VARIATION_CLIENT, edgecolor=BAR_EDGE_COLOR, label='Var % Cliente Pipeline: '+ str(p) +' '+labels[(lang,'Var MAT')])
        b2 = ax2.bar(np.arange(len(ran)) + 0.3, v2.values, 0.3, color=BAR_COLOR_VARIATION_NUMERATOR, edgecolor=BAR_EDGE_COLOR, label='Var % Numerator '+labels[(lang,'Var MAT')])

        ax.set_xticks(np.arange(len(ran)))
        ax.set_xticklabels(ran, rotation=30)
        ax2.tick_params(left=False, labelleft=False, top=False, labeltop=False,
            right=False, labelright=False, bottom=False, labelbottom=False)

        for v in [v1,v2]:
            for x,y in zip(np.arange(len(v1))+0.2,v):
                label = "{:.1f}%".format(y)
                bbox_props_white=dict(facecolor=ANNOTATION_BOX_FACE,edgecolor=ANNOTATION_BOX_EDGE)
                plt.annotate(f"{y*100:.1f}%", (x, y), textcoords="offset points", xytext=(0, 10), ha='center', color='red' if y < 0 else 'green', size=9, bbox=bbox_props_white)

        
        #Cria a legenda 
        lns = l1 + l2 + [b1, b2]
        labs = [l.get_label() for l in l1 + l2] + [b1.get_label(), b2.get_label()]
        ax.legend(lns, labs, loc='upper center', bbox_to_anchor=(0.475,-0.115), borderaxespad=0.05, frameon=True, prop={'size': 12}, ncol=2)
        img_stream = io.BytesIO()


        #Cria o titulo
        ax.set_title(labels[(lang,'MAT')]+' | ' + w[2:], size=18, pad=10)

        #Salva o grafico
        # Guardar respetando el tamaño definido (no usar bbox_inches='tight')
        fig.savefig(img_stream, format='png', transparent=True)

        #Insere o grafico
        img_stream.seek(0) 

        # Cerrar figura para liberar recursos y evitar reutilización
        try:
            plt.close(fig)
        except Exception:
            pass

        return img_stream

#Grafico de Lineas
def line_graf(df, p, title, c_fig,ven, width_emu=None, height_emu=None):
    # Ajusta o tamanho da figura para combinar com o espaco reservado no slide
    width_inches = emu_to_inches(width_emu) if width_emu is not None else None
    height_inches = emu_to_inches(height_emu) if height_emu is not None else None

    if width_inches is None and height_inches is None:
        # Ancho dinámico según cantidad de datos (en pulgadas)
        n_points = len(df)
        ancho = max(15, n_points * 0.5)  # 0.5 puede ajustarse para más/menos espacio (pulgadas)
        # Altura fija de 10 cm -> convertir a pulgadas
        altura = emu_to_inches(Cm(10))
        figsize = (ancho, altura)
    else:
        if width_inches is None:
            width_inches = height_inches * DEFAULT_LINE_CHART_RATIO
        elif height_inches is None:
            height_inches = width_inches / DEFAULT_LINE_CHART_RATIO
        figsize = (width_inches, height_inches)

    fig = plt.figure(num=c_fig, figsize=figsize)
    ax = fig.add_subplot(1, 1, 1)
    # Escalar elementos (fuentes, linewidth) según la altura para que el contenido se adapte
    # Se toma como referencia la altura original de 5 pulgadas usada previamente
    ref_height = 5.0
    actual_height = figsize[1]
    scale = actual_height / ref_height if ref_height else 1.0
    base_linewidth = 2.0 * max(0.6, scale)
    title_base_size = max(10, int(18 * scale))
    legend_base_size = max(8, int(10 * scale))
    xtick_size = max(8, int(10 * scale))
    
    aux = df.copy()
    ran = [x.strftime('%m-%y') for x in aux.iloc[:, 0]]

    if ven>1:
        colunas = aux.columns[1:]  
        comp=colunas[:len(colunas)//2]
        vent=colunas[len(colunas)//2:]
        cmap = plt.get_cmap('tab20')
        cor_map = {base: cmap(i % 20) for i, base in enumerate(comp)}
        df_cor = pd.DataFrame({
        'Compras': comp,
        'Ventas': vent,
        'Cor': [cor_map[base] for base in comp]
        })
    else:
        colunas = aux.columns[1:]  
        comp=colunas
        cmap = plt.get_cmap('tab20')
        cor_map = {base: cmap(i % 20) for i, base in enumerate(comp)}
        df_cor = pd.DataFrame({
        'Compras': comp,
        'Ventas': comp,
        'Cor': [cor_map[base] for base in comp]
        })
    lns = []

    for col in colunas:
        if ven>1: 
            estilo = '-' if '.v' in col.lower() else '--' 
        else: 
            estilo='-'

        tipo = 'Ventas' if '.v' in col.lower() else 'Compras'
        cor = df_cor['Cor'][df_cor[tipo]==col].iloc[0]

        y = aux[col].values
        line, = ax.plot(ran[p:], y[p:], color=cor, linewidth=base_linewidth, linestyle=estilo, label=col)
        lns.append(line)

    # Eixo X,Y
    plt.xticks(rotation=30, fontsize=xtick_size)
    plt.ylim(0)

    # Legenda
    labs = [l.get_label() for l in lns]
    ax.legend(lns, labs, loc='upper center', bbox_to_anchor=(0.5, -0.15), borderaxespad=0, frameon=False,
              prop={'size': legend_base_size}, ncol=max(1, len(colunas)/2))
    # plt.tight_layout() # Eliminado para respetar el tamaño definido

    # Título (escalado)
    plt.title(title, size=title_base_size, pad=10)

    # Salva imagem
    img_stream = io.BytesIO()
    fig.savefig(img_stream, format='png', transparent=True)
    img_stream.seek(0)

    # Cerrar figura para liberar memoria y evitar reutilización
    try:
        plt.close(fig)
    except Exception:
        pass

    return img_stream

#Função que cria a tabela de aporte
def aporte(df,p,lang,tipo):

        aux = df.copy()

        aux['Total'] = aux.iloc[:len(df),1:].sum(axis=1)

        apo=pd.DataFrame(columns=[tipo] + aux.columns[1:].tolist())
        #Vol ultimo MAT
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-12-p,labels[(lang,'Data')]].strftime('%b-%y') ] + [aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p,aux.shape[1]-1].sum() for col in range(1,len(aux.columns) )]
        #Vol MAT atual
        apo.loc[len(apo)] = [str('Vol' )+" "+aux.loc[len(aux)-1-p,labels[(lang,'Data')]].strftime('%b-%y') ] + [aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-12-p:len(aux)-p,aux.shape[1]-1].values.sum() for col in range(1,len(aux.columns))]
        #Var
        apo.loc[len(apo)] = ["Var %"] +  [ aux.iloc[len(aux)-12-p:len(aux)-p, col].sum() / aux.iloc[len(aux)-24-p:len(aux)-12-p, col].sum()-1 for col in range(1,len(aux.columns))]
        #Aporte
        apo.loc[len(apo)] = ["Aporte"] + [apo.iloc[-1, col] * apo.iloc[0, col] for col in range(1,len(aux.columns))]

        #Formatação do volume
        apo.iloc[:2, 1:] = apo.iloc[:2, 1:].applymap(lambda x: f"{round(x * 100, 1)}%")
        #Formatação da variação e aporte
        apo.iloc[2:, 1:] = apo.iloc[2:, 1:].applymap(lambda x: f"{round(x * 100, 2)}%")

        return apo

#Função que cria o gráfico tabela de aporte

def graf_apo(apo,c_fig):
    fig, ax = plt.subplots(num=c_fig)
    row_height = 0.45
    col_width = 1.0

    n_rows, n_cols = apo.shape
    fig_width = col_width * n_cols
    fig_height = row_height * (n_rows + 0.2)
    fig.set_size_inches(fig_width, fig_height)
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax.axis('off')
    ax.set_facecolor('white')

    table = ax.table(
        cellText=apo.values,
        colLabels=apo.columns,
        loc='center',
        cellLoc='center'
    )

    for i, _ in enumerate(apo.columns):
        table.auto_set_column_width(i)

    table.scale(1, 1.2)
    table.auto_set_font_size(False)
    table.set_fontsize(11)

    header_main = HEADER_COLOR_PRIMARY
    header_secondary = HEADER_COLOR_SECONDARY
    total_fill = HEADER_TOTAL_FILL
    first_col_fill = HEADER_FIRST_COL_FILL
    positive_color = TABLE_POSITIVE_COLOR
    negative_color = TABLE_NEGATIVE_COLOR
    bar_padding_ratio = 0.08
    bar_height_ratio = 0.55

    data_columns = [idx for idx in range(1, n_cols) if idx != n_cols - 1]
    start_rgb = np.array(mcolors.to_rgb(VOLUME_BAR_START))
    end_rgb = np.array(mcolors.to_rgb(VOLUME_BAR_END))

    def share_colors(count: int):
        if count <= 0:
            return []
        if count == 1:
            return [mcolors.to_hex(np.clip(start_rgb, 0, 1))]
        colors = []
        for i in range(count):
            ratio = i / (count - 1)
            rgb = np.clip(start_rgb + (end_rgb - start_rgb) * ratio, 0, 1)
            colors.append(mcolors.to_hex(rgb))
        return colors

    column_colors = share_colors(len(data_columns))
    white_rgb = np.array([1.0, 1.0, 1.0])
    volume_row_indexes = [idx for idx, label in enumerate(apo.iloc[:, 0]) if str(label).lower().startswith('vol')]

    def to_percentage(value):
        try:
            text_value = str(value).strip().replace('%', '').replace(',', '.')
            if not text_value:
                return None
            return max(0.0, float(text_value) / 100.0)
        except Exception:
            return None

    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor(TABLE_GRID_COLOR)
        cell.set_linewidth(1)
        cell.get_text().set_zorder(3)

        if row == 0:
            text = cell.get_text()
            text.set_weight('bold')
            text.set_color(HEADER_FONT_COLOR)
            if col == 0 or col == n_cols - 1:
                cell.set_facecolor(header_main)
            else:
                cell.set_facecolor(header_secondary)
        else:
            label = str(apo.iloc[row - 1, 0])
            text = cell.get_text()

            if col == 0:
                cell.set_facecolor(first_col_fill)
                text.set_weight('bold')
                text.set_color(TABLE_TEXT_PRIMARY)
            else:
                if col == n_cols - 1:
                    cell.set_facecolor(total_fill)
                else:
                    cell.set_facecolor('white')

                if label.lower().startswith('vol'):
                    text.set_color(TABLE_TEXT_PRIMARY)
                else:
                    value = apo.iloc[row - 1, col]
                    try:
                        numeric = float(str(value).replace('%', '').replace(',', '.'))
                        text.set_color(positive_color if numeric >= 0 else negative_color)
                    except Exception:
                        pass

    for df_row_idx in volume_row_indexes:
        table_row = df_row_idx + 1
        for rel_idx, col in enumerate(data_columns):
            raw_value = apo.iloc[df_row_idx, col]
            percent = to_percentage(raw_value)
            if percent is None:
                continue
            cell = table[(table_row, col)]
            bar_color = column_colors[rel_idx] if rel_idx < len(column_colors) else VOLUME_BAR_END
            base_rgb = np.array(mcolors.to_rgb(bar_color))
            intensity = min(percent, 1.0)
            blended_rgb = white_rgb * (1.0 - intensity) + base_rgb * intensity
            cell.set_facecolor(blended_rgb)
            cell.get_text().set_color(TABLE_TEXT_PRIMARY)

    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    table_bbox = table.get_tightbbox(renderer).transformed(fig.dpi_scale_trans.inverted())
    target_height_in = TABLE_TARGET_HEIGHT_CM / 2.54
    if table_bbox.height > 0:
        scale = target_height_in / table_bbox.height
        if abs(scale - 1.0) > 1e-3:
            current_width, current_height = fig.get_size_inches()
            fig.set_size_inches(current_width * scale, current_height * scale)
            fig.canvas.draw()
            renderer = fig.canvas.get_renderer()
            table_bbox = table.get_tightbbox(renderer).transformed(fig.dpi_scale_trans.inverted())
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=400, bbox_inches=table_bbox, pad_inches=0)
    buf.seek(0)
    return buf

def select_excel_file(base_dir: Path) -> str:
    excel_files = sorted([p for p in base_dir.glob('*.xlsx') if not p.name.startswith('~$')])

    def metadata_for(file_path: Path) -> str:
        try:
            parts = file_path.stem.split('_')
            country_code = int(parts[0]) if parts else None
            category_code = parts[1] if len(parts) > 1 else None
            pais_name = pais.loc[pais['cod'] == country_code, 'pais'].iloc[0] if country_code is not None else 'Pais N/D'
            cesta = categ.loc[categ['cod'] == category_code, 'cest'].iloc[0] if category_code else 'Cesta N/D'
            categoria = categ.loc[categ['cod'] == category_code, 'cat'].iloc[0] if category_code else 'Categoria N/D'
            return f"{pais_name} | {cesta} | {categoria}"
        except Exception:
            return 'Metadata no disponible'

    if not excel_files:
        raise FileNotFoundError(colorize(f'No se encontraron archivos .xlsx en {base_dir}', COLOR_RED))

    if len(excel_files) == 1:
        meta = metadata_for(excel_files[0])
        print_colored(f"Archivo encontrado: {excel_files[0].name} {colorize('[' + meta + ']', COLOR_YELLOW)}", COLOR_BLUE)
        return excel_files[0].name

    print_colored('Archivos Excel disponibles:', COLOR_QUESTION)
    for idx, archivo in enumerate(excel_files, 1):
        meta = metadata_for(archivo)
        print(f"{colorize(f'  {idx}. {archivo.name}', COLOR_BLUE)} {colorize('[' + meta + ']', COLOR_YELLOW)}")

    prompt = colorize(f"Seleccione el numero del archivo (1-{len(excel_files)}). Enter para {excel_files[0].name}: ", COLOR_QUESTION)
    while True:
        choice = input(prompt).strip()
        if not choice:
            return excel_files[0].name
        if choice.isdigit():
            idx = int(choice)
            if 1 <= idx <= len(excel_files):
                return excel_files[idx - 1].name
        matching = [archivo.name for archivo in excel_files if archivo.name.lower() == choice.lower()]
        if matching:
            return matching[0]
        print_colored('Entrada invalida. Intente nuevamente.', COLOR_RED)


def configure_cover_slide(presentation, lang):
    title_by_lang = {
        'P': 'Adicionais de Cobertura - 5W1H',
        'E': 'Adicionales de Cobertura - 5W1H'
    }
    cover_slide = presentation.slides[0]
    target_text = title_by_lang.get(lang, title_by_lang['E'])
    fallback_geometry = (Inches(0.42), Inches(2.0), Inches(10), Inches(0.8))
    geometry = None
    for shape in list(cover_slide.shapes):
        if not getattr(shape, 'has_text_frame', False):
            continue
        text_value = shape.text.strip()
        if 'Cobertura' in text_value and '5W1H' in text_value:
            geometry = (shape.left, shape.top, shape.width, shape.height)
            cover_slide.shapes._spTree.remove(shape._element)
    left, top, width, height = geometry if geometry else fallback_geometry
    text_box = cover_slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = target_text
    run = paragraph.runs[0]
    font = run.font
    font.name = 'Arial'
    font.size = Inches(0.45)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)


def plot_ven():
    options = [
        ("1 - Plotear Ventas y Compras Juntas", "1"),
        ("2 - Plotear Ventas y Compras Separadas", "2"),
        ("3 - No hay W con Ventas en esa base", "3"),
    ]

    print_colored('\nOpciones de grafico disponibles:', COLOR_QUESTION)
    for idx, (texto, _) in enumerate(options, 1):
        print_colored(f"  {idx}. {texto}", COLOR_BLUE)

    prompt = colorize(f"Seleccione el numero de la opcion (1-{len(options)}). Enter para {options[0][1]}: ", COLOR_QUESTION)
    while True:
        choice = input(prompt).strip()
        if not choice:
            return options[0][1]
        if choice.isdigit():
            idx = int(choice)
            if 1 <= idx <= len(options):
                return options[idx - 1][1]
        for _, valor in options:
            if choice == valor:
                return valor
        print_colored('Entrada invalida. Intente nuevamente.', COLOR_RED)

#Codigos paises
pais = pd.DataFrame(
    {'cod': [10, 54, 91, 55, 12, 56, 57, 93, 52, 51, 66, 63, 62, 64, 65, 67, 69],
     'pais': ['LatAm', 'Argentina', 'Bolivia', 'Brasil', 'CAM', 'Chile', 'Colombia', 'Ecuador', 'Mexico', 'Peru', 
              'Costa Rica', 'El Salvador', 'Guatemala', 'Honduras', 'Nicaragua', 'Panama', 'Republica Dominicana']
    }
)

#Codigos categorias
categ = pd.DataFrame({
    'cest': ['Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 
               'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Bebidas', 'Lacteos', 
               'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 
               'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Lacteos', 'Ropas y Calzados', 'Ropas y Calzados', 'Ropas y Calzados', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 'Alimentos', 
               'Alimentos', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 
               'Cuidado del Hogar', 'Cuidado del Hogar', 'Cuidado del Hogar', 'OTC', 'OTC', 'OTC', 'OTC', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 'Otros', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 
               'Cuidado Personal', 'Cuidado Personal', 'Cuidado Personal', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 'Material Escolar', 
               'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 'Diversos', 
               'Diversos', 'Diversos', 'Diversos', 'Diversos','Alimentos'],
    'cat': ['Bebidas Alcohólicas', 'Cervezas', 'Bebidas Gaseosas', 'Agua Gasificada', 'Água de Coco', 'Café-Consolidado de Café', 'Cross Category (Bebidas)', 'Bebidas Energéticas', 
                  'Bebidas Saborizadas Sin Gas', 'Café Tostado y Molido', 'Jugos Caseros', 'Té Helado', 'Café Instantáneo-Café Sucedáneo', 'Jugos y Nectares', 'Zumos de Vegetales', 'Agua Natural', 
                  'Gaseosas + Aguas', 'Mixta Café+Malta', 'Mixta Dolce Gusto-Mixta Té Helado + Café + Modificadores', 'Mixta Jugos y Leches', 'Mixta Jugos Líquidos + Bebidas de Soja', 'Mixta Té+Café', 
                  'Jugos Liquidos-Jugos Polvo', 'Refrescos en Polvo-Jugos - Bebidas Instantáneas En Polvo - Jugos Polvo', 'Bebidas Refrescantes', 'Refrescos Líquidos-Jugos Líquidos', 'Té Líquido - Listo para Tomar',
                  'Bebidas de Soja', 'Bebidas Isotónicas', 'Té e Infusiones-Te-Infusión Hierbas', 'Yerba Mate', 'Manteca', 'Queso Fresco y para Untar', 'Leche Condensada', 'Queso Untable', 'Yoghurt p-beber', 
                  'Leche Culinaria-Leche Evaporada', 'Leche Fermentada', 'Leche Líquida Saborizada-Leche Líquida Con Sabor', 'Fórmulas Infantiles', 'Leche Líquida', 'Leche Larga Vida', 'Margarina', 'Queso Fundido', 
                  'Crema de Leche', 'Mixta Lácteos-Postre+Leches+Yogurt', 'Mixta Leches', 'Mixta Yoghurt+Postres', 'Petit Suisse', 'Leche en Polvo', 'Yoghurt p-comer', 'Leche-Leche Líquida Blanca - Leche Liq. Natural', 
                  'Yoghurt', 'Ropas', 'Calzados', 'Medias-Calcetines', 'Arepas', 'Cereales Infantiles', 'Nutrición Infantil-Colados y Picados', 'Frijoles', 'Galletas', 'Caldos-Caldos y Sazonadores', 'Pan', 
                  'Apanados-Empanizadores', 'Empanados', 'Cereales-Cereales Desayuno-Avenas y Cereales', 'Hamburguesas', 'Mezclas Listas para Tortas-Preparados Base Harina Trigo', 'Queques-Ponques Industrializados', 
                  'Conservas De Pescado', 'Conservas de Frutas y Verduras', 'Dulce de Leche-Manjar', 'Alfajores', 'Barras de Cereal', 'Pollo', 'Chocolate', 'Chocolate de Taza-Achocolatados - Cocoas', 'Salsas Frías', 
                  'Compotas', 'Condimentos y Especias', 'Chocolate de Mesa', 'Aceite-Aceites Comestibles', 'Salsas Listas-Salsas Caseras Envasadas', 'Grano, Harina y Masa de Maíz', 'Fécula de Maíz', 'Harina De Maíz', 
                  'Ayudantes Culinarios', 'Postres Preparados', 'Jamón Endiablado', 'Semillas y Frutos Secos', 'Pan de Pascua', 'Huevos de Páscua', 'Huevos', 'Flash Cecinas', 'Harinas', 'Carne Fresca', 
                  'Platos Listos Congelados','Alimentos Congelados', 'Jamones', 'Cereales Calientes-Cereales Precocidos', 'Salsas Picantes', 'Helados', 'Pan Industrializado', 'Puré Instantáneo', 
                  'Fideos Instantáneos', 'Mermeladas', 'Ketchup', 'Jugo de Limon Adereso', 'Maltas', 'Adobos - Sazonadores', 'Mayonesa', 'Cárnicos', 'Modificadores de Leche-Saborizadores p-leche', 
                  'Mixta Cereales Infantiles+Avenas', 'Mixta Caldos + Saborizantes', 'Mixta Caldos + Sopas', 'Mixta Cereales + Cereales Calientes', 'Mixta Chocolate + Manjar', 'Galletas, snacks y mini tostadas', 
                  'Aceites + Mantecas', 'Aceites + Conservas De Pescado', 'Ayudantes Culinarios + Bolsa de Hornear', 'Mixta Huevos de Páscua + Chocolates', 'Mixta Platos Listos Congelados + Pasta', 
                  'Mixta Platos Congelados y Listos para Comer', 'Mixta Alimentos Congelados + Margarina', 'Mixta Modificadores + Cocoa', 'Mixta Pastas', 'Mixta Sopas+Cremas+Ramen', 
                  'Mixta Margarina + Mayonesa + Queso Crema', 'Mixta Azúcar+Endulzantes', 'Mostaza', 'Sustitutos de Crema', 'Fideos', 'Nuggets', 'Avena en hojuelas-liquidas', 'Aceitunas', 'Tortilla', 
                  'Panetón', 'Pastas', 'Salsas para Pasta', 'Turrón de maní', 'Carne Porcina', 'Postres en Polvo-Postres para Preparar - Horneables-Gelificables', 'Leche de Soya en Polvo', 'Cereales Precocidos', 
                  'Masas Frescas-Tapas Empanadas y Tarta', 'Pre-Pizzas', 'Meriendas listas', 'Arroz', 'Galletas de Arroz', 'Frijoles Procesados', 'Pratos Prontos - Comidas Listas', 'Aderezos para Ensalada', 
                  'Sal', 'Galletas Saladas-Galletas No Dulce', 'Sardina Envasada', 'Cecinas', 'Milanesas', 'Snacks', 'Fideos Sopa', 'Sopas-Sopas Cremas', 'Siyau', 'Tallarines-Spaguetti', 'Chocolate para Untar', 
                  'Azucar', 'Galletas Dulces', 'Untables Dulces', 'Endulzantes', 'Torradas - Tostadas', 'Salsas de Tomate', 'Atún Envasado', 'Leche Vegetal', 'Harinas de trigo', 
                  'Ambientadores-Desodorante Ambiental', 'Jabón en Barra-Jabón de lavar', 'Cloro-Lavandinas-Lejías-Blanqueadores', 'Pastillas para Inodoro', 'Guantes de látex', 
                  'Esponjas de Limpieza-Esponjas y paños', 'Utensilios de Limpieza', 'Filtros de Café', 'Cross Category (Limpiadores Domesticos)', 'Cross Category (Lavandería)', 
                  'Cross Category (Productos de Papel)', 'Lavavajillas-Lavaplatos - Lavalozas mano', 'Empaques domésticos-Bolsas plásticas-Plástico Adherente-Papel encerado-Papel aluminio', 
                  'Destapacañerias', 'Perfumantes para Ropa-Perfumes para Ropa', 'Cera p-pisos', 'Desodorante para Pies', 'Lustramuebles', 'Bolsas de Basura', 'Limpiadores verdes', 
                  'Limpiadores-Limpiadores y Desinfectantes', 'Insecticidas-Raticidas', 'Toallas de papel-Papel Toalla - Toallas de Cocina - Rollos Absorbentes de Papel', 'Detergentes para ropa', 
                  'Apresto', 'Mixta Pastillas para Inodoro + Limpiadores', 'Mixta Home Care-Cloro-Limpiadores-Ceras-Ambientadores', 'Mixta Limpiadores + Cloro', 'Mixta Detergentes + Cloro', 
                  'Mixta Detergentes + Lavavajillas', 'Pañitos + Papel Higienico', 'Servilletas', 'Film plastico e papel aluminio', 'Esponjas de Acero', 'Suavizantes de Ropa', 'Quitamanchas-Desmanchadores', 
                  'Papel Higiénico', 'Paños de Limpieza', 'Analgésicos-Painkillers', 'Suplementos alimentares', 'Gastrointestinales-Efervescentes', 'Vitaminas y Calcio', 'Categoría Desconocida', 
                  'Pilas-Baterías', 'Combustible Gas', 'Panel Financiero de Hogares', 'Panel Financiero de Hogares', 'Cartuchos de Tintas', 'Alimento para Mascota-Alim.p - perro - gato', 
                  'Telecomunicaciones - Convergencia', 'Tickets - Till Rolls', 'Tabaco - Cigarrillos', 'Incontinencia de Adultos', 'Shampoo Infantil', 'Maquinas de Afeitar', 'Cremas Corporales', 
                  'Paños Húmedos', 'Cremas para Peinar', 'Acondicionador-Bálsamo', 'Cross Category (Higiene)', 'Cross Category (Personal Care)', 'Desodorantes', 'Pañales-Pañales Desechables', 
                  'Cremas Faciales', 'Pañuelos Faciales', 'Protección Femenina-Toallas Femeninas', 'Fragancias', 'Cuidado del Cabello-Hair Care', 
                  'Tintes para el Cabello-Tintes - Tintura - Tintes y Coloración para el cabello', 'Depilación', 'Alisadores para el Cabello', 
                  'Fijadores para el Cabello-Modeladores-Gel-Fijadores para el cabello', 'Tratamientos para el Cabello', 'Óleo Calcáreo', 'Maquillaje-Cosméticos', 'Jabón Medicinal', 
                  'Pañitos + Pañales', 'Mixta Make Up+Tinturas', 'Enjuague Bucal-Refrescante Bucal', 'Cuidado Bucal', 'Protectores Femeninos', 'Toallas Femininas', 'Shampoo', 
                  'Afeitado-Crema afeitar-Loción de afeitar-Pord. Antes del afeitado', 'Cremas Faciales y Corporales-Cremas de Belleza - Cremas Cuerp y Faciales', 'Protección Solar', 
                  'Talcos-Talco para pies', 'Tampones Femeninos', 'Jabón de Tocador', 'Cepillos Dentales', 'Pastas Dentales', 'Morrales y MAletas Escoalres', 'Lapices de Colores', 'Lapices De Grafito', 
                  'Marcadores', 'Cuadernos', 'Útiles Escolares', 'Estudio de Categorías', 'Corporativa', 'Cross Category', 'Cross Category (Bebés)', 'Cross Category (Desayuno)-Yogurt, Cereal, Pan y Queso', 
                  'Cross Category (Diet y Light)', 'Cross Category (Alimentos Secos)', 'Cross Category (Alimentos)', 'Cross Category (Salsas)-Mayonesas-Ketchup - Salsas Frías', 'Cross Category (Snacks)', 
                  'Demo', 'Flash', 'Holistic View', 'Mezcla para café instantaneo y crema no láctea', 'Mezclas nutricionales y suplementos', 'Consolidado-Multicategory', 'Pantry Check', 'Inventario', 
                  'Leche y Cereales Calientes-Cereales Precocidos y Leche Líquida Blanca','Agua Saborizada'],
    'cod': ['ALCB', 'BEER', 'CARB', 'CWAT', 'COCW', 'COFF', 'CRBE', 'ENDR', 'FLBE', 'GCOF', 'HJUI', 'ITEA', 'ICOF', 'JUNE', 'VEJU', 'WATE', 'CSDW', 'MXCM', 'MXDG', 'MXJM', 'MXJS', 'MXTC', 'JUIC', 'PWDJ', 
               'RFDR', 'RTDJ', 'RTEA', 'SOYB', 'SPDR', 'TEAA', 'YERB', 'BUTT', 'CHEE', 'CMLK', 'CRCH', 'DYOG', 'EMLK', 'FRMM', 'FMLK', 'FRMK', 'LQDM', 'LLFM', 'MARG', 'MCHE', 'MKCR', 'MXDI', 'MXMI', 'MXYD', 
               'PTSS', 'PWDM', 'SYOG', 'MILK', 'YOGH', 'CLOT', 'FOOT', 'SOCK', 'AREP', 'BCER', 'BABF', 'BEAN', 'BISC', 'BOUI', 'BREA', 'BRCR', 'BRDC', 'CERE', 'BURG', 'CCMX', 'CAKE', 'FISH', 'CFAV', 'CRML', 
               'CMLC', 'CBAR', 'CHCK', 'CHOC', 'COCO', 'COLS', 'COMP', 'SPIC', 'CKCH', 'COIL', 'CSAU', 'CNML', 'CNST', 'CNFL', 'CAID', 'DESS', 'DHAM', 'DFNS', 'EBRE', 'EEGG', 'EGGS', 'FLSS', 'FLOU', 'MEAT', 
               'FRDS', 'FRFO', 'HAMS', 'HCER', 'HOTS', 'ICEC', 'IBRE', 'IMPO', 'INOO', 'JAMS', 'KETC', 'LJDR', 'MALT', 'SEAS', 'MAYO', 'MEAT', 'MLKM', 'MXCO', 'MXBS', 'MXSB', 'MXCH', 'MXCC', 'MXSN', 'COBT', 
               'COCF', 'CABB', 'MXEC', 'MXDP', 'MXFR', 'MXFM', 'MXMC', 'MXPS', 'MXSO', 'MXSP', 'MXSW', 'MUST', 'NDCR', 'NOOD', 'NUGG', 'OAFL', 'OLIV', 'PANC', 'PANE', 'PAST', 'PSAU', 'PNOU', 'PORK', 'PPMX', 
               'PWSM', 'PCCE', 'DOUG', 'PPIZ', 'REFR', 'RICE', 'RBIS', 'RTEB', 'RTEM', 'SDRE', 'SALT', 'SLTC', 'SARD', 'SAUS', 'SCHN', 'SNAC', 'SNOO', 'SOUP', 'SOYS', 'SPAG', 'SPCH', 'SUGA', 'SWCO', 'SWSP', 
               'SWEE', 'TOAS', 'TOMA', 'TUNA', 'VMLK', 'WFLO', 'AIRC', 'BARS', 'BLEA', 'CBLK', 'CGLO', 'CLSP', 'CLTO', 'FILT', 'CRHC', 'CRLA', 'CRPA', 'DISH', 'DPAC', 'DRUB', 'FBRF', 'FWAX', 'FDEO', 'FRNP', 
               'GBBG', 'GCLE', 'CLEA', 'INSE', 'KITT', 'LAUN', 'LSTA', 'MXBC', 'MXHC', 'MXCB', 'MXLB', 'MXLD', 'CRTO', 'NAPK', 'PLWF', 'SCOU', 'SOFT', 'STRM', 'TOIP', 'WIPE', 'ANLG', 'FSUP', 'GMED', 'VITA', 'nan', 
               'BATT', 'CGAS', 'PFHH', 'PFIN', 'INKC', 'PETF', 'TELE', 'TILL', 'TOBA', 'ADIP', 'BSHM', 'RAZO', 'BDCR', 'CWIP', 'COMB', 'COND', 'CRHY', 'CRPC', 'DEOD', 'DIAP', 'FCCR', 'FTIS', 'FEMI', 'FRAG', 'HAIR', 
               'HRCO', 'HREM', 'HRST', 'HSTY', 'HRTR', 'LINI', 'MAKE', 'MEDS', 'CRDT', 'MXMH', 'MOWA', 'ORAL', 'SPAD', 'STOW', 'SHAM', 'SHAV', 'SKCR', 'SUNP', 'TALC', 'TAMP', 'TOIL', 'TOOB', 'TOOT', 'BAGS', 'CLPC', 
               'GRPC', 'MRKR', 'NTBK', 'SCHS', 'CSTD', 'CORP', 'CROS', 'CRBA', 'CRBR', 'CRDT', 'CRDF', 'CRFO', 'CRSA', 'CRSN', 'DEMO', 'FLSH', 'HLVW', 'COCP', 'CRSN', 'MULT', 'PCHK', 'STCK', 'MIHC','FLWT'],
})

    #obtém o país,categoria,cesta e fabricante para template e ppt

base_dir = Path(__file__).resolve().parent

os.chdir(base_dir)

excel = select_excel_file(base_dir)



file = pd.ExcelFile(str(base_dir / excel))



W = file.sheet_names


#Obtém o pais cesta categoria fabricante marca e idioma para o qual se fará o estudo
land, cesta, cat, client = pais.loc[pais.cod==int(excel.split('_')[0]),'pais'].iloc[0], categ.loc[categ.cod==excel.split('_')[1],'cest'].iloc[0] ,categ.loc[categ.cod==excel.split('_')[1],'cat'].iloc[0] ,excel.split('_')[2].rsplit('.', 1)[0]

lang= "P" if land=='Brasil' else "E"

modelo_path = base_dir / 'Modelo_5W1H.pptx'

if not modelo_path.exists():

    raise FileNotFoundError(colorize(f'No se encontro el template {modelo_path.name} en {base_dir}', COLOR_RED))



ppt= Presentation(str(modelo_path))

configure_cover_slide(ppt, lang)



brand = W[0][2:]
#Dicionario com as correspondencias dos numeros e o respectivo W
c_w={
('P','1'):'1W - Quando?',
('P','3-1'):'3W - O Que? Tamanhos',
('P','3-2'):'3W - O Que? Marcas',
('P','3-3'):'3W - O Que? Sabores',
('P','4'):'4W - Quem? NSE',
('P','5-1'):'5W - Onde? Regiões',
('P','5-2'):'5W - Onde? Canais',
('P','6'):'Players',

('E','1'):'1W - ¿Cuando?',
('E','3-1'):'3W - ¿Lo Que? Tamaños',
('E','3-2'):'3W - ¿Lo Que? Marcas',
('E','3-3'):'3W - ¿Lo Que? Sabores',
('E','4'):'4W - ¿Quién? NSE',
('E','5-1'):'5W - ¿Dónde? Regiones',
('E','5-2'):'5W - ¿Dónde? Canales',
('E','6'):'Players'
}

#Etiquetas dos slides
labels  ={
('P','Data'):'Data',
('P','MAT'):'Avaliação em Ano Móvel Acumulado',
('P','Var MAT'):"em ano móvel",
('P','comp'):"Concorrência do mercado de: ",

('E','Data'):'Fecha',
('E','MAT'):'Evaluación en Año Móvil Acumulado',
('E','Var MAT'):"en año móvil",
('E','comp'):"Competencia para el mercado de: "}

class SeriesConfig(NamedTuple):
    data: pd.DataFrame
    raw_tipo: str
    display_tipo: str
    pipeline: int

def prepare_series_configs(df_list, lang, p_ventas):
    configs = []
    for original_df in df_list:
        df_local = original_df.copy()
        first_col = df_local.iloc[:,0]
        if first_col.dtype == object:
            df_local.iloc[:,0] = pd.to_datetime(first_col, format='%b-%y  ')
        raw_tipo = str(df_local.columns[0])
        df_local.rename(columns={df_local.columns[0]: labels[(lang,'Data')]}, inplace=True)
        is_compras = 'compras' in raw_tipo.lower()
        if not is_compras and len(df_local.columns) > 1:
            is_compras = '.c' in str(df_local.columns[1]).lower()
        display_tipo = 'Compras' if is_compras else 'Ventas'
        pipeline = 0 if is_compras else p_ventas
        configs.append(SeriesConfig(df_local, raw_tipo, display_tipo, pipeline))
    return configs








#Variavel de controle do numero de graficos
c_fig=0

plot=plot_ven()
last_reference_source = None
#---------------------------------------------------------------------------------------------------------------------
for w in W:

    #1- Quando, grafico de variações MAT
    if w[0]=='1':

        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[35])

        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = c_w[(lang,w[0])]+' '+ labels[(lang,'MAT')]+' | ' + w[2:] 
        t.font.bold = True
        t.font.size = Inches(0.35)

        #Obtém pipeline das vendas
        p=int(file.parse(w).columns[1].split("_")[1])

        #Cria a base
        mat=df_mat(file.parse(w),p)

        last_reference_source = mat
        
        #Elimina linhas com divisão com zero devido ao pipeline
        mat=mat[~np.isinf(mat.iloc[:, 3])]
      
        #Incrementa contador do gráfico
        c_fig+=1
        
        #Insere o gráfico do MAT
        left_margin = Inches(0.33)
        right_margin = Inches(0.33)
        available = available_width(ppt, left_margin, right_margin)
        pic=slide.shapes.add_picture(graf_mat(mat,c_fig,p), left_margin, Inches(1.15),width=available)

        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(5.8), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = "Comentário"
        t.font.size = Inches(0.28)

        #Limpa área de plotagem
        plt.clf()

        #mensaje de conclusion por cada slide
        print_colored(c_w[(lang,w[0])]+' realizado para '+ w[2:], COLOR_GREEN)

    #Outros
    else: 

        #Carrega a base
        df_start=file.parse(w)

        #Lista para as bases
        df_list=[]
        p_ventas=0

        #Verifica se hay dados de vendas para obter as bases e o pipeline
        try:
            start,p=[(i-1,col.split("_")[1]) for i,col in enumerate(df_start.columns) if "ventas" in col.lower()][0]
            p_ventas=int(p)
            compras=df_start.iloc[:,:start]
            ventas=df_start.iloc[:,start+1:]
            compras.columns= [compras.columns[0]]+[x+'.C' for x in compras.columns[1:]]
            ventas.columns= [ventas.columns[0][:-2]]+[x.replace('.1','.V') for x in ventas.columns[1:]]
            df_list.append(compras)
            df_list.append(ventas)
        except:
            df_list.append(df_start)

        series_configs = prepare_series_configs(df_list, lang, p_ventas)
        last_reference_source = series_configs[0].data if series_configs else df_start
        #Cria o slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[35])

        #Define o titulo
        txTitle = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(10), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        if w[0] in ['3','5']:
            titulo = c_w[(lang,w[0]+'-'+w[-1])]+' | ' + w[2:-2] 
            t.text = titulo
            t.font.size = Inches(0.35)
        elif w[0]=='6':
            titulo =  c_w[(lang,w[0])] + ' | ' + labels[(lang,'comp')] + cat
            t.text =  titulo
            t.font.size = Inches(0.33)
        else:
            titulo = c_w[(lang,w[0])]+' | '+ w[2:] 
            t.text = titulo 
            t.font.size = Inches(0.35)

        t.font.bold = True

        #Insere caixa de texto para comentário do slide
        txTitle = slide.shapes.add_textbox(Inches(11.07), Inches(6.33), Inches(2), Inches(0.5))
        tf = txTitle.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text ="Comentário"
        t.font.size= Inches(0.25)

        #Controle posicao tabela aporte
        target_table_height = Cm(TABLE_TARGET_HEIGHT_CM)
        bottom_margin = Cm(1.5)
        left_margin = Cm(1.2)
        vertical_spacing = Cm(0.2)

        for idx, serie in enumerate(series_configs):
            apo = aporte(serie.data.copy(), serie.pipeline, lang, serie.raw_tipo)
            c_fig += 1
            # Si hay dos tablas, ambas se posicionan a la misma altura y ocupan el espacio óptimo
            if len(series_configs) == 2:
                left_margin = Cm(1.2)
                right_margin = Cm(1.2)
                gap = Cm(0.5)
                top_position = ppt.slide_height - target_table_height - bottom_margin
                if serie.display_tipo.lower() == 'ventas':
                    left_position = ppt.slide_width / 2 + gap / 2
                else:
                    left_position = left_margin
                pic = slide.shapes.add_picture(graf_apo(apo, c_fig), left_position, top_position, height=target_table_height)
            else:
                top_position = ppt.slide_height - target_table_height - bottom_margin
                left_position = left_margin
                pic = slide.shapes.add_picture(graf_apo(apo, c_fig), left_position, top_position, height=target_table_height)
            plt.clf()

        if plot=="1" and len(series_configs)>1:
            df_full = series_configs[0].data.copy()
            for extra in series_configs[1:]:
                df_full = pd.concat([df_full, extra.data.iloc[:,1:]], axis=1)
            pipeline_combined = max((cfg.pipeline for cfg in series_configs), default=0)
            c_fig+=1
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available_line_width = available_width(ppt, left_margin, right_margin)
            pic=slide.shapes.add_picture(line_graf(df_full,pipeline_combined,titulo,c_fig,len(series_configs), width_emu=available_line_width, height_emu=Cm(10)), left_margin, Inches(1.15),width=available_line_width,height=Cm(10))
        elif plot=="2" and len(series_configs)>1:
            ven_param = max(len(series_configs)-1, 1)
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available = available_width(ppt, left_margin, right_margin)
            count = len(series_configs)
            gap = Inches(0.1) if count > 1 else Inches(0)
            total_gap = int(gap) * (count - 1)
            effective_width = available - total_gap
            if count:
                if effective_width <= 0:
                    chart_width = available // count
                else:
                    chart_width = effective_width // count
            else:
                chart_width = available
            chart_width = max(chart_width, 1)
            for idx, serie in enumerate(series_configs):
                c_fig+=1
                left_position = left_margin + idx * (chart_width + int(gap))
                pic=slide.shapes.add_picture(line_graf(serie.data,serie.pipeline,titulo+' '+serie.display_tipo,c_fig,ven_param, width_emu=chart_width, height_emu=Cm(10)), left_position, Inches(1.15),width=chart_width,height=Cm(10))
                plt.clf()
        elif series_configs:
            c_fig+=1
            left_margin = Inches(0.33)
            right_margin = Inches(0.33)
            available_single_width = available_width(ppt, left_margin, right_margin)
            pic=slide.shapes.add_picture(line_graf(series_configs[0].data,series_configs[0].pipeline,titulo,c_fig,len(series_configs), width_emu=available_single_width, height_emu=Cm(10)), left_margin, Inches(1.15),width=available_single_width,height=Cm(10))
            plt.clf()
        if w[0] in ['3','5']:
            print_colored(c_w[(lang,w[0]+'-'+w[-1])]+' realizado para ' + w[2:-2], COLOR_GREEN) 
        elif w[0]=='6':
            print_colored(c_w[(lang,w[0])] + ' realizado para ' + labels[(lang,'comp')] + cat, COLOR_GREEN)
        else:
            print_colored(c_w[(lang,w[0])]+' realizado para '+ w[2:], COLOR_GREEN)


#Referencia da base
ref_source = last_reference_source if last_reference_source is not None else file.parse(W[0])
last_value = ref_source.iloc[-1, 0] if not ref_source.empty else None
if last_value is None:
    ref = 'NA'
    ref_display = 'NA'
else:
    if isinstance(last_value, str):
        try:
            last_dt = dt.strptime(last_value, '%b-%y  ')
        except ValueError:
            last_dt = pd.to_datetime(last_value).to_pydatetime()
    else:
        last_dt = pd.to_datetime(last_value).to_pydatetime()
    ref = last_dt.strftime('%m-%y')
    month_list = MONTH_NAMES.get(lang, MONTH_NAMES['E'])
    month_name = month_list[last_dt.month - 1]
    ref_display = f"{month_name}-{last_dt.year}"
slide=ppt.slides[0]

#Fabricante
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(3.77), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
t.text = "Numerator vs " + client
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.44), True
font.color.rgb = RGBColor(255, 255, 255)

#Referencia
txBox = slide.shapes.add_textbox(Inches(0.42), Inches(5.98), Inches(10), Inches(0.5))
tf = txBox.text_frame
tf.clear()
t = tf.paragraphs[0]
reference_label = 'Corte em' if lang == 'P' else 'Corte a'
display_ref = ref_display if ref_display != 'NA' else ref
t.text = f"{reference_label} {display_ref}"
run = t.runs[0]
font = run.font
font.name, font.size, font.bold = 'Arial', Inches(0.32), True
font.color.rgb = RGBColor(255, 255, 255)

ppt.save(land+'-'+cat+'-'+client+'-'+brand+'-5W1H-'+ref+'.pptx')

fim = dt.now()

t = int((fim - agora).total_seconds())
print_colored(f'Tiempo de ejecucion : {t//60} min {t%60} s' if t >= 60 else f'Tiempo de ejecucion : {t} s', COLOR_BLUE)

